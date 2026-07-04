#!/usr/bin/env python3
"""Export an HTML slide deck to a widescreen PPTX.

The v1 export is fidelity-first: each rendered HTML slide is captured as a
1920x1080 PNG and placed full-bleed on a PowerPoint slide. The resulting deck
opens everywhere PowerPoint files are supported, but slide content is flattened
into images rather than converted into editable Office shapes.
"""

from __future__ import annotations

import argparse
import json
import os
from pathlib import Path
import shutil
import subprocess
import sys
import tempfile
from textwrap import dedent

from pptx import Presentation
from pptx.util import Inches


DEFAULT_WIDTH = 1920
DEFAULT_HEIGHT = 1080
PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Export an HTML presentation to a screenshot-backed PPTX."
    )
    parser.add_argument("html", help="Path to the HTML presentation.")
    parser.add_argument(
        "output",
        nargs="?",
        help="Output .pptx path. Defaults to the HTML filename with .pptx.",
    )
    parser.add_argument(
        "--screenshots-dir",
        help=(
            "Use existing PNG screenshots instead of launching a browser. "
            "Useful for tests or custom capture pipelines."
        ),
    )
    parser.add_argument(
        "--keep-screenshots",
        action="store_true",
        help="Keep generated screenshots next to the output PPTX.",
    )
    parser.add_argument(
        "--width",
        type=int,
        default=DEFAULT_WIDTH,
        help=f"Capture width in pixels. Default: {DEFAULT_WIDTH}.",
    )
    parser.add_argument(
        "--height",
        type=int,
        default=DEFAULT_HEIGHT,
        help=f"Capture height in pixels. Default: {DEFAULT_HEIGHT}.",
    )
    return parser.parse_args()


def natural_pngs(directory: Path) -> list[Path]:
    def sort_key(path: Path) -> tuple[int, str]:
        digits = "".join(ch for ch in path.stem if ch.isdigit())
        return (int(digits) if digits else 0, path.name)

    return sorted(directory.glob("*.png"), key=sort_key)


def run(command: list[str], *, cwd: Path | None = None) -> None:
    try:
        subprocess.run(command, cwd=cwd, check=True)
    except FileNotFoundError as exc:
        raise SystemExit(f"Missing executable: {command[0]}") from exc
    except subprocess.CalledProcessError as exc:
        raise SystemExit(f"Command failed ({exc.returncode}): {' '.join(command)}") from exc


def _path_under(path: Path, root: Path) -> Path | None:
    try:
        return path.resolve().relative_to(root.resolve())
    except ValueError:
        return None


def resolve_html_path(raw_path: str) -> Path:
    path = Path(raw_path).expanduser()
    if path.is_absolute():
        return path.resolve()

    workspace_root = os.environ.get("HELPUDOC_WORKSPACE_ROOT")
    if workspace_root:
        candidate = (Path(workspace_root) / path).resolve()
        if candidate.exists():
            return candidate

    return path.resolve()


def resolve_output_path(raw_output: str | None, html_path: Path) -> tuple[Path, str | None]:
    workspace_root_raw = os.environ.get("HELPUDOC_WORKSPACE_ROOT")
    workspace_output_raw = os.environ.get("HELPUDOC_WORKSPACE_OUTPUT_ROOT")
    workspace_root = Path(workspace_root_raw).resolve() if workspace_root_raw else None
    workspace_output_root = Path(workspace_output_raw).resolve() if workspace_output_raw else None

    if raw_output:
        requested = Path(raw_output).expanduser()
        if requested.is_absolute():
            requested_abs = requested.resolve()
            rel_output = _path_under(requested_abs, workspace_root) if workspace_root else None
            if workspace_output_root and rel_output is not None:
                return (workspace_output_root / rel_output).resolve(), rel_output.as_posix()
            return requested_abs, rel_output.as_posix() if rel_output is not None else None
        rel_output = requested
    else:
        html_rel = _path_under(html_path, workspace_root) if workspace_root else None
        rel_output = html_rel.with_suffix(".pptx") if html_rel is not None else Path(html_path.with_suffix(".pptx").name)

    if workspace_output_root:
        return (workspace_output_root / rel_output).resolve(), rel_output.as_posix()
    output_path = (Path.cwd() / rel_output).resolve() if not rel_output.is_absolute() else rel_output.resolve()
    return output_path, rel_output.as_posix() if not rel_output.is_absolute() else None


def emit_tool_artifact(output_path: Path, workspace_rel_path: str | None, slide_count: int) -> None:
    sandbox_run_dir_raw = os.environ.get("HELPUDOC_SANDBOX_RUN_DIR")
    if not sandbox_run_dir_raw or not workspace_rel_path:
        return
    payload_path = Path(sandbox_run_dir_raw) / "out" / "tool_artifacts.json"
    payload_path.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "files": [
            {
                "path": workspace_rel_path,
                "mimeType": PPTX_MIME_TYPE,
                "size": output_path.stat().st_size if output_path.exists() else 0,
                "metadata": {
                    "kind": "frontend-slides-pptx",
                    "slideCount": slide_count,
                },
            }
        ]
    }
    payload_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")


def node_capture_script() -> str:
    return dedent(
        r"""
        import { chromium } from 'playwright';
        import { createServer } from 'http';
        import { readFileSync, existsSync, mkdirSync } from 'fs';
        import { join, extname, resolve } from 'path';

        const serveDir = process.argv[2];
        const htmlFile = process.argv[3];
        const outDir = process.argv[4];
        const width = Number(process.argv[5] || 1920);
        const height = Number(process.argv[6] || 1080);

        const mimeTypes = {
          '.html': 'text/html',
          '.css': 'text/css',
          '.js': 'application/javascript',
          '.json': 'application/json',
          '.png': 'image/png',
          '.jpg': 'image/jpeg',
          '.jpeg': 'image/jpeg',
          '.gif': 'image/gif',
          '.svg': 'image/svg+xml',
          '.webp': 'image/webp',
          '.woff': 'font/woff',
          '.woff2': 'font/woff2',
          '.ttf': 'font/ttf',
        };

        mkdirSync(outDir, { recursive: true });

        const server = createServer((req, res) => {
          const decoded = decodeURIComponent(req.url || '/');
          const relative = decoded === '/' ? htmlFile : decoded.replace(/^\/+/, '');
          const filePath = resolve(join(serveDir, relative));
          if (!filePath.startsWith(resolve(serveDir))) {
            res.writeHead(403);
            res.end('Forbidden');
            return;
          }
          try {
            const content = readFileSync(filePath);
            const ext = extname(filePath).toLowerCase();
            res.writeHead(200, { 'Content-Type': mimeTypes[ext] || 'application/octet-stream' });
            res.end(content);
          } catch {
            res.writeHead(404);
            res.end('Not found');
          }
        });

        await new Promise((resolveListen) => server.listen(0, '127.0.0.1', resolveListen));
        const port = server.address().port;

        const browser = await chromium.launch({ headless: true });
        const page = await browser.newPage({ viewport: { width, height }, deviceScaleFactor: 1 });
        await page.goto(`http://127.0.0.1:${port}/${encodeURIComponent(htmlFile)}`, {
          waitUntil: 'networkidle',
        });

        await page.evaluate(() => document.fonts && document.fonts.ready);
        await page.waitForTimeout(300);

        const slideCount = await page.evaluate(() => document.querySelectorAll('.slide').length);
        if (!slideCount) {
          await browser.close();
          server.close();
          throw new Error('No .slide elements found in the presentation.');
        }

        for (let index = 0; index < slideCount; index += 1) {
          await page.evaluate((activeIndex) => {
            const deckStage = document.querySelector('deck-stage');
            if (deckStage) deckStage.setAttribute('noscale', '');

            const stage = document.querySelector('#deckStage, .deck-stage');
            if (stage) {
              stage.style.transform = 'none';
              stage.style.left = '0';
              stage.style.top = '0';
            }

            const slides = Array.from(document.querySelectorAll('.slide'));
            slides.forEach((slide, slideIndex) => {
              const active = slideIndex === activeIndex;
              slide.classList.toggle('active', active);
              slide.classList.toggle('visible', active);
              slide.style.visibility = active ? 'visible' : 'hidden';
              slide.style.opacity = active ? '1' : '0';
              slide.style.pointerEvents = active ? 'auto' : 'none';
              slide.style.zIndex = active ? '1' : '0';
            });

            window.dispatchEvent(new CustomEvent('slidechange', { detail: { index: activeIndex } }));
          }, index);
          await page.waitForTimeout(150);
          const path = join(outDir, `slide-${String(index + 1).padStart(3, '0')}.png`);
          await page.screenshot({ path, fullPage: false });
        }

        await browser.close();
        server.close();
        console.log(JSON.stringify({ slideCount, outDir }));
        """
    ).strip()


def capture_screenshots(html_path: Path, output_dir: Path, width: int, height: int) -> list[Path]:
    if not shutil.which("npm"):
        raise SystemExit("Node.js/npm is required for browser capture.")

    with tempfile.TemporaryDirectory(prefix="frontend-slides-pptx-node-") as tmp:
        tmp_path = Path(tmp)
        (tmp_path / "package.json").write_text(
            '{"name":"frontend-slides-pptx-export","private":true,"type":"module"}\n',
            encoding="utf-8",
        )
        script_path = tmp_path / "capture-slides.mjs"
        script_path.write_text(node_capture_script(), encoding="utf-8")

        run(["npm", "install", "--silent", "playwright"], cwd=tmp_path)
        run(["npx", "playwright", "install", "chromium"], cwd=tmp_path)
        run(
            [
                "node",
                str(script_path),
                str(html_path.parent),
                html_path.name,
                str(output_dir),
                str(width),
                str(height),
            ],
            cwd=tmp_path,
        )

    screenshots = natural_pngs(output_dir)
    if not screenshots:
        raise SystemExit(f"No screenshots were generated in {output_dir}")
    return screenshots


def build_pptx(screenshots: list[Path], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for screenshot in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(screenshot),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> int:
    args = parse_args()
    html_path = resolve_html_path(args.html)
    if not html_path.exists():
        print(f"HTML file not found: {html_path}", file=sys.stderr)
        return 1

    output_path, workspace_rel_output = resolve_output_path(args.output, html_path)

    if args.screenshots_dir:
        screenshots_dir = Path(args.screenshots_dir).expanduser().resolve()
        screenshots = natural_pngs(screenshots_dir)
        if not screenshots:
            print(f"No PNG screenshots found in {screenshots_dir}", file=sys.stderr)
            return 1
        build_pptx(screenshots, output_path)
        emit_tool_artifact(output_path, workspace_rel_output, len(screenshots))
        print(f"Exported {len(screenshots)} slides to {output_path}")
        return 0

    with tempfile.TemporaryDirectory(prefix="frontend-slides-pptx-") as tmp:
        screenshots_dir = Path(tmp) / "screenshots"
        screenshots_dir.mkdir(parents=True, exist_ok=True)
        screenshots = capture_screenshots(html_path, screenshots_dir, args.width, args.height)
        build_pptx(screenshots, output_path)

        if args.keep_screenshots:
            kept_dir = output_path.with_suffix("")
            kept_dir = kept_dir.parent / f"{kept_dir.name}-screenshots"
            if kept_dir.exists():
                shutil.rmtree(kept_dir)
            shutil.copytree(screenshots_dir, kept_dir)
            print(f"Saved screenshots to {kept_dir}")

    emit_tool_artifact(output_path, workspace_rel_output, len(screenshots))
    print(f"Exported {len(screenshots)} slides to {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
