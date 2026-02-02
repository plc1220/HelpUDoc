"""
Export utilities for Paper2Slides outputs.
"""
import io
import json
import logging
import os
import shutil
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches

from .pptx_builder import PPTXBuilder


class _SimpleTextStyle:
    def __init__(
        self,
        font_color_rgb: Optional[List[int]] = None,
        text_alignment: Optional[str] = None,
        is_bold: bool = False,
        is_italic: bool = False,
        is_underline: bool = False,
    ) -> None:
        self.font_color_rgb = font_color_rgb
        self.text_alignment = text_alignment
        self.is_bold = is_bold
        self.is_italic = is_italic
        self.is_underline = is_underline


logger = logging.getLogger(__name__)


class ExportService:
    """Service for exporting slides to PPTX."""

    @staticmethod
    def create_pptx_from_images(image_paths: List[str], output_file: Optional[str] = None) -> Optional[bytes]:
        """
        Create a PPTX file from image paths.

        Args:
            image_paths: Ordered list of image paths.
            output_file: Optional output file path (if None, returns bytes).

        Returns:
            PPTX bytes if output_file is None, otherwise None.
        """
        if not image_paths:
            raise ValueError("No image paths provided for PPTX export")

        # Validate and preserve order
        valid_paths = []
        for image_path in image_paths:
            if os.path.exists(image_path):
                valid_paths.append(image_path)
            else:
                logger.warning("Image not found for PPTX export: %s", image_path)

        if not valid_paths:
            raise ValueError("No valid images found for PPTX export")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        blank_slide_layout = prs.slide_layouts[6]
        for image_path in valid_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        if output_file:
            prs.save(output_file)
            return None

        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        return pptx_bytes.getvalue()

    @staticmethod
    def create_pptx_from_pdf(
        pdf_path: str,
        output_file: Optional[str] = None,
        slide_width_px: int = 1920,
        slide_height_px: int = 1080,
    ) -> Optional[bytes]:
        """
        Render PDF pages to images and create a PPTX deck.

        Args:
            pdf_path: Path to the PDF file.
            output_file: Optional output file path (if None, returns bytes).
            slide_width_px: Target slide width in pixels for rendering.
            slide_height_px: Target slide height in pixels for rendering.
        """
        pdf_file = Path(pdf_path)
        if not pdf_file.exists():
            raise FileNotFoundError(f"PDF not found: {pdf_path}")

        image_paths, temp_dir = ExportService._render_pdf_backgrounds(
            pdf_file,
            slide_width_px,
            slide_height_px,
        )
        if not image_paths:
            raise RuntimeError(
                "No images rendered from PDF. Ensure PyMuPDF (fitz) is installed."
            )

        try:
            return ExportService.create_pptx_from_images(image_paths, output_file)
        finally:
            if temp_dir:
                shutil.rmtree(temp_dir, ignore_errors=True)

    @staticmethod
    def create_editable_pptx_from_images_via_mineru(
        image_paths: List[str],
        output_file: str,
        lang: Optional[str] = None,
        include_background: bool = True,
    ) -> None:
        """
        Create editable PPTX from slide images by running MinerU OCR per image.
        """
        if not image_paths:
            raise ValueError("No image paths provided for editable PPTX export")

        from paper2slides.raganything.parser import MineruParser

        try:
            from PIL import Image
        except Exception as exc:
            raise RuntimeError(
                "Pillow is required to read image dimensions for PPTX export."
            ) from exc

        valid_paths = []
        for image_path in image_paths:
            if os.path.exists(image_path):
                valid_paths.append(image_path)
            else:
                logger.warning("Image not found for editable PPTX export: %s", image_path)

        if not valid_paths:
            raise ValueError("No valid images found for editable PPTX export")

        mineru = MineruParser()
        builder = PPTXBuilder()
        builder.create_presentation()

        slide_width_px = None
        slide_height_px = None
        content_by_slide: List[List[Dict[str, Any]]] = []

        temp_dir = Path(tempfile.mkdtemp(prefix="mineru-slides-"))
        try:
            for image_path in valid_paths:
                with Image.open(image_path) as img:
                    width_px, height_px = img.size

                if slide_width_px is None or slide_height_px is None:
                    slide_width_px = width_px
                    slide_height_px = height_px
                    builder.setup_presentation_size(slide_width_px, slide_height_px)

                content_list = mineru.parse_image(
                    image_path=image_path,
                    output_dir=str(temp_dir),
                    lang=lang,
                )
                content_by_slide.append(content_list)

            if slide_width_px is None or slide_height_px is None:
                raise ValueError("Unable to determine slide dimensions from images")

            for idx, image_path in enumerate(valid_paths):
                slide = builder.add_blank_slide()

                if include_background:
                    slide.shapes.add_picture(
                        image_path,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height,
                    )

                slide_items = content_by_slide[idx]
                if not slide_items:
                    continue

                with Image.open(image_path) as img:
                    width_px, height_px = img.size

                scale_x = slide_width_px / width_px
                scale_y = slide_height_px / height_px

                text_items = []
                image_items = []
                for item in slide_items:
                    item_type = item.get("type", "")
                    if item_type in ["text", "title", "header", "footer"]:
                        text_items.append(item)
                    elif item_type in ["image", "table", "equation"]:
                        image_items.append(item)

                for img_item in image_items:
                    ExportService._add_mineru_image_to_slide(
                        builder, slide, img_item, temp_dir, scale_x, scale_y
                    )

                for text_item in text_items:
                    ExportService._add_mineru_text_to_slide(
                        builder, slide, text_item, scale_x, scale_y
                    )
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

        builder.save(output_file)

    @staticmethod
    def create_editable_pptx_from_slide_assets(
        slide_dirs: List[str],
        output_file: str,
        include_background: bool = True,
    ) -> None:
        """
        Build editable PPTX from per-slide asset folders.

        Each slide directory should include layout.json and optional assets/background images.
        """
        if not slide_dirs:
            raise ValueError("No slide directories provided for PPTX export")

        builder = PPTXBuilder()
        builder.create_presentation()

        slide_width_px: Optional[int] = None
        slide_height_px: Optional[int] = None

        for slide_dir in slide_dirs:
            slide_path = Path(slide_dir)
            layout_path = slide_path / "layout.json"
            if not layout_path.exists():
                logger.warning("Missing layout.json in %s", slide_dir)
                continue

            try:
                with open(layout_path, "r", encoding="utf-8") as handle:
                    layout = json.load(handle)
            except Exception as exc:
                logger.warning("Failed to read layout.json in %s: %s", slide_dir, exc)
                continue

            canvas = layout.get("canvas", {}) if isinstance(layout, dict) else {}
            width_px = ExportService._coerce_int(canvas.get("width"))
            height_px = ExportService._coerce_int(canvas.get("height"))
            if width_px is None or height_px is None:
                logger.warning("Invalid canvas size in %s", slide_dir)
                continue

            if slide_width_px is None or slide_height_px is None:
                slide_width_px = width_px
                slide_height_px = height_px
                builder.setup_presentation_size(slide_width_px, slide_height_px)

            scale_x = slide_width_px / width_px
            scale_y = slide_height_px / height_px

            slide = builder.add_blank_slide()

            if include_background:
                bg_path = ExportService._resolve_slide_path(slide_path, layout.get("background"))
                if not bg_path:
                    bg_path = ExportService._resolve_slide_path(
                        slide_path, layout.get("source_image")
                    )
                if bg_path and os.path.exists(bg_path):
                    slide.shapes.add_picture(
                        bg_path,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height,
                    )

            ExportService._render_layout_elements(
                builder, slide, layout, slide_path, scale_x, scale_y
            )

        builder.save(output_file)

    @staticmethod
    def _render_layout_elements(
        builder: PPTXBuilder,
        slide,
        layout: Dict[str, Any],
        slide_path: Path,
        scale_x: float,
        scale_y: float,
    ) -> None:
        elements = layout.get("elements")
        if not isinstance(elements, list):
            return

        for idx, element in enumerate(elements, start=1):
            if not isinstance(element, dict):
                continue
            element_type = str(element.get("type", "")).lower().strip()
            bbox = element.get("bbox")
            scaled_bbox = ExportService._scale_bbox(bbox, scale_x, scale_y)
            if not scaled_bbox:
                continue

            if element_type == "text":
                text = str(element.get("text", "")).strip()
                if not text:
                    continue
                align = str(element.get("align", "left")).lower()
                font_size = element.get("font_size")
                color = element.get("color_rgb")
                style = _SimpleTextStyle(
                    font_color_rgb=color if isinstance(color, list) else None,
                    text_alignment=None,
                    is_bold=bool(element.get("bold", False)),
                    is_italic=bool(element.get("italic", False)),
                    is_underline=bool(element.get("underline", False)),
                )
                builder.add_text_element(
                    slide=slide,
                    text=text,
                    bbox=scaled_bbox,
                    align=align,
                    font_size=font_size,
                    text_style=style,
                )
                continue

            if element_type == "table" and element.get("cells"):
                builder.add_table_element(
                    slide=slide,
                    html_table=element.get("cells"),
                    bbox=scaled_bbox,
                )
                continue

            asset_path = element.get("asset_path")
            if not asset_path and slide_path:
                candidate = slide_path / "assets" / f"asset-{idx:03d}.png"
                if candidate.exists():
                    asset_path = str(candidate)
            resolved_path = ExportService._resolve_slide_path(slide_path, asset_path)
            if not resolved_path:
                continue

            builder.add_image_element(
                slide=slide,
                image_path=resolved_path,
                bbox=scaled_bbox,
            )

    @staticmethod
    def _scale_bbox(
        bbox: Any, scale_x: float, scale_y: float
    ) -> Optional[List[int]]:
        if not isinstance(bbox, (list, tuple)) or len(bbox) != 4:
            return None
        try:
            x0, y0, x1, y1 = bbox
            return [
                int(round(x0 * scale_x)),
                int(round(y0 * scale_y)),
                int(round(x1 * scale_x)),
                int(round(y1 * scale_y)),
            ]
        except Exception:
            return None

    @staticmethod
    def _resolve_slide_path(slide_path: Path, value: Any) -> Optional[str]:
        if not value or not isinstance(value, str):
            return None
        candidate = Path(value)
        if not candidate.is_absolute():
            candidate = slide_path / candidate
        return str(candidate.resolve())

    @staticmethod
    def _coerce_int(value: Any) -> Optional[int]:
        try:
            return int(value)
        except (TypeError, ValueError):
            return None

    @staticmethod
    def find_mineru_result_dir(search_root: str) -> Optional[str]:
        """
        Find the most recent MinerU result directory containing *_content_list.json.
        """
        root = Path(search_root)
        if not root.exists():
            return None

        candidates = list(root.rglob("*_content_list.json"))
        if not candidates:
            return None

        def score(path: Path) -> float:
            mtime = path.stat().st_mtime
            bonus = 1 if "auto" in path.parts else 0
            return mtime + bonus

        best = max(candidates, key=score)
        return str(best.parent)

    @staticmethod
    def _infer_page_dimensions(content_list: List[Dict[str, Any]]) -> Optional[Tuple[int, int]]:
        max_x = 0
        max_y = 0
        for item in content_list:
            bbox = item.get("bbox")
            if not bbox or len(bbox) != 4:
                continue
            max_x = max(max_x, int(bbox[2]))
            max_y = max(max_y, int(bbox[3]))
        if max_x <= 0 or max_y <= 0:
            return None
        return max_x, max_y

    @staticmethod
    def _find_source_file(mineru_dir: Path, base_stem: str, extensions: List[str]) -> Optional[Path]:
        for parent in [mineru_dir, mineru_dir.parent, mineru_dir.parent.parent]:
            for ext in extensions:
                candidate = parent / f"{base_stem}{ext}"
                if candidate.exists():
                    return candidate
        return None

    @staticmethod
    def _render_pdf_backgrounds(
        pdf_path: Path,
        target_width_px: int,
        target_height_px: int,
    ) -> Tuple[List[str], Optional[Path]]:
        try:
            import fitz  # type: ignore
        except Exception as exc:
            logger.warning("PyMuPDF not available for background rendering: %s", exc)
            return [], None

        doc = fitz.open(str(pdf_path))
        temp_dir = Path(tempfile.mkdtemp(prefix="paper2slides-bg-"))
        image_paths = []
        for page_index in range(len(doc)):
            page = doc[page_index]
            rect = page.rect
            scale_x = target_width_px / rect.width
            scale_y = target_height_px / rect.height
            matrix = fitz.Matrix(scale_x, scale_y)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            path = temp_dir / f"page_{page_index + 1:03d}.png"
            pix.save(str(path))
            image_paths.append(str(path))
        return image_paths, temp_dir

    @staticmethod
    def create_editable_pptx_from_mineru(
        mineru_result_dir: str,
        output_file: str = None,
        slide_width_pixels: int = 1920,
        slide_height_pixels: int = 1080,
        background_images: List[str] = None,
    ) -> Optional[bytes]:
        """
        Create editable PPTX file from MinerU parsing results.
        """
        from paper2slides.utils.pptx_builder import PPTXBuilder

        mineru_dir = Path(mineru_result_dir)

        content_list_files = list(mineru_dir.glob("*_content_list.json"))
        if not content_list_files:
            raise FileNotFoundError(f"No content_list.json found in {mineru_result_dir}")

        content_list_file = content_list_files[0]
        logger.info("Loading MinerU content from: %s", content_list_file)

        with open(content_list_file, "r", encoding="utf-8") as f:
            content_list = json.load(f)

        if not content_list:
            raise ValueError("Empty content list from MinerU")

        logger.info("Loaded %s items from MinerU content_list", len(content_list))

        base_stem = content_list_file.name.replace("_content_list.json", "")

        layout_file = mineru_dir / "layout.json"
        layout_data = None
        actual_page_width = slide_width_pixels
        actual_page_height = slide_height_pixels
        use_layout_coords = False

        if layout_file.exists():
            try:
                with open(layout_file, "r", encoding="utf-8") as f:
                    layout_data = json.load(f)
                    if "pdf_info" in layout_data and len(layout_data["pdf_info"]) > 0:
                        page_size = layout_data["pdf_info"][0].get("page_size")
                        if page_size and len(page_size) == 2:
                            actual_page_width, actual_page_height = page_size
                            use_layout_coords = True
                            logger.info(
                                "Using layout.json for coordinates: %sx%s",
                                actual_page_width,
                                actual_page_height,
                            )
                        else:
                            logger.warning("page_size not found in layout.json")
                    else:
                        logger.warning("pdf_info not found in layout.json")
            except Exception as exc:
                logger.warning("Failed to read layout.json: %s", exc)
        else:
            logger.warning("layout.json not found, using content_list coordinates")

        inferred = ExportService._infer_page_dimensions(content_list)
        if not use_layout_coords and inferred:
            actual_page_width, actual_page_height = inferred
            if slide_width_pixels == 1920 and slide_height_pixels == 1080:
                slide_width_pixels, slide_height_pixels = inferred

        logger.info("Target slide dimensions: %sx%s", slide_width_pixels, slide_height_pixels)
        logger.info("Actual page dimensions: %sx%s", actual_page_width, actual_page_height)

        temp_dir = None
        if background_images is None:
            background_images = []
            pdf_path = ExportService._find_source_file(mineru_dir, base_stem, [".pdf"])
            if pdf_path:
                background_images, temp_dir = ExportService._render_pdf_backgrounds(
                    pdf_path,
                    slide_width_pixels,
                    slide_height_pixels,
                )
            if not background_images:
                image_path = ExportService._find_source_file(
                    mineru_dir,
                    base_stem,
                    [".png", ".jpg", ".jpeg", ".webp"],
                )
                if image_path:
                    background_images = [str(image_path)]

        text_level_map = {}
        for item in content_list:
            if item.get("type") == "text" and "text" in item:
                text = item["text"].strip()
                text_level_map[text] = item.get("text_level")

        pages_content: Dict[int, List[Dict[str, Any]]] = {}

        if use_layout_coords and layout_data:
            logger.info("Using layout.json coordinates (accurate)")
            for page_info in layout_data.get("pdf_info", []):
                page_idx = page_info.get("page_idx", 0)
                pages_content[page_idx] = []

                for block in page_info.get("para_blocks", []):
                    block_type = block.get("type", "text")
                    bbox = block.get("bbox")
                    if not bbox:
                        continue

                    if block_type in ["text", "title"] and block.get("lines"):
                        for line in block["lines"]:
                            for span in line.get("spans", []):
                                if span.get("type") == "text" and span.get("content"):
                                    text = span["content"].strip()
                                    text_level = text_level_map.get(text)
                                    pages_content[page_idx].append(
                                        {
                                            "type": block_type,
                                            "text": text,
                                            "text_level": text_level,
                                            "bbox": bbox,
                                            "page_idx": page_idx,
                                        }
                                    )
                    elif block_type in ["image", "table", "equation"]:
                        img_path = block.get("image_path")
                        if not img_path and block.get("blocks"):
                            for sub_block in block["blocks"]:
                                for line in sub_block.get("lines", []):
                                    for span in line.get("spans", []):
                                        if span.get("image_path"):
                                            img_path = span["image_path"]
                                            break
                                    if img_path:
                                        break
                                if img_path:
                                    break
                        if img_path:
                            pages_content[page_idx].append(
                                {
                                    "type": block_type,
                                    "img_path": img_path,
                                    "bbox": bbox,
                                    "page_idx": page_idx,
                                    "html_table": block.get("html_table"),
                                }
                            )
        else:
            logger.info("Using content_list.json coordinates (need scaling)")
            for item in content_list:
                page_idx = item.get("page_idx", 0)
                pages_content.setdefault(page_idx, []).append(item)

        scale_x = slide_width_pixels / actual_page_width
        scale_y = slide_height_pixels / actual_page_height

        builder = PPTXBuilder()
        builder.create_presentation()
        builder.setup_presentation_size(slide_width_pixels, slide_height_pixels)

        for page_idx in sorted(pages_content.keys()):
            slide = builder.add_blank_slide()

            if background_images and page_idx < len(background_images):
                bg_image_path = background_images[page_idx]
                if bg_image_path and os.path.exists(bg_image_path):
                    try:
                        slide.shapes.add_picture(
                            bg_image_path,
                            left=0,
                            top=0,
                            width=builder.prs.slide_width,
                            height=builder.prs.slide_height,
                        )
                    except Exception as exc:
                        logger.error("Failed to add background image: %s", exc)
                else:
                    logger.warning("Background image missing for page %s", page_idx)

            page_items = pages_content[page_idx]

            text_items = []
            image_items = []
            for item in page_items:
                item_type = item.get("type", "")
                if item_type in ["text", "title", "header", "footer"]:
                    text_items.append(item)
                elif item_type in ["image", "table", "equation"]:
                    image_items.append(item)

            for img_item in image_items:
                ExportService._add_mineru_image_to_slide(
                    builder, slide, img_item, mineru_dir, scale_x, scale_y
                )

            for text_item in text_items:
                ExportService._add_mineru_text_to_slide(
                    builder, slide, text_item, scale_x, scale_y
                )

        try:
            if output_file:
                builder.save(output_file)
                return None

            pptx_bytes = io.BytesIO()
            builder.get_presentation().save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
        finally:
            if temp_dir and temp_dir.exists():
                shutil.rmtree(temp_dir, ignore_errors=True)

    @staticmethod
    def _add_mineru_text_to_slide(
        builder,
        slide,
        text_item: Dict[str, Any],
        scale_x: float = 1.0,
        scale_y: float = 1.0,
    ):
        text = text_item.get("text", "").strip()
        if not text:
            return

        bbox = text_item.get("bbox")
        if not bbox or len(bbox) != 4:
            logger.warning("Invalid bbox for text item: %s", text_item)
            return

        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y),
        ]

        item_type = text_item.get("type", "text")
        text_level = text_item.get("text_level")
        level = "title" if item_type == "title" or text_level == 1 else "default"

        try:
            builder.add_text_element(
                slide=slide,
                text=text,
                bbox=bbox,
                text_level=level,
            )
        except Exception as exc:
            logger.error("Failed to add text element: %s", exc)

    @staticmethod
    def _add_mineru_image_to_slide(
        builder,
        slide,
        image_item: Dict[str, Any],
        mineru_dir: Path,
        scale_x: float = 1.0,
        scale_y: float = 1.0,
    ):
        bbox = image_item.get("bbox")
        if not bbox or len(bbox) != 4:
            logger.warning("Invalid bbox for image item: %s", image_item)
            return

        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y),
        ]

        html_table = image_item.get("html_table") or image_item.get("table_body")
        item_type = image_item.get("type", "image")

        if html_table and item_type == "table":
            try:
                builder.add_table_element(
                    slide=slide,
                    html_table=html_table,
                    bbox=bbox,
                )
                return
            except Exception as exc:
                logger.error("Failed to add table: %s", exc)

        img_path_str = (
            image_item.get("img_path")
            or image_item.get("table_img_path")
            or image_item.get("equation_img_path")
            or ""
        )
        if not img_path_str:
            logger.warning("No image path in item: %s", image_item)
            return

        possible_paths = [
            Path(img_path_str),
            mineru_dir / img_path_str,
            mineru_dir / "images" / Path(img_path_str).name,
            mineru_dir / Path(img_path_str).name,
        ]

        image_path = None
        for path in possible_paths:
            if path.exists():
                image_path = str(path)
                break

        if not image_path:
            logger.warning("Image file not found: %s", img_path_str)
            builder.add_image_placeholder(slide, bbox)
            return

        try:
            builder.add_image_element(
                slide=slide,
                image_path=image_path,
                bbox=bbox,
            )
        except Exception as exc:
            logger.error("Failed to add image element: %s", exc)
