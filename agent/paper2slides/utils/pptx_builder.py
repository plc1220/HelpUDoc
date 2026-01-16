"""
PPTX builder utilities for MinerU outputs.
"""
from __future__ import annotations

import logging
import os
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

try:  # pragma: no cover - optional dependency
    from PIL import ImageFont
except Exception:  # pragma: no cover - optional dependency
    ImageFont = None


class HTMLTableParser(HTMLParser):
    """Parse HTML table into row/column data."""

    def __init__(self) -> None:
        super().__init__()
        self.table_data: List[List[str]] = []
        self.current_row: List[str] = []
        self.current_cell: List[str] = []
        self.in_table = False
        self.in_row = False
        self.in_cell = False

    def handle_starttag(self, tag, attrs):
        if tag == "table":
            self.in_table = True
            self.table_data = []
        elif tag == "tr":
            self.in_row = True
            self.current_row = []
        elif tag in ("td", "th"):
            self.in_cell = True
            self.current_cell = []

    def handle_endtag(self, tag):
        if tag == "table":
            self.in_table = False
        elif tag == "tr":
            self.in_row = False
            if self.current_row:
                self.table_data.append(self.current_row)
        elif tag in ("td", "th"):
            self.in_cell = False
            cell_text = "".join(self.current_cell).strip()
            self.current_row.append(cell_text)

    def handle_data(self, data):
        if self.in_cell:
            self.current_cell.append(data)

    @staticmethod
    def parse_html_table(html: str) -> List[List[str]]:
        parser = HTMLTableParser()
        parser.feed(html)
        return parser.table_data


class PPTXBuilder:
    """Builder class for creating editable PPTX files from structured content."""

    DEFAULT_SLIDE_WIDTH_INCHES = 10
    DEFAULT_SLIDE_HEIGHT_INCHES = 5.625
    DEFAULT_DPI = 96

    MAX_SLIDE_WIDTH_INCHES = 56.0
    MAX_SLIDE_HEIGHT_INCHES = 56.0
    MIN_SLIDE_WIDTH_INCHES = 1.0
    MIN_SLIDE_HEIGHT_INCHES = 1.0

    MIN_FONT_SIZE = 6
    MAX_FONT_SIZE = 200
    DEFAULT_FONT_NAME = "Calibri"

    FONT_PATH = os.path.join(
        os.path.dirname(__file__),
        "..",
        "fonts",
        "NotoSansSC-Regular.ttf",
    )

    _font_cache: Dict[float, Any] = {}

    @classmethod
    def _get_font(cls, size_pt: float) -> Optional[Any]:
        if ImageFont is None:
            return None

        cache_key = round(size_pt * 2) / 2
        if cache_key not in cls._font_cache:
            try:
                cls._font_cache[cache_key] = ImageFont.truetype(
                    cls.FONT_PATH, int(size_pt)
                )
            except Exception as exc:
                logger.warning("Failed to load font %s: %s", cls.FONT_PATH, exc)
                return None
        return cls._font_cache[cache_key]

    @classmethod
    def _measure_text_width(cls, text: str, font_size_pt: float) -> Optional[float]:
        font = cls._get_font(font_size_pt)
        if font is None:
            return None

        try:
            bbox = font.getbbox(text)
            width_px = bbox[2] - bbox[0]
            return width_px
        except Exception as exc:
            logger.warning("Failed to measure text: %s", exc)
            return None

    def __init__(
        self,
        slide_width_inches: float | None = None,
        slide_height_inches: float | None = None,
    ) -> None:
        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = (
            slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        )
        self.prs: Optional[Presentation] = None
        self.current_slide = None

    def create_presentation(self) -> Presentation:
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        return self.prs

    def setup_presentation_size(
        self, width_pixels: int, height_pixels: int, dpi: int | None = None
    ) -> None:
        dpi = dpi or self.DEFAULT_DPI

        width_inches = width_pixels / dpi
        height_inches = height_pixels / dpi

        scale_factor = 1.0

        if width_inches > self.MAX_SLIDE_WIDTH_INCHES:
            scale_factor = self.MAX_SLIDE_WIDTH_INCHES / width_inches
            logger.warning(
                "Slide width %.2f\" exceeds python-pptx limit (%.2f\"), scaling by %.3fx",
                width_inches,
                self.MAX_SLIDE_WIDTH_INCHES,
                scale_factor,
            )

        if height_inches > self.MAX_SLIDE_HEIGHT_INCHES:
            height_scale = self.MAX_SLIDE_HEIGHT_INCHES / height_inches
            if height_scale < scale_factor:
                scale_factor = height_scale
                logger.warning(
                    "Slide height %.2f\" exceeds python-pptx limit (%.2f\"), scaling by %.3fx",
                    height_inches,
                    self.MAX_SLIDE_HEIGHT_INCHES,
                    scale_factor,
                )

        if scale_factor < 1.0:
            width_inches *= scale_factor
            height_inches *= scale_factor
            logger.info(
                "Final slide dimensions after scaling: %.2f\" x %.2f\" (from %sx%s px @ %s DPI)",
                width_inches,
                height_inches,
                width_pixels,
                height_pixels,
                dpi,
            )

        width_inches = max(self.MIN_SLIDE_WIDTH_INCHES, width_inches)
        height_inches = max(self.MIN_SLIDE_HEIGHT_INCHES, height_inches)

        self.slide_width_inches = width_inches
        self.slide_height_inches = height_inches

        if self.prs:
            self.prs.slide_width = Inches(self.slide_width_inches)
            self.prs.slide_height = Inches(self.slide_height_inches)

    def add_blank_slide(self):
        if not self.prs:
            self.create_presentation()

        blank_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide

    def pixels_to_inches(self, pixels: float, dpi: int | None = None) -> float:
        dpi = dpi or self.DEFAULT_DPI
        return pixels / dpi

    def calculate_font_size(
        self,
        bbox: List[int],
        text: str,
        text_level: Any = None,
        dpi: int | None = None,
    ) -> float:
        dpi = dpi or self.DEFAULT_DPI

        width_px = bbox[2] - bbox[0]
        height_px = bbox[3] - bbox[1]

        width_pt = (width_px / dpi) * 72
        height_pt = (height_px / dpi) * 72

        usable_width_pt = width_pt
        usable_height_pt = height_pt

        if usable_width_pt <= 0 or usable_height_pt <= 0:
            logger.warning(
                "Bbox too small for text: %sx%s px, text: '%s...'",
                width_px,
                height_px,
                text[:30],
            )
            return self.MIN_FONT_SIZE

        line_height_ratio = 1.0
        use_precise = os.path.exists(self.FONT_PATH)
        best_size = self.MIN_FONT_SIZE

        for font_size in range(int(self.MAX_FONT_SIZE), int(self.MIN_FONT_SIZE) - 1, -1):
            font_size = float(font_size)

            lines = text.split("\n")
            total_required_lines = 0

            for line in lines:
                if not line:
                    total_required_lines += 1
                    continue

                if use_precise:
                    line_width_pt = self._measure_text_width(line, font_size)
                    if line_width_pt is None:
                        use_precise = False

                if not use_precise:
                    cjk_count = sum(
                        1
                        for c in line
                        if "\u4e00" <= c <= "\u9fff"
                        or "\u3040" <= c <= "\u30ff"
                        or "\uac00" <= c <= "\ud7af"
                    )
                    non_cjk_count = len(line) - cjk_count
                    line_width_pt = (cjk_count * 1.0 + non_cjk_count * 0.5) * font_size

                lines_needed = max(1, -(-int(line_width_pt) // int(usable_width_pt)))
                total_required_lines += lines_needed

            line_height_pt = font_size * line_height_ratio
            total_height_pt = total_required_lines * line_height_pt

            if total_height_pt <= usable_height_pt:
                best_size = font_size
                break

        if best_size == self.MIN_FONT_SIZE and len(text) > 3:
            logger.warning(
                "Text may overflow: '%s...' in bbox %sx%s px",
                text[:50],
                width_px,
                height_px,
            )

        logger.debug(
            "Font size calc: '%s%s' bbox=%sx%s px -> %s pt",
            text[:20],
            "..." if len(text) > 20 else "",
            width_px,
            height_px,
            best_size,
        )

        return best_size

    def add_text_element(
        self,
        slide,
        text: str,
        bbox: List[int],
        text_level: Any = None,
        dpi: int | None = None,
        align: str = "left",
        font_size: float | None = None,
        text_style: Any = None,
    ) -> None:
        dpi = dpi or self.DEFAULT_DPI

        has_colored_segments = (
            text_style
            and hasattr(text_style, "colored_segments")
            and text_style.colored_segments
            and len(text_style.colored_segments) > 0
        )

        if has_colored_segments:
            actual_text = "".join(seg.text for seg in text_style.colored_segments)
        else:
            actual_text = text

        expand_ratio = 0.01
        bbox_width = bbox[2] - bbox[0]
        bbox_height = bbox[3] - bbox[1]
        expand_w = bbox_width * expand_ratio
        expand_h = bbox_height * expand_ratio

        left = Inches(self.pixels_to_inches(bbox[0] - expand_w / 2, dpi))
        top = Inches(self.pixels_to_inches(bbox[1] - expand_h / 2, dpi))
        width = Inches(self.pixels_to_inches(bbox_width + expand_w, dpi))
        height = Inches(self.pixels_to_inches(bbox_height + expand_h, dpi))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)

        def replace_some_chars(value: str) -> str:
            if value.lstrip().startswith("·"):
                return value.replace("·", "•", 1)
            return value

        actual_text = replace_some_chars(actual_text)
        if font_size is None:
            font_size = self.calculate_font_size(bbox, actual_text, text_level, dpi)
        else:
            font_size = float(font_size)
            font_size = max(self.MIN_FONT_SIZE, min(self.MAX_FONT_SIZE, font_size))

        effective_align = align
        if text_style and getattr(text_style, "text_alignment", None):
            effective_align = text_style.text_alignment

        is_bold = False
        is_italic = False
        is_underline = False
        if text_style:
            is_bold = getattr(text_style, "is_bold", False)
            is_italic = getattr(text_style, "is_italic", False)
            is_underline = getattr(text_style, "is_underline", False)

        if text_level == 1 or text_level == "title":
            is_bold = True

        if has_colored_segments:
            paragraph = text_frame.paragraphs[0]
            paragraph.clear()
            for seg in text_style.colored_segments:
                run = paragraph.add_run()
                run.text = replace_some_chars(seg.text)
                run.font.size = Pt(font_size)
                run.font.bold = is_bold
                run.font.underline = is_underline
                r, g, b = seg.color_rgb
                run.font.color.rgb = RGBColor(r, g, b)
                run.font.name = self.DEFAULT_FONT_NAME
                if getattr(seg, "is_latex", False):
                    run.font.italic = True
                else:
                    run.font.italic = is_italic
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 1.0
        else:
            text_frame.text = actual_text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.bold = is_bold
                paragraph.font.italic = is_italic
                paragraph.font.underline = is_underline
                paragraph.font.name = self.DEFAULT_FONT_NAME
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                paragraph.line_spacing = 1.0

                if text_style and getattr(text_style, "font_color_rgb", None):
                    r, g, b = text_style.font_color_rgb
                    paragraph.font.color.rgb = RGBColor(r, g, b)

                for run in paragraph.runs:
                    run.font.name = self.DEFAULT_FONT_NAME

        paragraphs = text_frame.paragraphs if not has_colored_segments else [paragraph]
        for para in paragraphs:
            if effective_align == "center":
                para.alignment = PP_ALIGN.CENTER
            elif effective_align == "right":
                para.alignment = PP_ALIGN.RIGHT
            elif effective_align == "justify":
                para.alignment = PP_ALIGN.JUSTIFY
            else:
                para.alignment = PP_ALIGN.LEFT

        logger.debug(
            "Text: '%s' | box: %sx%s px | font: %.1f pt | chars: %s",
            actual_text[:35],
            bbox_width,
            bbox_height,
            font_size,
            len(actual_text),
        )

    def add_image_element(
        self, slide, image_path: str, bbox: List[int], dpi: int | None = None
    ) -> None:
        dpi = dpi or self.DEFAULT_DPI

        if not os.path.exists(image_path):
            logger.warning("Image not found: %s, adding placeholder", image_path)
            self.add_image_placeholder(slide, bbox, dpi)
            return

        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))

        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as exc:
            logger.error("Failed to add image %s: %s", image_path, exc)
            self.add_image_placeholder(slide, bbox, dpi)

    def add_image_placeholder(
        self, slide, bbox: List[int], dpi: int | None = None
    ) -> None:
        dpi = dpi or self.DEFAULT_DPI

        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = RGBColor(203, 213, 225)
        shape.text_frame.text = "Image not found"

    def add_table_element(
        self, slide, html_table: Any, bbox: List[int], dpi: int | None = None
    ) -> None:
        dpi = dpi or self.DEFAULT_DPI

        if isinstance(html_table, list):
            table_data = [[str(cell) for cell in row] for row in html_table]
        else:
            try:
                table_data = HTMLTableParser.parse_html_table(html_table)
            except Exception as exc:
                logger.error("Failed to parse HTML table: %s", exc)
                return

        if not table_data or not table_data[0]:
            logger.warning("Empty table data")
            return

        rows = len(table_data)
        cols = len(table_data[0])

        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))

        try:
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            cell_height_px = (bbox[3] - bbox[1]) / rows
            font_size = min(18, max(8, cell_height_px * 0.3))

            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_text
                        text_frame = cell.text_frame
                        text_frame.word_wrap = True
                        for paragraph in text_frame.paragraphs:
                            paragraph.font.size = Pt(font_size)
                            paragraph.font.name = self.DEFAULT_FONT_NAME
                            paragraph.space_before = Pt(0)
                            paragraph.space_after = Pt(0)
                            paragraph.line_spacing = 1.0
                            paragraph.alignment = PP_ALIGN.CENTER
                            if row_idx == 0:
                                paragraph.font.bold = True
        except Exception as exc:
            logger.error("Failed to create table: %s", exc)

    def save(self, output_path: str) -> None:
        if not self.prs:
            raise ValueError("No presentation to save")

        output_path_obj = Path(output_path)
        output_dir = output_path_obj.parent
        if str(output_dir) != ".":
            output_dir.mkdir(parents=True, exist_ok=True)

        self.prs.save(output_path)
        logger.info("Saved presentation to: %s", output_path)

    def get_presentation(self) -> Presentation:
        return self.prs
