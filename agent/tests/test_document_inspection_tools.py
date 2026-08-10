from __future__ import annotations

import json
import re
import zipfile
from pathlib import Path

from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from helpudoc_agent.state import WorkspaceState
from helpudoc_agent.tools.workspace.builtins.document_inspection import (
    build_document_inspection_tools,
)
from helpudoc_agent.tools.workspace.builtins.knowledge_navigation import (
    build_knowledge_navigation_tools,
)


def _strip_dimension_elements(path: Path, only: set[str] | None = None) -> None:
    """Remove ``<dimension/>`` like Google Drive/Sheets exports do."""
    with zipfile.ZipFile(path) as source:
        entries = [(item, source.read(item.filename)) for item in source.infolist()]
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as target:
        for item, data in entries:
            name = item.filename
            if (
                name.startswith("xl/worksheets/")
                and name.endswith(".xml")
                and (only is None or name in only)
            ):
                data = re.sub(rb"<dimension[^>]*/>", b"", data)
            target.writestr(item, data)


def _set_pptx_alt_text(shape, *, title: str, description: str) -> None:
    for child in shape._element.iter():
        if str(child.tag).rsplit("}", 1)[-1] == "cNvPr":
            child.set("title", title)
            child.set("descr", description)
            return
    raise AssertionError("shape has no cNvPr node")


def _create_inspection_deck(path: Path) -> dict[str, int]:
    presentation = Presentation()
    presentation.core_properties.title = "Launch review"
    presentation.core_properties.author = "HelpUDoc QA"

    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Malaysia launch"
    body = first.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5), Inches(1))
    body.text_frame.text = "Revenue grows by thirty percent."
    table_shape = first.shapes.add_table(2, 2, Inches(0.8), Inches(3), Inches(5), Inches(1.5))
    table_shape.table.cell(0, 0).text = "Market"
    table_shape.table.cell(0, 1).text = "Owner"
    table_shape.table.cell(1, 0).text = "Malaysia"
    table_shape.table.cell(1, 1).text = "Aisha"
    first.notes_slide.notes_text_frame.text = (
        "Speaker-only launch date is 15 September. " + ("context " * 180) + "notes-tail-token"
    )

    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    visual = second.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(3),
        Inches(2),
    )
    visual.text = "Roadmap visual"
    _set_pptx_alt_text(
        visual,
        title="Accessible roadmap",
        description=(
            "Confidential roadmap illustration for enterprise rollout "
            + ("context " * 180)
            + "alt-tail-token"
        ),
    )
    presentation.save(path)
    return {"body": body.shape_id, "table": table_shape.shape_id, "visual": visual.shape_id}


def test_document_tools_search_and_inspect_original_workbook(tmp_path) -> None:
    workbook_path = tmp_path / "forecast.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Forecast"
    sheet.append(["Region", "Revenue"])
    sheet.append(["Malaysia", 125000])
    sheet.append(["Singapore", 98000])
    workbook.save(workbook_path)

    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    inventory = json.loads(tools["inspect_document"].invoke({"file_path": "forecast.xlsx"}))
    assert inventory["file"] == "/forecast.xlsx"
    assert inventory["sheets"][0]["name"] == "Forecast"

    matches = json.loads(
        tools["search_document"].invoke(
            {"file_path": "forecast.xlsx", "query": "Malaysia"},
        )
    )
    assert matches["results"][0]["location"] == "sheet:Forecast:cell:A2"

    cells = json.loads(
        tools["inspect_document"].invoke(
            {
                "file_path": "forecast.xlsx",
                "sheet_name": "Forecast",
                "cell_range": "A1:B2",
            },
        )
    )
    assert cells["cells"][1][1]["value"] == 125000
    single_cell = json.loads(
        tools["inspect_document"].invoke(
            {
                "file_path": "forecast.xlsx",
                "sheet_name": "Forecast",
                "cell_range": "B2",
            },
        )
    )
    assert single_cell["cells"][0][0]["value"] == 125000

    oversized = json.loads(
        tools["inspect_document"].invoke(
            {
                "file_path": "forecast.xlsx",
                "sheet_name": "Forecast",
                "cell_range": "A1:XFD1048576",
            },
        )
    )
    assert oversized["status"] == "error"
    assert oversized["errorCode"] == "RANGE_TOO_LARGE"
    assert oversized["retryable"] is False
    assert oversized["message"].startswith(
        "Document inspection failed: Spreadsheet inspection range is too large"
    )
    assert oversized["suggestedNextCall"]


def test_document_tools_read_unsized_google_drive_export(tmp_path) -> None:
    workbook_path = tmp_path / "drive-export.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Region", "Revenue"])
    sheet.append(["Malaysia", 125000])
    notes = workbook.create_sheet("Notes")
    notes["A1"] = "prepared by finance"
    workbook.save(workbook_path)
    _strip_dimension_elements(workbook_path)

    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    # A targeted range read must not depend on dimension discovery.
    targeted = json.loads(
        tools["inspect_document"].invoke(
            {
                "file_path": "drive-export.xlsx",
                "sheet_name": "Data",
                "cell_range": "A1:C3",
            }
        )
    )
    assert targeted["status"] == "ok"
    assert targeted["selectedRange"] == "A1:C3"
    assert targeted["cells"][0][0]["cell"] == "A1"
    assert targeted["cells"][1][1]["value"] == 125000
    # Padded cells past the stored data stay addressable instead of raising.
    assert targeted["cells"][0][2]["cell"] == "C1"
    assert targeted["cells"][0][2]["value"] is None
    data_entry = next(item for item in targeted["sheets"] if item["name"] == "Data")
    assert data_entry["dimensions"] is None
    assert "unsized" in data_entry["dimensionsError"].lower()
    assert {item["name"] for item in targeted["sheets"]} == {"Data", "Notes"}

    # The inventory-only call recovers dimensions through the force fallback.
    inventory = json.loads(
        tools["inspect_document"].invoke({"file_path": "drive-export.xlsx"})
    )
    assert inventory["status"] == "ok"
    recovered = next(item for item in inventory["sheets"] if item["name"] == "Data")
    assert recovered["dimensions"] == "A1:B2"
    assert "dimensionsError" not in recovered
    assert recovered["maxRow"] == 2
    assert recovered["maxColumn"] == 2

    matches = json.loads(
        tools["search_document"].invoke(
            {"file_path": "drive-export.xlsx", "query": "Malaysia"}
        )
    )
    assert matches["status"] == "ok"
    assert matches["results"][0]["location"] == "sheet:Data:cell:A2"


def test_document_tools_preserve_partial_sheet_inventory(tmp_path) -> None:
    workbook_path = tmp_path / "mixed.xlsx"
    workbook = Workbook()
    sized = workbook.active
    sized.title = "Sized"
    sized.append(["a", "b"])
    unsized = workbook.create_sheet("Unsized")
    unsized["A1"] = "value"
    workbook.save(workbook_path)
    _strip_dimension_elements(workbook_path, only={"xl/worksheets/sheet2.xml"})

    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    payload = json.loads(
        tools["inspect_document"].invoke(
            {
                "file_path": "mixed.xlsx",
                "sheet_name": "Sized",
                "cell_range": "A1:B1",
            }
        )
    )

    assert payload["status"] == "ok"
    assert payload["cells"][0][0]["value"] == "a"
    entries = {item["name"]: item for item in payload["sheets"]}
    assert entries["Sized"]["dimensions"] == "A1:B1"
    assert "dimensionsError" not in entries["Sized"]
    assert entries["Unsized"]["dimensions"] is None
    assert entries["Unsized"]["dimensionsError"]


def test_document_tools_return_structured_error_envelopes(tmp_path) -> None:
    workbook_path = tmp_path / "book.xlsx"
    workbook = Workbook()
    workbook.active.title = "Only"
    workbook.active["A1"] = 1
    workbook.save(workbook_path)
    (tmp_path / "notes.rtf").write_text("unsupported", encoding="utf-8")

    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    missing = json.loads(tools["inspect_document"].invoke({"file_path": "nope.xlsx"}))
    assert missing["status"] == "error"
    assert missing["errorCode"] == "FILE_NOT_FOUND"
    assert missing["retryable"] is False
    assert missing["tool"] == "inspect_document"
    assert missing["suggestedNextCall"]

    unknown_sheet = json.loads(
        tools["inspect_document"].invoke(
            {"file_path": "book.xlsx", "sheet_name": "Missing"}
        )
    )
    assert unknown_sheet["errorCode"] == "UNKNOWN_SHEET"
    assert "Only" in unknown_sheet["message"]

    bad_range = json.loads(
        tools["inspect_document"].invoke(
            {"file_path": "book.xlsx", "sheet_name": "Only", "cell_range": "not-a-range"}
        )
    )
    assert bad_range["errorCode"] == "INVALID_RANGE"

    unsupported = json.loads(tools["inspect_document"].invoke({"file_path": "notes.rtf"}))
    assert unsupported["errorCode"] == "UNSUPPORTED_DOCUMENT_TYPE"

    no_query = json.loads(
        tools["search_document"].invoke({"file_path": "book.xlsx", "query": "  "})
    )
    assert no_query["status"] == "error"
    assert no_query["errorCode"] == "MISSING_QUERY"
    assert no_query["tool"] == "search_document"
    assert no_query["message"].startswith("Document search failed: query is required")


def test_document_tools_keep_docx_search_locations_addressable(tmp_path) -> None:
    document_path = tmp_path / "policy.docx"
    document = Document()
    document.add_paragraph("")
    document.add_heading("Returns", level=1)
    document.add_paragraph("Customers have 30 days to return an item.")
    document.save(document_path)

    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}
    matches = json.loads(
        tools["search_document"].invoke({"file_path": "policy.docx", "query": "30 days"})
    )
    paragraph = matches["results"][0]["paragraph"]

    inspected = json.loads(
        tools["inspect_document"].invoke(
            {
                "file_path": "policy.docx",
                "item_start": paragraph,
                "item_end": paragraph,
            }
        )
    )
    assert inspected["selectedItems"][0]["index"] == paragraph
    assert "30 days" in inspected["selectedItems"][0]["text"]


def test_document_tools_inspect_bounded_pptx_slides_and_deck_metadata(tmp_path) -> None:
    deck_path = tmp_path / "launch.pptx"
    shape_ids = _create_inspection_deck(deck_path)
    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    inspected = json.loads(
        tools["inspect_document"].invoke(
            {"file_path": "launch.pptx", "slide_start": 1, "slide_end": 1}
        )
    )

    assert inspected["status"] == "ok"
    assert inspected["kind"] == "pptx"
    assert inspected["file"] == "/launch.pptx"
    assert inspected["slideCount"] == 2
    assert inspected["selectedSlides"] == [1, 1]
    assert inspected["metadata"]["title"] == "Launch review"
    assert inspected["slideSize"]["widthInches"] > inspected["slideSize"]["heightInches"]
    assert [entry["slide"] for entry in inspected["slideInventory"]] == [1, 2]
    assert inspected["slideInventory"][0]["title"] == "Malaysia launch"
    assert len(inspected["slides"]) == 1

    selected = inspected["slides"][0]
    body = next(item for item in selected["shapes"] if item["shapeId"] == str(shape_ids["body"]))
    assert body["paragraphs"][0]["location"] == (
        f"slide:1:shape:{shape_ids['body']}:paragraph:1"
    )
    table = next(item for item in selected["shapes"] if item["shapeId"] == str(shape_ids["table"]))
    assert table["table"]["rows"][1]["location"] == (
        f"slide:1:shape:{shape_ids['table']}:table:row:2"
    )
    assert table["table"]["rows"][1]["values"] == ["Malaysia", "Aisha"]
    assert selected["speakerNotes"][0]["location"] == "slide:1:note:paragraph:1"
    assert "15 September" in selected["speakerNotes"][0]["text"]
    assert inspected["visualInspection"]["rendered"] is False
    assert inspected["visualInspection"]["renderRequiredForVisualQuestions"] is True
    assert inspected["visualInspection"]["slides"] == [1]


def test_document_tools_search_pptx_text_tables_notes_and_alt_text(tmp_path) -> None:
    deck_path = tmp_path / "launch.pptx"
    shape_ids = _create_inspection_deck(deck_path)
    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    text = json.loads(
        tools["search_document"].invoke({"file_path": "launch.pptx", "query": "thirty percent"})
    )
    assert text["status"] == "ok"
    assert text["results"][0]["kind"] == "slideText"
    assert text["results"][0]["location"] == f"slide:1:shape:{shape_ids['body']}:paragraph:1"

    table = json.loads(
        tools["search_document"].invoke({"file_path": "launch.pptx", "query": "Malaysia Aisha"})
    )
    assert table["results"][0]["kind"] == "tableRow"
    assert table["results"][0]["location"] == f"slide:1:shape:{shape_ids['table']}:table:row:2"

    notes = json.loads(
        tools["search_document"].invoke({"file_path": "launch.pptx", "query": "15 September"})
    )
    assert notes["results"][0]["kind"] == "speakerNote"
    assert notes["results"][0]["location"] == "slide:1:note:paragraph:1"
    notes_tail = json.loads(
        tools["search_document"].invoke({"file_path": "launch.pptx", "query": "notes-tail-token"})
    )
    assert notes_tail["results"][0]["kind"] == "speakerNote"

    alt_text = json.loads(
        tools["search_document"].invoke(
            {"file_path": "launch.pptx", "query": "confidential enterprise rollout"}
        )
    )
    assert alt_text["results"][0]["kind"] == "altText"
    assert alt_text["results"][0]["location"] == f"slide:2:shape:{shape_ids['visual']}:alt-text"
    assert alt_text["results"][0]["altTextFields"] == ["title", "description"]
    assert "Confidential roadmap" in alt_text["results"][0]["snippet"]
    alt_tail = json.loads(
        tools["search_document"].invoke({"file_path": "launch.pptx", "query": "alt-tail-token"})
    )
    assert alt_tail["results"][0]["kind"] == "altText"


def test_document_tools_bound_pptx_slide_ranges_and_describe_support(tmp_path, monkeypatch) -> None:
    from helpudoc_agent.tools.workspace.builtins import document_inspection

    deck_path = tmp_path / "launch.pptx"
    _create_inspection_deck(deck_path)
    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_document_inspection_tools(workspace)}

    past_end = json.loads(
        tools["inspect_document"].invoke(
            {"file_path": "launch.pptx", "slide_start": 3, "slide_end": 3}
        )
    )
    assert past_end["status"] == "error"
    assert past_end["errorCode"] == "INVALID_RANGE"
    assert "2 slides" in past_end["message"]
    assert "slide_start=1" in past_end["suggestedNextCall"]

    monkeypatch.setattr(document_inspection, "_MAX_PPTX_INSPECT_SLIDES", 1)
    oversized = json.loads(
        tools["inspect_document"].invoke(
            {"file_path": "launch.pptx", "slide_start": 1, "slide_end": 2}
        )
    )
    assert oversized["status"] == "error"
    assert oversized["errorCode"] == "RANGE_TOO_LARGE"
    assert "at most 1 slides" in oversized["message"]
    assert "slide_end=1" in oversized["suggestedNextCall"]
    assert "PPTX" in tools["inspect_document"].description
    assert "speaker-note" in tools["search_document"].description

    monkeypatch.setattr(document_inspection, "_MAX_PPTX_SCAN_UNITS", 1)
    truncated = json.loads(
        tools["search_document"].invoke(
            {"file_path": "launch.pptx", "query": "not present anywhere"}
        )
    )
    assert truncated["status"] == "ok"
    assert truncated["results"] == []
    assert truncated["scanTruncated"] is True

    monkeypatch.setattr(document_inspection, "_MAX_PPTX_INSPECT_SLIDES", 20)
    monkeypatch.setattr(document_inspection, "_MAX_OUTPUT_CHARS", 3000)
    bounded_raw = tools["inspect_document"].invoke({"file_path": "launch.pptx"})
    bounded = json.loads(bounded_raw)
    assert bounded["status"] == "ok"
    assert bounded["outputTruncated"] is True
    assert len(bounded_raw) < 4000


def test_document_tools_reject_paths_outside_workspace(tmp_path) -> None:
    workspace = WorkspaceState(workspace_id="docs", root_path=tmp_path / "workspace")
    tool = {tool.name: tool for tool in build_document_inspection_tools(workspace)}["inspect_document"]

    result = json.loads(tool.invoke({"file_path": "../secret.txt"}))

    assert result["status"] == "error"
    assert result["errorCode"] == "PATH_OUTSIDE_WORKSPACE"
    assert result["retryable"] is False
    assert result["message"].startswith(
        "Document inspection failed: Path must remain inside the workspace"
    )


def test_knowledge_tools_navigate_published_okf_bundle(tmp_path) -> None:
    bundle = tmp_path / ".system" / "knowledge" / "7"
    concepts = bundle / "concepts"
    concepts.mkdir(parents=True)
    (bundle / "index.md").write_text(
        "---\nokf_version: \"0.2\"\n---\n\n# Product Handbook\n\n"
        "* [Returns](concepts/returns.md) - Return policy\n",
        encoding="utf-8",
    )
    (concepts / "returns.md").write_text(
        "---\ntitle: \"Returns\"\n---\n\n# Returns\n\nCustomers have 30 days.\n",
        encoding="utf-8",
    )

    workspace = WorkspaceState(workspace_id="knowledge", root_path=tmp_path)
    workspace.context["knowledge_refs"] = [
        {"id": 7, "bundleRoot": str(bundle), "snapshotHash": "snapshot-7"}
    ]
    tools = {tool.name: tool for tool in build_knowledge_navigation_tools(workspace)}

    listing = json.loads(tools["knowledge_search"].invoke({"query": ""}))
    assert listing["results"][0]["path"] == "knowledge://7/index.md"

    search = json.loads(tools["knowledge_search"].invoke({"query": "30 days"}))
    assert search["results"][0]["path"] == "knowledge://7/concepts/returns.md"
    assert search["results"][0]["line"] == 7

    content = json.loads(
        tools["knowledge_read"].invoke(
            {
                "path": "knowledge://7/concepts/returns.md",
                "start_line": 7,
                "end_line": 7,
            }
        )
    )
    assert content["content"] == "Customers have 30 days."
    assert content["selectedLines"] == [7, 7]


def test_knowledge_tools_require_an_explicit_tagged_bundle(tmp_path) -> None:
    bundle = tmp_path / "standalone-bundle"
    bundle.mkdir()
    (bundle / "index.md").write_text("# Secret Knowledge\n", encoding="utf-8")
    workspace = WorkspaceState(workspace_id="knowledge", root_path=tmp_path)
    tools = {tool.name: tool for tool in build_knowledge_navigation_tools(workspace)}

    search = json.loads(tools["knowledge_search"].invoke({"query": "Secret"}))
    assert search["resultCount"] == 0
    assert "@ command" in search["message"]
    assert tools["knowledge_read"].invoke({"path": "knowledge://7/index.md"}).startswith(
        "Knowledge read failed: Knowledge bundle was not tagged"
    )
