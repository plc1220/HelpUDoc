"""Attachment parsing and Gemini-backed understanding helpers."""
from __future__ import annotations

import base64
import json
import logging
import mimetypes
import os
import re
import tempfile
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


def _env_flag(name: str, default: bool) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _env_int(name: str, default: int) -> int:
    raw = os.getenv(name)
    if raw is None or not raw.strip():
        return default
    try:
        return int(raw.strip())
    except (TypeError, ValueError):
        return default


def _docling_enabled() -> bool:
    return _env_flag("DOCLING_ENABLED", True)


def _docling_force() -> bool:
    return _env_flag("DOCLING_FORCE", False)


def _docling_generate_images() -> bool:
    return _env_flag("DOCLING_GENERATE_IMAGES", False)


def _docling_max_pages() -> int:
    return _env_int("DOCLING_MAX_PAGES", 15)


def _docling_max_bytes() -> int:
    return _env_int("DOCLING_MAX_BYTES", 5 * 1024 * 1024)


def docling_timeout_s() -> int:
    return _env_int("DOCLING_TIMEOUT_S", 600)


def docling_parse_concurrency() -> int:
    return max(1, _env_int("DOCLING_PARSE_CONCURRENCY", 1))


def _extract_json_block(text: str) -> str:
    fenced = re.search(r"```json\s*(.*?)```", text or "", flags=re.IGNORECASE | re.DOTALL)
    if fenced:
        return fenced.group(1).strip()
    return (text or "").strip()


def _response_text(response: Any) -> str:
    text = getattr(response, "text", None)
    if isinstance(text, str) and text.strip():
        return text.strip()
    parts: List[str] = []
    for candidate in getattr(response, "candidates", []) or []:
        content = getattr(candidate, "content", None)
        if not content:
            continue
        for part in getattr(content, "parts", []) or []:
            value = getattr(part, "text", None)
            if isinstance(value, str) and value.strip():
                parts.append(value.strip())
    return "\n".join(parts).strip()


def _lc_ai_message_text(ai_message: Any) -> str:
    text = getattr(ai_message, "text", None)
    if isinstance(text, str) and text.strip():
        return text.strip()
    content = getattr(ai_message, "content", None)
    if isinstance(content, str) and content.strip():
        return content.strip()
    if isinstance(content, list):
        parts: List[str] = []
        for block in content:
            if isinstance(block, dict) and block.get("type") == "text":
                fragment = block.get("text")
                if isinstance(fragment, str) and fragment.strip():
                    parts.append(fragment.strip())
            elif isinstance(block, str) and block.strip():
                parts.append(block.strip())
        if parts:
            return "\n".join(parts).strip()
    return ""


def _split_sections_from_markdown(markdown: str) -> List[Dict[str, str]]:
    sections: List[Dict[str, str]] = []
    heading = "Overview"
    body_lines: List[str] = []
    for line in (markdown or "").splitlines():
        match = re.match(r"^#{1,6}\s+(.*)$", line.strip())
        if match:
            if body_lines:
                sections.append({"heading": heading, "body": "\n".join(body_lines).strip()})
            heading = match.group(1).strip() or "Section"
            body_lines = []
            continue
        body_lines.append(line)
    if body_lines:
        sections.append({"heading": heading, "body": "\n".join(body_lines).strip()})
    return [section for section in sections if section.get("body")]


def _markdown_from_attachment_payload(payload: Dict[str, Any]) -> str:
    title = str(payload.get("title") or "Attachment").strip() or "Attachment"
    summary = str(payload.get("summary") or "").strip()
    outline = [str(item).strip() for item in (payload.get("outline") or []) if str(item).strip()]
    raw_sections = payload.get("sections") or []
    lines: List[str] = [f"# {title}", ""]
    if summary:
        lines.extend(["## Summary", "", summary, ""])
    if outline:
        lines.extend(["## Outline", ""])
        lines.extend(f"- {item}" for item in outline)
        lines.append("")
    sections_added = False
    if isinstance(raw_sections, list):
        for item in raw_sections:
            if not isinstance(item, dict):
                continue
            heading = str(item.get("heading") or "").strip() or "Section"
            body = str(item.get("body") or "").strip()
            if not body:
                continue
            lines.extend([f"## {heading}", "", body, ""])
            sections_added = True
    normalized_body = str(payload.get("normalizedBody") or "").strip()
    if normalized_body and not sections_added:
        lines.extend(["## Details", "", normalized_body, ""])
    return "\n".join(lines).strip()


def _build_partial_attachment_payload(file_name: str, text: str, *, effective_mode: str) -> Dict[str, Any]:
    cleaned_lines = [line.strip() for line in (text or "").splitlines() if line.strip()]
    summary = cleaned_lines[0] if cleaned_lines else f"Extracted content from {file_name}"
    excerpt = "\n".join(cleaned_lines[:80]).strip() or f"Unable to extract detailed content from {file_name}."
    sections = [{"heading": "Source Excerpt", "body": excerpt}]
    payload = {
        "title": file_name,
        "summary": summary,
        "outline": ["Source Excerpt"],
        "sections": sections,
        "effectiveMode": effective_mode,
        "status": "partial",
    }
    payload["markdown"] = _markdown_from_attachment_payload(payload)
    return payload


def _docling_available() -> bool:
    try:
        __import__("docling")
        __import__("docling.document_converter")
    except Exception:
        return False
    return True


def _docling_markdown_to_payload(file_name: str, markdown: str, *, effective_mode: str = "parser") -> Dict[str, Any]:
    cleaned_markdown = (markdown or "").strip()
    sections = _split_sections_from_markdown(cleaned_markdown)
    title = file_name
    title_match = re.search(r"^#\s+(.+)$", cleaned_markdown, flags=re.MULTILINE)
    if title_match and title_match.group(1).strip():
        title = title_match.group(1).strip()
    summary = ""
    for section in sections:
        body = str(section.get("body") or "").strip()
        if body:
            summary = body.splitlines()[0].strip()
            if summary:
                break
    outline = [str(section.get("heading") or "").strip() for section in sections if str(section.get("heading") or "").strip()]
    payload = {
        "title": title,
        "summary": summary or f"Parsed content from {file_name}",
        "outline": outline,
        "sections": sections,
        "normalizedBody": cleaned_markdown,
        "effectiveMode": effective_mode,
        "status": "ready",
    }
    payload["markdown"] = cleaned_markdown or _markdown_from_attachment_payload(payload)
    return payload


def _pdf_page_count(buffer: bytes) -> Optional[int]:
    try:
        from pypdf import PdfReader  # type: ignore

        return len(PdfReader(BytesIO(buffer)).pages)
    except Exception:
        return None


def pdf_has_extractable_text(buffer: bytes, *, min_chars: int = 32, sample_pages: int = 5) -> bool:
    """Return True if the PDF has a usable text layer (so OCR can be skipped).

    Errs on the side of True (skip OCR) when uncertain, since OCR is the heavy path.
    """
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(BytesIO(buffer))
        total = 0
        for page in reader.pages[:sample_pages]:
            total += len((page.extract_text() or "").strip())
            if total >= min_chars:
                return True
        return total >= min_chars
    except Exception:
        return True


def should_use_docling(kind: str, buffer: bytes, *, force: bool = False) -> bool:
    """Gate the heavy Docling path: only small docs (or explicit opt-in) use it."""
    if not _docling_enabled() or not _docling_available():
        return False
    if force or _docling_force():
        return True
    if len(buffer) > _docling_max_bytes():
        return False
    if kind == "pdf":
        pages = _pdf_page_count(buffer)
        if pages is not None and pages > _docling_max_pages():
            return False
    return True


@lru_cache(maxsize=8)
def _get_docling_converter(is_pdf: bool, do_ocr: bool, generate_images: bool):
    """Build (and cache) a DocumentConverter per option-set to avoid per-request model reloads."""
    from docling.document_converter import DocumentConverter

    converter_kwargs: Dict[str, Any] = {}
    if is_pdf:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.document_converter import PdfFormatOption

            pdf_options = PdfPipelineOptions()
            pdf_options.do_ocr = do_ocr
            pdf_options.generate_picture_images = generate_images
            converter_kwargs["format_options"] = {
                InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_options)
            }
        except Exception:
            logger.debug("Docling PDF pipeline options unavailable; continuing with default conversion", exc_info=True)

    return DocumentConverter(**converter_kwargs)


def _extract_docling_payload(file_name: str, mime_type: str, buffer: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    suffix = Path(file_name).suffix.lower()
    guessed_suffix = mimetypes.guess_extension(mime_type or "") or ""
    if not suffix and guessed_suffix:
        suffix = guessed_suffix.lower()
    if not suffix:
        suffix = ".bin"
    is_pdf = suffix == ".pdf"
    # OCR auto by text layer: skip OCR (the heavy path) for digital PDFs; enable only for scanned.
    do_ocr = (not pdf_has_extractable_text(buffer)) if is_pdf else False
    if is_pdf:
        logger.info("Docling parsing %s (do_ocr=%s, generate_images=%s)", file_name, do_ocr, _docling_generate_images())
    with tempfile.TemporaryDirectory(prefix="helpudoc-docling-") as temp_dir:
        temp_root = Path(temp_dir)
        temp_input = temp_root / f"source{suffix}"
        temp_input.write_bytes(buffer)
        output_dir = temp_root / "parsed"
        output_dir.mkdir(parents=True, exist_ok=True)

        converter = _get_docling_converter(is_pdf, do_ocr, _docling_generate_images())
        conversion = converter.convert(temp_input)
        document = conversion.document
        md_path = output_dir / f"{temp_input.stem}.md"
        try:
            from docling_core.types.doc import ImageRefMode

            document.save_as_markdown(md_path, image_mode=ImageRefMode.REFERENCED)
            markdown = md_path.read_text(encoding="utf-8", errors="replace").strip()
        except Exception:
            markdown = str(document.export_to_markdown() or "").strip()
            md_path.write_text(markdown, encoding="utf-8")
        if not markdown:
            raise RuntimeError(f"Docling produced empty markdown for {file_name}")

        extracted_assets: List[Dict[str, Any]] = []
        for asset_path in sorted(output_dir.rglob("*")):
            if not asset_path.is_file() or asset_path == md_path:
                continue
            mime = mimetypes.guess_type(asset_path.name)[0] or ""
            if not mime.startswith("image/"):
                continue
            try:
                source_path = str(asset_path.relative_to(output_dir)).replace("\\", "/")
            except Exception:
                source_path = asset_path.name
            extracted_assets.append(
                {
                    "name": asset_path.name,
                    "mimeType": mime,
                    "contentB64": base64.b64encode(asset_path.read_bytes()).decode("ascii"),
                    "sourcePath": source_path,
                    "caption": None,
                    "footnote": None,
                }
            )
        return markdown, extracted_assets


def _extract_text_from_docx(buffer: bytes) -> str:
    from docx import Document  # type: ignore

    document = Document(BytesIO(buffer))
    chunks: List[str] = []
    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            chunks.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                chunks.append(" | ".join(cells))
    return "\n".join(chunks).strip()


def _extract_text_from_pptx(buffer: bytes) -> str:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(BytesIO(buffer))
    chunks: List[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: List[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                slide_lines.append(text.strip())
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                notes_text = ""
            if notes_text and notes_text.strip():
                slide_lines.append(f"Notes: {notes_text.strip()}")
        chunks.append("\n".join(slide_lines))
    return "\n\n".join(chunks).strip()


def _extract_text_from_pdf(buffer: bytes) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(BytesIO(buffer))
    chunks: List[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            chunks.append(f"Page {index}\n{text}")
    return "\n\n".join(chunks).strip()


def _guess_attachment_strategy(file_name: str, mime_type: str) -> Tuple[str, str]:
    suffix = Path(file_name).suffix.lower()
    normalized_mime = (mime_type or "").lower()
    if normalized_mime.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}:
        return "image", "part"
    if suffix == ".docx" or "wordprocessingml" in normalized_mime:
        return "docx", "parser"
    if suffix == ".pptx" or "presentationml" in normalized_mime:
        return "pptx", "parser"
    if suffix == ".pdf" or normalized_mime == "application/pdf":
        return "pdf", "parser"
    return "text", "part"


def _attachment_understanding_prompt(file_name: str, extracted_text: str | None, kind: str) -> str:
    base = (
        "Return strict JSON only with keys: "
        "title, summary, outline, sections, normalizedBody, effectiveMode, status. "
        "sections must be an array of objects with heading and body. "
        "effectiveMode must be one of part, parser, hybrid. "
        "status must be ready or partial. "
        "Preserve the document structure and factual detail instead of over-summarizing. "
        f"File name: {file_name}. "
        f"Kind: {kind}."
    )
    if extracted_text is None:
        return (
            f"{base} Understand the attached file directly. "
            "If details are uncertain, say so briefly in the summary instead of inventing content."
        )
    excerpt = extracted_text[:30000]
    return f"{base}\n\nSource content:\n{excerpt}"


def _parse_attachment_payload(raw_text: str, file_name: str, *, fallback_text: str, fallback_mode: str) -> Dict[str, Any]:
    try:
        parsed = json.loads(_extract_json_block(raw_text))
    except Exception:
        return _build_partial_attachment_payload(file_name, fallback_text, effective_mode=fallback_mode)
    if not isinstance(parsed, dict):
        return _build_partial_attachment_payload(file_name, fallback_text, effective_mode=fallback_mode)
    payload = {
        "title": str(parsed.get("title") or file_name).strip() or file_name,
        "summary": str(parsed.get("summary") or "").strip(),
        "outline": [str(item).strip() for item in (parsed.get("outline") or []) if str(item).strip()],
        "sections": [
            {
                "heading": str(item.get("heading") or "").strip() or "Section",
                "body": str(item.get("body") or "").strip(),
            }
            for item in (parsed.get("sections") or [])
            if isinstance(item, dict) and str(item.get("body") or "").strip()
        ],
        "normalizedBody": str(parsed.get("normalizedBody") or "").strip(),
        "effectiveMode": str(parsed.get("effectiveMode") or fallback_mode).strip() or fallback_mode,
        "status": "partial" if str(parsed.get("status") or "").strip().lower() == "partial" else "ready",
    }
    payload["markdown"] = _markdown_from_attachment_payload(payload)
    return payload
