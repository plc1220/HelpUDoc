"""FastAPI surface for the DeepAgents service."""
from __future__ import annotations

import asyncio
from typing import Any, Dict, List, AsyncGenerator, Iterable, Sequence, Set, Optional, Tuple
import base64
import html as html_lib
import importlib
import json
import logging
import os
import fnmatch
import mimetypes
import re
import shutil
import sys
import tempfile
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from uuid import uuid4
from dotenv import load_dotenv

from fastapi import FastAPI, HTTPException, Body, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field, model_validator
from google.genai import types as genai_types

from .configuration import describe_workspace_root, load_settings
from .graph import AgentRegistry
from .state import AgentRuntimeState
from .tools_and_schemas import ToolFactory, GeminiClientManager
from .mcp_manager import describe_mcp_servers
from .interrupt_payloads import (
    extract_interrupt_payload_from_tool_call,
    extract_interrupt_payload_from_tool_text,
    normalize_interrupt_payload_value,
)
from .utils import SourceTracker
from langchain_core.callbacks.base import AsyncCallbackHandler
from .rag_worker import RagIndexWorker
from .rag_indexer import _safe_join_workspace
from .skills_registry import (
    activate_skill_context,
    build_loaded_skill_text,
    collect_tool_names,
    find_skill,
    is_skill_allowed,
    load_skills,
    read_skill_content,
)
from .paper2slides_runner import run_paper2slides, export_pptx_from_pdf
from .jwt_utils import decode_and_verify_hs256_jwt
from .langfuse_callbacks import langfuse_langchain_callbacks
from langgraph.errors import GraphInterrupt
from langgraph.types import Command
from paper2slides.raganything.parser import DoclingParser


logger = logging.getLogger(__name__)
_INTERRUPT_TOOL_NAMES: Set[str] = {"request_clarification", "request_human_action"}
_LOCAL_DEV_AGENT_JWT_SECRET = "helpudoc-local-dev-agent-jwt-secret"
_RAG_PREFETCHABLE_EXTENSIONS: Set[str] = {".pdf", ".doc", ".docx", ".md", ".html", ".htm"}
_TAGGED_HTML_EXTENSIONS: Set[str] = {".html", ".htm"}
_TAGGED_RAG_CONTEXT_CHAR_BUDGET = 6000


def _get_agent_jwt_secret() -> str:
    configured = os.getenv("AGENT_JWT_SECRET", "").strip()
    if configured:
        return configured
    env = os.getenv("NODE_ENV", "").strip().lower()
    if not env or env == "development":
        return _LOCAL_DEV_AGENT_JWT_SECRET
    return ""


def _format_exception(exc: BaseException) -> str:
    if isinstance(exc, BaseExceptionGroup):
        parts = [_format_exception(inner) for inner in exc.exceptions]
        message = "; ".join(part for part in parts if part)
        return message or (str(exc) or repr(exc))
    return str(exc) or repr(exc)


def _filter_rag_prefetchable_tagged_files(tagged_paths: Sequence[str]) -> List[str]:
    candidates: List[str] = []
    for raw in tagged_paths:
        if not isinstance(raw, str):
            continue
        cleaned = raw.strip()
        if not cleaned:
            continue
        suffix = Path(cleaned).suffix.lower()
        if suffix in _RAG_PREFETCHABLE_EXTENSIONS:
            candidates.append(cleaned)
    return candidates


def _normalize_tagged_file_paths(tagged_paths: Sequence[str]) -> List[str]:
    normalized: List[str] = []
    for raw in tagged_paths:
        if not isinstance(raw, str):
            continue
        cleaned = raw.strip().replace("\\", "/")
        if not cleaned:
            continue
        if not cleaned.startswith("/"):
            cleaned = f"/{cleaned.lstrip('/')}"
        normalized.append(cleaned)
    return sorted(set(normalized))


def _allow_basename_match_for_path(path_value: str) -> bool:
    normalized = str(path_value or "").strip().replace("\\", "/").lstrip("/")
    if not normalized:
        return False
    if normalized.startswith(".system/"):
        return False
    return "/" not in normalized


def _build_tagged_rag_keywords(prompt: str, tagged_paths: Sequence[str]) -> List[str]:
    keywords: List[str] = []
    if isinstance(prompt, str) and prompt.strip():
        keywords.append(prompt.strip())
    for item in _normalize_tagged_file_paths(tagged_paths):
        keywords.append(item)
        name = Path(item).name
        if name:
            keywords.append(name)
    deduped: List[str] = []
    seen: Set[str] = set()
    for item in keywords:
        if item not in seen:
            seen.add(item)
            deduped.append(item)
    return deduped


def _filter_rag_chunks_to_tagged_paths(chunks: Sequence[Dict[str, Any]], tagged_paths: Sequence[str]) -> List[Dict[str, Any]]:
    normalized = _normalize_tagged_file_paths(tagged_paths)
    if not normalized:
        return list(chunks)
    basenames = {Path(item).name for item in normalized if _allow_basename_match_for_path(item)}
    filtered: List[Dict[str, Any]] = []
    for chunk in chunks:
        file_path = str(chunk.get("file_path") or "").strip().replace("\\", "/")
        if file_path and not file_path.startswith("/"):
            file_path = f"/{file_path.lstrip('/')}"
        if file_path in normalized or Path(file_path).name in basenames:
            filtered.append(chunk)
    return filtered


def _compress_tagged_context_lines(lines: Sequence[str], *, max_chars: int = _TAGGED_RAG_CONTEXT_CHAR_BUDGET) -> str | None:
    collected: List[str] = []
    total = 0
    for raw in lines:
        content = str(raw or "").strip()
        if not content:
            continue
        piece = content if not collected else f"\n\n{content}"
        if total + len(piece) > max_chars:
            remaining = max_chars - total
            if remaining > 64:
                collected.append(piece[:remaining].rstrip() + "\n\n[Truncated]")
            break
        collected.append(piece if not collected else content)
        total += len(piece)
    if not collected:
        return None
    return "\n\n".join(collected)[:max_chars]


def _strip_html_fragment(fragment: str) -> str:
    text = re.sub(r"(?is)<[^>]+>", " ", fragment or "")
    text = html_lib.unescape(text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _extract_html_outline_from_path(path: Path, *, max_chars: int = _TAGGED_RAG_CONTEXT_CHAR_BUDGET) -> str | None:
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        logger.exception("Failed reading tagged HTML outline: %s", path)
        return None
    sanitized = re.sub(r"(?is)<script\b[^>]*>.*?</script>", " ", raw)
    sanitized = re.sub(r"(?is)<style\b[^>]*>.*?</style>", " ", sanitized)
    sanitized = re.sub(r"(?is)<!--.*?-->", " ", sanitized)

    outline_parts: List[str] = []
    title_match = re.search(r"(?is)<title[^>]*>(.*?)</title>", sanitized)
    title_text = _strip_html_fragment(title_match.group(1)) if title_match else ""
    if title_text:
        outline_parts.append(f"TITLE: {title_text}")

    headings = [
        _strip_html_fragment(match)
        for match in re.findall(r"(?is)<h[1-3][^>]*>(.*?)</h[1-3]>", sanitized)
    ]
    headings = [item for item in headings if item]
    if headings:
        outline_parts.append("HEADINGS:")
        outline_parts.extend(f"- {item}" for item in headings[:12])

    paragraphs = [
        _strip_html_fragment(match)
        for match in re.findall(r"(?is)<p[^>]*>(.*?)</p>", sanitized)
    ]
    paragraphs = [item for item in paragraphs if item]
    if paragraphs:
        outline_parts.append("EXCERPTS:")
        outline_parts.extend(paragraphs[:8])

    if not outline_parts:
        fallback_text = _strip_html_fragment(sanitized)
        if fallback_text:
            outline_parts.append(fallback_text)

    return _compress_tagged_context_lines(outline_parts, max_chars=max_chars)


def _append_tagged_file_guidance(prompt: str, tagged_paths: Sequence[str]) -> str:
    if not prompt:
        return prompt
    if "Tagged file guidance:" in prompt:
        return prompt
    has_html = any(
        isinstance(raw, str) and Path(raw.strip()).suffix.lower() in _TAGGED_HTML_EXTENSIONS
        for raw in tagged_paths
    )
    if not has_html:
        return prompt
    guidance = (
        "Tagged file guidance:\n"
        "- Treat tagged .html files as reference artifacts, not raw context to ingest in full.\n"
        "- Do not read an entire report HTML unless absolutely necessary.\n"
        "- Prefer the canonical dataset as the source of truth and inspect only targeted report sections if needed."
    )
    return f"{prompt.rstrip()}\n\n{guidance}"


def _append_artifact_first_guidance(
    prompt: str,
    file_context_refs: Sequence[Dict[str, Any]],
    tagged_paths: Sequence[str],
    *,
    multimodal_active: bool,
) -> str:
    if not prompt:
        return prompt
    if "Artifact-first guidance:" in prompt:
        return prompt
    if not file_context_refs:
        return prompt
    ready_refs = [
        item
        for item in file_context_refs
        if str(item.get("status") or "").strip().lower() in {"ready", "partial"}
    ]
    if not ready_refs:
        return prompt
    binary_ready = [
        item
        for item in ready_refs
        if not str(item.get("sourceMimeType") or "").strip().lower().startswith("text/")
    ]
    if not binary_ready:
        return prompt
    guidance_lines = [
        "Artifact-first guidance:",
        "- Ready derived artifacts are available for the attached files and are the primary source of truth for this turn.",
        "- Prefer the tagged derived artifact paths over the original source file when answering.",
        "- Do not call read_file on the original binary source (.docx, .pptx, .pdf, etc.) if a ready derived artifact is already available.",
    ]
    if multimodal_active:
        guidance_lines.append("- Use the current-turn multimodal attachment only for additional grounding; keep follow-up reasoning anchored to the derived artifact.")
    else:
        guidance_lines.append("- If you need to inspect content, read the derived artifact markdown first rather than the original binary file.")
    if tagged_paths:
        guidance_lines.append("- Tagged derived artifacts:")
        guidance_lines.extend(f"  - {path}" for path in tagged_paths)
    return f"{prompt.rstrip()}\n\n" + "\n".join(guidance_lines)


_EXTERNAL_FRESHNESS_HINTS = (
    "latest",
    "current",
    "today",
    "up-to-date",
    "up to date",
    "as of ",
    "verify",
    "check",
    "pricing",
    "price",
    "cost",
    "thinking token",
    "thinking tokens",
    "model version",
)


def _should_force_tagged_files_rag_only(
    *,
    env_tagged_files_rag_only: bool,
    explicit_artifact_paths: Sequence[str],
    message_content: Sequence[Dict[str, Any]] | None,
    preferred_mcp_server: str | None,
    prompt_for_tagged_files: str,
) -> bool:
    """Keep artifact-first turns narrow, except when the user explicitly asks for fresh external checks."""
    if env_tagged_files_rag_only:
        return True
    if not explicit_artifact_paths or message_content:
        return False
    if str(preferred_mcp_server or "").strip():
        return False
    lowered_prompt = str(prompt_for_tagged_files or "").strip().lower()
    if lowered_prompt and any(marker in lowered_prompt for marker in _EXTERNAL_FRESHNESS_HINTS):
        return False
    return True


class ChatRequest(BaseModel):
    message: str
    history: List[Dict[str, Any]] | None = None
    forceReset: bool = False
    fileContextRefs: List[Dict[str, Any]] | None = None
    messageContent: List[Dict[str, Any]] | None = None


class ChatResponse(BaseModel):
    reply: Any


class Action(BaseModel):
    name: str
    args: Dict[str, Any] = Field(default_factory=dict)


class Decision(BaseModel):
    type: str
    edited_action: Optional[Action] = None
    message: Optional[str] = None


class ResumeChatRequest(BaseModel):
    decisions: List[Decision]


class InterruptResponseRequest(BaseModel):
    message: Optional[str] = None
    selectedChoiceIds: List[str] = Field(default_factory=list)
    selectedValues: List[str] = Field(default_factory=list)
    answersByQuestionId: Dict[str, str | List[str]] = Field(default_factory=dict)


class InterruptAction(BaseModel):
    id: str
    value: Optional[str] = None
    payload: Dict[str, Any] = Field(default_factory=dict)
    text: Optional[str] = None


class InterruptActionRequest(BaseModel):
    action: InterruptAction


class AttachmentUnderstandingRequest(BaseModel):
    fileName: str
    mimeType: str
    contentB64: str = ""
    workspaceId: Optional[str] = None
    relativePath: Optional[str] = None

    @model_validator(mode="after")
    def validate_attachment_source(self) -> "AttachmentUnderstandingRequest":
        has_b64 = bool((self.contentB64 or "").strip())
        has_path = bool((self.workspaceId or "").strip() and (self.relativePath or "").strip())
        if not has_b64 and not has_path:
            raise ValueError("Either contentB64 or workspaceId+relativePath is required")
        return self


class AttachmentUnderstandingSection(BaseModel):
    heading: str
    body: str


class AttachmentUnderstandingAsset(BaseModel):
    name: str
    mimeType: str
    contentB64: str
    sourcePath: Optional[str] = None
    caption: Optional[str] = None
    footnote: Optional[str] = None


class AttachmentUnderstandingResponse(BaseModel):
    title: str
    summary: str
    outline: List[str] = Field(default_factory=list)
    markdown: str
    sections: List[AttachmentUnderstandingSection] = Field(default_factory=list)
    extractedAssets: List[AttachmentUnderstandingAsset] = Field(default_factory=list)
    effectiveMode: str = "part"
    status: str = "ready"


class RagQueryRequest(BaseModel):
    query: str
    mode: str = "local"
    onlyNeedContext: bool = True
    includeReferences: bool = False


class RagQueryResponse(BaseModel):
    response: str


class RagStatusRequest(BaseModel):
    files: List[str]


class RagStatusResponse(BaseModel):
    statuses: Dict[str, Any]


class Paper2SlidesFile(BaseModel):
    name: str
    contentB64: str


class Paper2SlidesOptions(BaseModel):
    output: str | None = None
    content: str | None = None
    style: str | None = None
    length: str | None = None
    mode: str | None = None
    parallel: int | bool | None = None
    fromStage: str | None = None
    exportPptx: bool | None = None


class Paper2SlidesImage(BaseModel):
    name: str
    contentB64: str


class Paper2SlidesRunRequest(BaseModel):
    files: List[Paper2SlidesFile]
    options: Paper2SlidesOptions = Field(default_factory=Paper2SlidesOptions)


class Paper2SlidesRunResponse(BaseModel):
    pdfB64: str | None = None
    pptxB64: str | None = None
    images: List[Paper2SlidesImage] = []


class Paper2SlidesExportRequest(BaseModel):
    fileName: str
    contentB64: str


class Paper2SlidesExportResponse(BaseModel):
    pptxB64: str


def _extract_json_block(text: str) -> str:
    fenced = re.search(r"```json\s*(.*?)```", text or "", flags=re.IGNORECASE | re.DOTALL)
    if fenced:
        return fenced.group(1).strip()
    return (text or "").strip()


def _response_text(response: Any) -> str:
    text = getattr(response, "text", None)
    if isinstance(text, str) and text.strip():
        return text.strip()
    parts: List[str] = []
    for candidate in getattr(response, "candidates", []) or []:
        content = getattr(candidate, "content", None)
        if not content:
            continue
        for part in getattr(content, "parts", []) or []:
            value = getattr(part, "text", None)
            if isinstance(value, str) and value.strip():
                parts.append(value.strip())
    return "\n".join(parts).strip()


def _split_sections_from_markdown(markdown: str) -> List[Dict[str, str]]:
    sections: List[Dict[str, str]] = []
    heading = "Overview"
    body_lines: List[str] = []
    for line in (markdown or "").splitlines():
        match = re.match(r"^#{1,6}\s+(.*)$", line.strip())
        if match:
            if body_lines:
                sections.append({"heading": heading, "body": "\n".join(body_lines).strip()})
            heading = match.group(1).strip() or "Section"
            body_lines = []
            continue
        body_lines.append(line)
    if body_lines:
        sections.append({"heading": heading, "body": "\n".join(body_lines).strip()})
    return [section for section in sections if section.get("body")]


def _markdown_from_attachment_payload(payload: Dict[str, Any]) -> str:
    title = str(payload.get("title") or "Attachment").strip() or "Attachment"
    summary = str(payload.get("summary") or "").strip()
    outline = [str(item).strip() for item in (payload.get("outline") or []) if str(item).strip()]
    raw_sections = payload.get("sections") or []
    lines: List[str] = [f"# {title}", ""]
    if summary:
        lines.extend(["## Summary", "", summary, ""])
    if outline:
        lines.extend(["## Outline", ""])
        lines.extend(f"- {item}" for item in outline)
        lines.append("")
    sections_added = False
    if isinstance(raw_sections, list):
        for item in raw_sections:
            if not isinstance(item, dict):
                continue
            heading = str(item.get("heading") or "").strip() or "Section"
            body = str(item.get("body") or "").strip()
            if not body:
                continue
            lines.extend([f"## {heading}", "", body, ""])
            sections_added = True
    normalized_body = str(payload.get("normalizedBody") or "").strip()
    if normalized_body and not sections_added:
        lines.extend(["## Details", "", normalized_body, ""])
    return "\n".join(lines).strip()


def _build_partial_attachment_payload(file_name: str, text: str, *, effective_mode: str) -> Dict[str, Any]:
    cleaned_lines = [line.strip() for line in (text or "").splitlines() if line.strip()]
    summary = cleaned_lines[0] if cleaned_lines else f"Extracted content from {file_name}"
    excerpt = "\n".join(cleaned_lines[:80]).strip() or f"Unable to extract detailed content from {file_name}."
    sections = [{"heading": "Source Excerpt", "body": excerpt}]
    payload = {
        "title": file_name,
        "summary": summary,
        "outline": ["Source Excerpt"],
        "sections": sections,
        "effectiveMode": effective_mode,
        "status": "partial",
    }
    payload["markdown"] = _markdown_from_attachment_payload(payload)
    return payload


def _docling_available() -> bool:
    try:
        __import__("docling")
    except Exception:
        return False
    cli_path = shutil.which("docling")
    if cli_path:
        return True
    venv_candidates = [
        Path(sys.executable).parent / "docling",
        Path(sys.prefix) / "bin" / "docling",
    ]
    return any(candidate.exists() for candidate in venv_candidates)


def _docling_markdown_to_payload(file_name: str, markdown: str, *, effective_mode: str = "parser") -> Dict[str, Any]:
    cleaned_markdown = (markdown or "").strip()
    sections = _split_sections_from_markdown(cleaned_markdown)
    title = file_name
    title_match = re.search(r"^#\s+(.+)$", cleaned_markdown, flags=re.MULTILINE)
    if title_match and title_match.group(1).strip():
        title = title_match.group(1).strip()
    summary = ""
    for section in sections:
        body = str(section.get("body") or "").strip()
        if body:
            summary = body.splitlines()[0].strip()
            if summary:
                break
    outline = [str(section.get("heading") or "").strip() for section in sections if str(section.get("heading") or "").strip()]
    payload = {
        "title": title,
        "summary": summary or f"Parsed content from {file_name}",
        "outline": outline,
        "sections": sections,
        "normalizedBody": cleaned_markdown,
        "effectiveMode": effective_mode,
        "status": "ready",
    }
    payload["markdown"] = cleaned_markdown or _markdown_from_attachment_payload(payload)
    return payload


def _extract_docling_payload(file_name: str, mime_type: str, buffer: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    suffix = Path(file_name).suffix.lower()
    guessed_suffix = mimetypes.guess_extension(mime_type or "") or ""
    if not suffix and guessed_suffix:
        suffix = guessed_suffix.lower()
    if not suffix:
        suffix = ".bin"
    with tempfile.TemporaryDirectory(prefix="helpudoc-docling-") as temp_dir:
        temp_root = Path(temp_dir)
        temp_input = temp_root / f"source{suffix}"
        temp_input.write_bytes(buffer)
        output_dir = temp_root / "parsed"
        parser = DoclingParser()
        content_list = parser.parse_document(temp_input, output_dir=str(output_dir))
        stem = temp_input.stem
        docling_dir = output_dir / stem / "docling"
        md_path = docling_dir / f"{stem}.md"
        if not md_path.exists():
            raise RuntimeError(f"Docling did not produce markdown for {file_name}")
        markdown = md_path.read_text(encoding="utf-8", errors="replace").strip()
        if not markdown:
            raise RuntimeError(f"Docling produced empty markdown for {file_name}")
        extracted_assets: List[Dict[str, Any]] = []
        seen_paths: set[str] = set()
        for item in content_list:
            if not isinstance(item, dict):
                continue
            if str(item.get("type") or "").strip().lower() != "image":
                continue
            img_path_raw = str(item.get("img_path") or "").strip()
            if not img_path_raw or img_path_raw in seen_paths:
                continue
            seen_paths.add(img_path_raw)
            img_path = Path(img_path_raw)
            if not img_path.exists():
                continue
            mime = mimetypes.guess_type(img_path.name)[0] or "image/png"
            try:
                source_path = str(img_path.relative_to(docling_dir)).replace("\\", "/")
            except Exception:
                source_path = img_path.name
            extracted_assets.append(
                {
                    "name": img_path.name,
                    "mimeType": mime,
                    "contentB64": base64.b64encode(img_path.read_bytes()).decode("ascii"),
                    "sourcePath": source_path,
                    "caption": str(item.get("image_caption") or "").strip() or None,
                    "footnote": str(item.get("image_footnote") or "").strip() or None,
                }
            )
        return markdown, extracted_assets


def _extract_text_from_docx(buffer: bytes) -> str:
    from docx import Document  # type: ignore

    document = Document(BytesIO(buffer))
    chunks: List[str] = []
    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            chunks.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                chunks.append(" | ".join(cells))
    return "\n".join(chunks).strip()


def _extract_text_from_pptx(buffer: bytes) -> str:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(BytesIO(buffer))
    chunks: List[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: List[str] = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                slide_lines.append(text.strip())
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                notes_text = ""
            if notes_text and notes_text.strip():
                slide_lines.append(f"Notes: {notes_text.strip()}")
        chunks.append("\n".join(slide_lines))
    return "\n\n".join(chunks).strip()


def _extract_text_from_pdf(buffer: bytes) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(BytesIO(buffer))
    chunks: List[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            chunks.append(f"Page {index}\n{text}")
    return "\n\n".join(chunks).strip()


def _guess_attachment_strategy(file_name: str, mime_type: str) -> Tuple[str, str]:
    suffix = Path(file_name).suffix.lower()
    normalized_mime = (mime_type or "").lower()
    if normalized_mime.startswith("image/") or suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}:
        return "image", "part"
    if suffix == ".docx" or "wordprocessingml" in normalized_mime:
        return "docx", "parser"
    if suffix == ".pptx" or "presentationml" in normalized_mime:
        return "pptx", "parser"
    if suffix == ".pdf" or normalized_mime == "application/pdf":
        return "pdf", "parser"
    return "text", "part"


def _attachment_understanding_prompt(file_name: str, extracted_text: str | None, kind: str) -> str:
    base = (
        "Return strict JSON only with keys: "
        "title, summary, outline, sections, normalizedBody, effectiveMode, status. "
        "sections must be an array of objects with heading and body. "
        "effectiveMode must be one of part, parser, hybrid. "
        "status must be ready or partial. "
        "Preserve the document structure and factual detail instead of over-summarizing. "
        f"File name: {file_name}. "
        f"Kind: {kind}."
    )
    if extracted_text is None:
        return (
            f"{base} Understand the attached file directly. "
            "If details are uncertain, say so briefly in the summary instead of inventing content."
        )
    excerpt = extracted_text[:30000]
    return f"{base}\n\nSource content:\n{excerpt}"


def _parse_attachment_payload(raw_text: str, file_name: str, *, fallback_text: str, fallback_mode: str) -> Dict[str, Any]:
    try:
        parsed = json.loads(_extract_json_block(raw_text))
    except Exception:
        return _build_partial_attachment_payload(file_name, fallback_text, effective_mode=fallback_mode)
    if not isinstance(parsed, dict):
        return _build_partial_attachment_payload(file_name, fallback_text, effective_mode=fallback_mode)
    payload = {
        "title": str(parsed.get("title") or file_name).strip() or file_name,
        "summary": str(parsed.get("summary") or "").strip(),
        "outline": [str(item).strip() for item in (parsed.get("outline") or []) if str(item).strip()],
        "sections": [
            {
                "heading": str(item.get("heading") or "").strip() or "Section",
                "body": str(item.get("body") or "").strip(),
            }
            for item in (parsed.get("sections") or [])
            if isinstance(item, dict) and str(item.get("body") or "").strip()
        ],
        "normalizedBody": str(parsed.get("normalizedBody") or "").strip(),
        "effectiveMode": str(parsed.get("effectiveMode") or fallback_mode).strip() or fallback_mode,
        "status": "partial" if str(parsed.get("status") or "").strip().lower() == "partial" else "ready",
    }
    payload["markdown"] = _markdown_from_attachment_payload(payload)
    return payload


def _copy_content_block(block: Any) -> Any:
    if isinstance(block, dict):
        return dict(block)
    return block


def _extract_text_from_content(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: List[str] = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
                continue
            if isinstance(item, dict):
                if isinstance(item.get("text"), str):
                    parts.append(item["text"])
                    continue
                if item.get("type") == "text-plain" and isinstance(item.get("text"), str):
                    parts.append(item["text"])
                    continue
                if "content" in item and isinstance(item.get("content"), str):
                    parts.append(item["content"])
        return "\n".join(part for part in parts if part).strip()
    return str(content or "").strip()


def _replace_content_text(content: Any, text: str) -> Any:
    if isinstance(content, str):
        return text
    if isinstance(content, list):
        updated: List[Any] = []
        replaced = False
        for item in content:
            if isinstance(item, dict):
                copied = dict(item)
                block_type = str(copied.get("type") or "").strip().lower()
                if block_type in {"text", "text-plain"} or isinstance(copied.get("text"), str):
                    copied["text"] = text
                    updated.append(copied)
                    replaced = True
                    continue
                if isinstance(copied.get("content"), str):
                    copied["content"] = text
                    updated.append(copied)
                    replaced = True
                    continue
                updated.append(copied)
                continue
            updated.append(item)
        if replaced:
            return updated
        return [{"type": "text", "text": text}, *updated]
    return text


def _host_datetime_context_block() -> str:
    """Wall-clock snapshot for the model; agent system prompts are cached and omit real time."""
    utc_now = datetime.now(timezone.utc)
    local_now = datetime.now().astimezone()
    local_tz_name = local_now.tzname() or "local"
    return (
        "[Host time]\n"
        "Use this block as the authoritative current date/time for this turn. "
        "If the user says \"today\", \"tomorrow\", \"this week\", a calendar year, a deadline, or asks for a dated filename, "
        "use these timestamps instead of model priors.\n"
        f"Authoritative local date: {local_now.date().isoformat()}\n"
        f"UTC: {utc_now.isoformat(timespec='seconds')}\n"
        f"Server local ({local_tz_name}): {local_now.isoformat(timespec='seconds')}"
    )


def _inject_host_datetime_context(payload: List[Dict[str, Any]]) -> None:
    """Inject a fresh host-time system message and also prefix the latest user message for extra salience."""
    block = _host_datetime_context_block()
    payload.insert(0, {"role": "system", "content": block})
    for index in range(len(payload) - 1, -1, -1):
        role = str(payload[index].get("role") or "").strip().lower()
        if role not in {"user", "human"}:
            continue
        content = payload[index].get("content")
        if isinstance(content, str):
            stripped = content.strip()
            payload[index]["content"] = f"{block}\n\n{stripped}" if stripped else block
        elif isinstance(content, list):
            payload[index]["content"] = [{"type": "text", "text": block}, *content]
        else:
            text = str(content or "").strip()
            payload[index]["content"] = f"{block}\n\n{text}" if text else block
        return


class EmbeddedDirective(BaseModel):
    kind: str
    skillId: Optional[str] = None
    serverId: Optional[str] = None


BASE_DIR = Path(__file__).resolve().parent


_FILE_RESULT_PATTERNS = [
    re.compile(r"Updated file (?P<path>/[^\s]+)"),
    re.compile(r"in '(?P<path>/[^']+)'"),
    re.compile(r"Appended (?P<src>/[^\s]+) to (?P<dst>/[^\s]+)"),
    re.compile(r"Created PDF (?P<path>/[^\s]+)"),
]
_DIRECTIVE_BLOCK_RE = re.compile(
    r"^\s*<<<HELPUDOC_DIRECTIVE\s*\n(?P<payload>\{.*?\})\n>>>\s*(?P<rest>[\s\S]*)$",
    re.DOTALL,
)
_RAW_SKILL_DIRECTIVE_RE = re.compile(
    r"^\s*/skill\s+(?P<skill_id>[^\s]+)(?:\s+(?P<prompt>[\s\S]*))?$",
    re.IGNORECASE,
)
_RAW_MCP_DIRECTIVE_RE = re.compile(
    r"^\s*/mcp\s+(?P<server_id>[^\s]+)(?:\s+(?P<prompt>[\s\S]*))?$",
    re.IGNORECASE,
)
_LEGACY_SKILL_PROMPT_RE = re.compile(
    r'^\s*Use the "(?P<skill_id>[^"]+)" skill for this task\.\s*'
    r'First call load_skill with "(?P=skill_id)" to load the skill instructions, then follow that skill closely\.\s*'
    r'(?:User request:\s*(?P<prompt>[\s\S]*))?$',
    re.IGNORECASE,
)
_LEGACY_MCP_PROMPT_RE = re.compile(
    r'^\s*Prefer tools from the MCP server "(?P<server_id>[^"]+)" for this task\.\s*'
    r'[\s\S]*?(?:User request:\s*(?P<prompt>[\s\S]*))?$',
    re.IGNORECASE,
)
_LINE_SKILL_DIRECTIVE_RE = re.compile(
    r"^\s*(?:(?:please\s+)?use\s+)?/skill\s+(?P<skill_id>[^\s]+)(?:\s+(?P<prompt>[\s\S]*))?$",
    re.IGNORECASE,
)
_LINE_MCP_DIRECTIVE_RE = re.compile(
    r"^\s*(?:(?:please\s+)?use\s+)?/mcp\s+(?P<server_id>[^\s]+)(?:\s+(?P<prompt>[\s\S]*))?$",
    re.IGNORECASE,
)


def _infer_mime_type(file_path: str) -> str:
    guessed, _ = mimetypes.guess_type(file_path)
    return guessed or "application/octet-stream"


def _extract_output_files_from_tool_result(name: str, text: str) -> List[Dict[str, Any]]:
    if not text:
        return []
    outputs: List[Dict[str, Any]] = []
    if name == "write_file":
        match = _FILE_RESULT_PATTERNS[0].search(text)
        if match:
            path = match.group("path")
            outputs.append({"path": path.lstrip("/"), "mimeType": _infer_mime_type(path)})
        return outputs
    if name == "edit_file":
        match = _FILE_RESULT_PATTERNS[1].search(text)
        if match:
            path = match.group("path")
            outputs.append({"path": path.lstrip("/"), "mimeType": _infer_mime_type(path)})
        return outputs
    if name == "append_to_report":
        match = _FILE_RESULT_PATTERNS[2].search(text)
        if match:
            path = match.group("dst")
            outputs.append({"path": path.lstrip("/"), "mimeType": _infer_mime_type(path)})
        return outputs
    if name == "create_pdf_from_images":
        match = _FILE_RESULT_PATTERNS[3].search(text)
        if match:
            path = match.group("path")
            outputs.append({"path": path.lstrip("/"), "mimeType": _infer_mime_type(path)})
        return outputs
    return outputs


def _load_env_files() -> None:
    """Load environment variables from known locations.

    We prioritize the agent's .env (agent/.env) and then allow any existing
    process-level env vars to remain.
    """
    env_file = os.getenv("ENV_FILE")
    if env_file:
        load_dotenv(env_file)
        return
    load_dotenv(BASE_DIR.parent / ".env")


def create_app() -> FastAPI:
    _load_env_files()
    config_path: Optional[Path] = None
    env_config_path = os.getenv("AGENT_CONFIG_PATH")
    if env_config_path:
        candidate = Path(env_config_path).expanduser()
        if not candidate.is_absolute():
            candidate = (BASE_DIR.parent.parent / candidate).resolve()
        if candidate.exists():
            config_path = candidate
        else:
            logger.warning("AGENT_CONFIG_PATH is set but file does not exist; falling back to built-in config", extra={"path": env_config_path})

    settings = load_settings(config_path)
    workspace_root_diagnostic = describe_workspace_root(settings)
    workspace_root_message = (
        f"[agent] Workspace root: {workspace_root_diagnostic['resolved_path']} "
        f"(source={workspace_root_diagnostic['source']} "
        f"raw={workspace_root_diagnostic['raw_value'] or '<config>'})"
    )
    print(workspace_root_message)
    logger.info(workspace_root_message)
    file_understanding_mode = (os.getenv("FILE_UNDERSTANDING_MODE", "part-first") or "part-first").strip()
    rag_parser_pipeline = (os.getenv("RAG_PARSER_PIPELINE", "raganything") or "raganything").strip().lower()
    raganything_parser = (os.getenv("RAGANYTHING_PARSER", "docling") or "docling").strip().lower()
    parser_enrichment_mode = (
        os.getenv("PARSER_ENRICHMENT_MODE")
        or os.getenv("PARSER")
        or raganything_parser
    ).strip() or raganything_parser
    dependency_diag = {
        "lightrag": True,
        "raganything": True,
        "docling": _docling_available(),
    }
    try:
        importlib.import_module("lightrag")
    except Exception:
        dependency_diag["lightrag"] = False
    try:
        importlib.import_module("raganything")
    except Exception:
        try:
            importlib.import_module("paper2slides.raganything")
        except Exception:
            dependency_diag["raganything"] = False
    if rag_parser_pipeline in {"raganything", "rag_anything", "rag-everything", "rageverything"} and raganything_parser == "docling":
        if not dependency_diag["docling"]:
            raise RuntimeError(
                "Docling is configured as the global parser, but the docling package/CLI is unavailable. "
                "Install docling and ensure the `docling` command is on PATH."
            )
    logger.info(
        "[agent] File understanding: mode=%s parserPipeline=%s ragParser=%s parserEnrichment=%s deps=%s python=%s",
        file_understanding_mode,
        rag_parser_pipeline,
        raganything_parser,
        parser_enrichment_mode,
        dependency_diag,
        sys.executable,
    )
    source_tracker = SourceTracker()
    gemini_manager = GeminiClientManager(settings)
    tool_factory = ToolFactory(settings, source_tracker, gemini_manager)
    registry = AgentRegistry(settings, tool_factory)
    agent_jwt_secret = _get_agent_jwt_secret()

    app = FastAPI(title="DeepAgents Service", version="0.2.0")
    rag_worker = RagIndexWorker(settings.backend.workspace_root)

    @app.on_event("startup")
    async def _startup() -> None:
        try:
            await rag_worker.start()
        except Exception:
            logger.exception("Failed to start RAG index worker (continuing without it)")

    @app.on_event("shutdown")
    async def _shutdown() -> None:
        try:
            await rag_worker.stop()
        except Exception:
            logger.exception("Failed to stop RAG index worker cleanly")

    @app.post("/rag/workspaces/{workspace_id}/query", response_model=RagQueryResponse)
    async def rag_query(workspace_id: str, req: RagQueryRequest = Body(...)):
        if not req.query or not req.query.strip():
            raise HTTPException(status_code=400, detail="query is required")
        try:
            response = await rag_worker.store.query(
                workspace_id,
                req.query,
                mode=req.mode,
                only_need_context=req.onlyNeedContext,
                include_references=req.includeReferences,
            )
            return RagQueryResponse(response=response)
        except FileNotFoundError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc

    @app.post("/rag/workspaces/{workspace_id}/status", response_model=RagStatusResponse)
    async def rag_status(workspace_id: str, req: RagStatusRequest = Body(...)):
        if not req.files:
            raise HTTPException(status_code=400, detail="files is required")
        statuses: Dict[str, Any] = {}
        for name in req.files:
            if not isinstance(name, str) or not name.strip():
                continue
            relative = name.strip().lstrip("/")
            try:
                status = await rag_worker.store.get_doc_status(workspace_id, relative)
            except Exception as exc:
                status = {"status": "error", "error": str(exc)}
            if status is None:
                statuses[name] = {"status": "not_indexed"}
            else:
                statuses[name] = {
                    "status": status.get("status", "unknown"),
                    "updatedAt": status.get("updated_at"),
                    "error": status.get("error_msg"),
                }
        return RagStatusResponse(statuses=statuses)

    @app.post("/paper2slides/run", response_model=Paper2SlidesRunResponse)
    async def paper2slides_run(req: Paper2SlidesRunRequest = Body(...)):
        if not req.files:
            raise HTTPException(status_code=400, detail="files is required")
        try:
            payload_files = [file.model_dump() for file in req.files]
            options = req.options.model_dump()
            result = await asyncio.to_thread(run_paper2slides, payload_files, options)
            return Paper2SlidesRunResponse(**result)
        except Exception as exc:
            raise HTTPException(status_code=500, detail=str(exc)) from exc

    @app.post("/paper2slides/export-pptx", response_model=Paper2SlidesExportResponse)
    async def paper2slides_export(req: Paper2SlidesExportRequest = Body(...)):
        if not req.contentB64:
            raise HTTPException(status_code=400, detail="contentB64 is required")
        try:
            result = await asyncio.to_thread(export_pptx_from_pdf, req.fileName, req.contentB64)
            return Paper2SlidesExportResponse(**result)
        except Exception as exc:
            raise HTTPException(status_code=500, detail=str(exc)) from exc

    @app.post("/attachments/understand", response_model=AttachmentUnderstandingResponse)
    async def understand_attachment(req: AttachmentUnderstandingRequest = Body(...)):
        if not req.fileName.strip():
            raise HTTPException(status_code=400, detail="fileName is required")
        buffer: bytes
        ws_id = (req.workspaceId or "").strip()
        rel = (req.relativePath or "").strip()
        if ws_id and rel:
            try:
                abs_path = _safe_join_workspace(
                    Path(settings.backend.workspace_root),
                    ws_id,
                    rel,
                )
            except ValueError as exc:
                raise HTTPException(status_code=400, detail=str(exc)) from exc
            if not abs_path.is_file():
                raise HTTPException(
                    status_code=404,
                    detail=f"Workspace file not found for workspaceId={ws_id!r} relativePath={rel!r}",
                )
            buffer = await asyncio.to_thread(abs_path.read_bytes)
        else:
            if not req.contentB64.strip():
                raise HTTPException(status_code=400, detail="contentB64 is required")
            try:
                buffer = base64.b64decode(req.contentB64)
            except Exception as exc:
                raise HTTPException(status_code=400, detail="contentB64 must be valid base64") from exc

        kind, fallback_mode = _guess_attachment_strategy(req.fileName, req.mimeType)
        extracted_text: str | None = None
        try:
            contents: List[Any]
            if kind in {"docx", "pptx", "pdf"}:
                try:
                    markdown, extracted_assets = await asyncio.to_thread(
                        _extract_docling_payload,
                        req.fileName,
                        req.mimeType,
                        buffer,
                    )
                    payload = _docling_markdown_to_payload(req.fileName, markdown, effective_mode="parser")
                    return AttachmentUnderstandingResponse(
                        title=str(payload.get("title") or req.fileName),
                        summary=str(payload.get("summary") or ""),
                        outline=[str(item) for item in (payload.get("outline") or [])],
                        markdown=str(payload.get("markdown") or ""),
                        sections=[
                            AttachmentUnderstandingSection(
                                heading=str(item.get("heading") or "Section"),
                                body=str(item.get("body") or ""),
                            )
                            for item in (payload.get("sections") or [])
                            if isinstance(item, dict) and str(item.get("body") or "").strip()
                        ],
                        extractedAssets=[
                            AttachmentUnderstandingAsset(
                                name=str(item.get("name") or "image.png"),
                                mimeType=str(item.get("mimeType") or "image/png"),
                                contentB64=str(item.get("contentB64") or ""),
                                sourcePath=str(item.get("sourcePath") or "") or None,
                                caption=str(item.get("caption") or "") or None,
                                footnote=str(item.get("footnote") or "") or None,
                            )
                            for item in extracted_assets
                            if str(item.get("contentB64") or "").strip()
                        ],
                        effectiveMode="parser",
                        status="ready",
                    )
                except Exception:
                    logger.exception("Docling extraction failed for %s; falling back to legacy extraction", req.fileName)
                    if kind == "docx":
                        extracted_text = _extract_text_from_docx(buffer)
                    elif kind == "pptx":
                        extracted_text = _extract_text_from_pptx(buffer)
                    else:
                        extracted_text = _extract_text_from_pdf(buffer)
                    contents = [_attachment_understanding_prompt(req.fileName, extracted_text, kind)]
            elif kind == "image":
                contents = [
                    genai_types.Part.from_bytes(data=buffer, mime_type=req.mimeType),
                    _attachment_understanding_prompt(req.fileName, None, kind),
                ]
            else:
                try:
                    extracted_text = buffer.decode("utf-8", errors="replace")
                except Exception:
                    extracted_text = ""
                contents = [_attachment_understanding_prompt(req.fileName, extracted_text, kind)]

            fallback_text = extracted_text or f"Unable to extract textual content from {req.fileName}."
            raw_response_text = ""
            for attempt in range(2):
                response = gemini_manager.client.models.generate_content(
                    model=gemini_manager.model_name,
                    contents=contents,
                    config=genai_types.GenerateContentConfig(
                        candidate_count=1,
                        temperature=0.1,
                    ),
                )
                raw_response_text = _response_text(response)
                payload = _parse_attachment_payload(
                    raw_response_text,
                    req.fileName,
                    fallback_text=fallback_text,
                    fallback_mode=fallback_mode,
                )
                if payload.get("status") != "partial" or attempt == 1:
                    return AttachmentUnderstandingResponse(
                        title=str(payload.get("title") or req.fileName),
                        summary=str(payload.get("summary") or ""),
                        outline=[str(item) for item in (payload.get("outline") or [])],
                        markdown=str(payload.get("markdown") or ""),
                        sections=[
                            AttachmentUnderstandingSection(
                                heading=str(item.get("heading") or "Section"),
                                body=str(item.get("body") or ""),
                            )
                            for item in (payload.get("sections") or [])
                            if isinstance(item, dict) and str(item.get("body") or "").strip()
                        ],
                        effectiveMode=str(payload.get("effectiveMode") or fallback_mode),
                        status=str(payload.get("status") or "partial"),
                    )
                contents = [
                    "Repair the previous response into strict JSON only using the required schema.",
                    raw_response_text,
                ]
        except Exception:
            logger.exception("Attachment understanding failed for %s", req.fileName)

        partial = _build_partial_attachment_payload(
            req.fileName,
            extracted_text or f"Unable to understand {req.fileName}.",
            effective_mode=fallback_mode,
        )
        return AttachmentUnderstandingResponse(
            title=str(partial.get("title") or req.fileName),
            summary=str(partial.get("summary") or ""),
            outline=[str(item) for item in (partial.get("outline") or [])],
            markdown=str(partial.get("markdown") or ""),
            sections=[
                AttachmentUnderstandingSection(
                    heading=str(item.get("heading") or "Section"),
                    body=str(item.get("body") or ""),
                )
                for item in (partial.get("sections") or [])
                if isinstance(item, dict) and str(item.get("body") or "").strip()
            ],
            effectiveMode=str(partial.get("effectiveMode") or fallback_mode),
            status="partial",
        )

    @app.get("/agents")
    def list_agents():
        skills = load_skills(settings.backend.skills_root) if settings.backend.skills_root else []
        tool_names = collect_tool_names(skills)
        if tool_names:
            tool_names = [name for name in tool_names if name in settings.tools]
        if not tool_names:
            tool_names = list(settings.tools.keys())
        else:
            for extra in ("list_skills", "load_skill"):
                if extra in settings.tools and extra not in tool_names:
                    tool_names.append(extra)
        return {
            "agents": [
                {
                    "name": "fast",
                    "displayName": "Fast",
                    "description": "General assistant optimized for speed (Gemini Flash).",
                    "tools": tool_names,
                    "subagents": [],
                },
                {
                    "name": "pro",
                    "displayName": "Pro",
                    "description": "General assistant optimized for quality (Gemini Pro).",
                    "tools": tool_names,
                    "subagents": [],
                },
                {
                    "name": "skill-builder",
                    "displayName": "Skill Builder",
                    "description": "Admin-oriented assistant for creating and updating skills with structured actions.",
                    "tools": tool_names,
                    "subagents": [],
                },
            ],
            "mcpServers": describe_mcp_servers(settings),
        }

    def _prepare_payload(message: ChatRequest) -> List[Dict[str, Any]]:
        payload: List[Dict[str, Any]] = []
        if message.history:
            for item in message.history:
                if isinstance(item, dict):
                    payload.append(dict(item))
                else:
                    payload.append({"role": "user", "content": str(item)})
        if message.messageContent:
            copied_blocks = [_copy_content_block(block) for block in message.messageContent]
            for index in range(len(payload) - 1, -1, -1):
                role = str(payload[index].get("role") or "").strip().lower()
                if role in {"user", "human"}:
                    payload[index]["content"] = copied_blocks
                    return payload
            payload.append({"role": "user", "content": copied_blocks})
            return payload
        if payload:
            return payload
        return [{"role": "user", "content": message.message}]

    def _extract_directive_from_text(text: str) -> Tuple[EmbeddedDirective | None, str]:
        if not text:
            return None, ""

        block_match = _DIRECTIVE_BLOCK_RE.match(text)
        if block_match:
            payload = block_match.group("payload")
            rest = (block_match.group("rest") or "").strip()
            try:
                directive = EmbeddedDirective.model_validate_json(payload)
            except Exception:
                return None, rest or text
            return directive, rest

        raw_skill_match = _RAW_SKILL_DIRECTIVE_RE.match(text)
        if raw_skill_match:
            return (
                EmbeddedDirective(kind="skill", skillId=(raw_skill_match.group("skill_id") or "").strip()),
                (raw_skill_match.group("prompt") or "").strip(),
            )

        raw_mcp_match = _RAW_MCP_DIRECTIVE_RE.match(text)
        if raw_mcp_match:
            return (
                EmbeddedDirective(kind="mcp", serverId=(raw_mcp_match.group("server_id") or "").strip()),
                (raw_mcp_match.group("prompt") or "").strip(),
            )

        legacy_skill_match = _LEGACY_SKILL_PROMPT_RE.match(text)
        if legacy_skill_match:
            return (
                EmbeddedDirective(kind="skill", skillId=(legacy_skill_match.group("skill_id") or "").strip()),
                (legacy_skill_match.group("prompt") or "").strip(),
            )

        legacy_mcp_match = _LEGACY_MCP_PROMPT_RE.match(text)
        if legacy_mcp_match:
            return (
                EmbeddedDirective(kind="mcp", serverId=(legacy_mcp_match.group("server_id") or "").strip()),
                (legacy_mcp_match.group("prompt") or "").strip(),
            )

        lines = text.splitlines()
        for index, raw_line in enumerate(lines):
            if not raw_line.strip():
                continue
            skill_line_match = _LINE_SKILL_DIRECTIVE_RE.match(raw_line)
            if skill_line_match:
                remainder = "\n".join(lines[index + 1 :]).strip()
                prompt_parts = [
                    (skill_line_match.group("prompt") or "").strip(),
                    remainder,
                ]
                prompt = "\n\n".join(part for part in prompt_parts if part)
                return (
                    EmbeddedDirective(kind="skill", skillId=(skill_line_match.group("skill_id") or "").strip()),
                    prompt,
                )

            mcp_line_match = _LINE_MCP_DIRECTIVE_RE.match(raw_line)
            if mcp_line_match:
                remainder = "\n".join(lines[index + 1 :]).strip()
                prompt_parts = [
                    (mcp_line_match.group("prompt") or "").strip(),
                    remainder,
                ]
                prompt = "\n\n".join(part for part in prompt_parts if part)
                return (
                    EmbeddedDirective(kind="mcp", serverId=(mcp_line_match.group("server_id") or "").strip()),
                    prompt,
                )
            break

        return None, text

    globals()["_extract_directive_from_text"] = _extract_directive_from_text

    def _build_preloaded_skill_prompt(
        runtime: AgentRuntimeState,
        skill_id: str,
        user_request: str,
    ) -> str:
        fallback_request = user_request.strip() or "Continue with the selected skill."
        skill = find_skill(settings.backend.skills_root, skill_id)
        if skill is None:
            return (
                f"The user explicitly selected skill '{skill_id}', but it was not found in the configured skills registry.\n\n"
                f"User request:\n{fallback_request}"
            )
        if not is_skill_allowed(skill, runtime.workspace_state.context):
            return (
                f"The user explicitly selected skill '{skill.skill_id}', but it is not allowed for this user.\n\n"
                f"User request:\n{fallback_request}"
            )
        try:
            content = read_skill_content(skill)
        except Exception as exc:
            return (
                f"The user explicitly selected skill '{skill_id}', but the skill could not be read: {exc}\n\n"
                f"User request:\n{fallback_request}"
            )

        activate_skill_context(runtime.workspace_state.context, skill)
        return "\n\n".join(
            [
                f"The selected skill '{skill.skill_id}' is already loaded and active for this turn.",
                "Do not call list_skills or load_skill again unless you need to switch to a different skill.",
                build_loaded_skill_text(skill, content),
                f"User request:\n{fallback_request}",
            ]
        )

    def _build_preferred_mcp_prompt(server_id: str, user_request: str) -> str:
        normalized_server_id = str(server_id or "").strip()
        fallback_request = user_request.strip() or f"Use MCP server '{normalized_server_id}' for this task."
        return "\n\n".join(
            [
                f"The preferred MCP server for this turn is '{normalized_server_id}'.",
                "Prefer tools from that server before unrelated MCP servers or general web search when they can satisfy the request.",
                f"User request:\n{fallback_request}",
            ]
        )

    def _apply_embedded_directives(
        runtime: AgentRuntimeState,
        payload: List[Dict[str, Any]],
    ) -> Tuple[List[Dict[str, Any]], str]:
        latest_user_text = ""
        for index in range(len(payload) - 1, -1, -1):
            message = payload[index]
            role = str(message.get("role") or "").strip().lower()
            if role not in {"user", "human"}:
                continue
            content = message.get("content")
            latest_text = _extract_text_from_content(content)
            if not latest_text:
                break
            directive, stripped_text = _extract_directive_from_text(latest_text)
            latest_user_text = stripped_text
            if directive is None:
                message["content"] = _replace_content_text(content, stripped_text)
                break
            if directive.kind == "skill" and directive.skillId:
                runtime.workspace_state.context.pop("preferred_mcp_server", None)
                message["content"] = _replace_content_text(
                    content,
                    _build_preloaded_skill_prompt(runtime, directive.skillId, stripped_text),
                )
            elif directive.kind == "mcp" and directive.serverId:
                runtime.workspace_state.context["preferred_mcp_server"] = directive.serverId
                message["content"] = _replace_content_text(
                    content,
                    _build_preferred_mcp_prompt(directive.serverId, stripped_text),
                )
            else:
                message["content"] = _replace_content_text(content, stripped_text)
            break
        return payload, latest_user_text

    def _extract_request_context(request: Request) -> Dict[str, Any]:
        """Extract backend-provided context (RBAC policy, user id) from JWT."""
        if not agent_jwt_secret:
            return {}
        raw_auth = request.headers.get("authorization") or ""
        token = ""
        if raw_auth.lower().startswith("bearer "):
            token = raw_auth.split(" ", 1)[1].strip()
        if not token:
            return {}
        payload = decode_and_verify_hs256_jwt(token, agent_jwt_secret)
        if not payload:
            return {}
        context: Dict[str, Any] = {}
        user_id = payload.get("userId") or payload.get("sub")
        if isinstance(user_id, str) and user_id.strip():
            context["user_id"] = user_id.strip()
        skill_allow_ids = payload.get("skillAllowIds") or []
        if isinstance(skill_allow_ids, list):
            context["skill_allow_ids"] = [str(x).strip() for x in skill_allow_ids if str(x).strip()]
        allow_ids = payload.get("mcpServerAllowIds") or []
        deny_ids = payload.get("mcpServerDenyIds") or []
        is_admin = bool(payload.get("isAdmin", False))
        allow_script_runner = bool(payload.get("allowScriptRunner") or payload.get("allow_script_runner"))
        if isinstance(allow_ids, list) or isinstance(deny_ids, list) or isinstance(is_admin, bool):
            context["mcp_policy"] = {
                "allowIds": [str(x) for x in (allow_ids or []) if str(x).strip()],
                "denyIds": [str(x) for x in (deny_ids or []) if str(x).strip()],
                "isAdmin": is_admin,
            }
        mcp_auth = payload.get("mcpAuth") or {}
        if isinstance(mcp_auth, dict):
            normalized_auth: Dict[str, Dict[str, str]] = {}
            for server_name, headers in mcp_auth.items():
                if not isinstance(server_name, str) or not server_name.strip():
                    continue
                if not isinstance(headers, dict):
                    continue
                normalized_headers: Dict[str, str] = {}
                for header_name, header_value in headers.items():
                    if not isinstance(header_name, str) or not header_name.strip():
                        continue
                    if isinstance(header_value, str) and header_value.strip():
                        normalized_headers[header_name] = header_value
                if normalized_headers:
                    normalized_auth[server_name.strip()] = normalized_headers
            if normalized_auth:
                context["mcp_auth"] = normalized_auth
        mcp_auth_fingerprint = payload.get("mcpAuthFingerprint")
        if isinstance(mcp_auth_fingerprint, str) and mcp_auth_fingerprint.strip():
            context["mcp_auth_fingerprint"] = mcp_auth_fingerprint.strip()
        if allow_script_runner:
            context["allow_script_runner"] = True
        if isinstance(payload.get("skipPlanApprovals"), bool):
            context["skip_plan_approvals"] = payload["skipPlanApprovals"]
        return context

    def _seed_initial_skill_context(initial_context: Dict[str, Any], message: ChatRequest) -> Dict[str, Any]:
        seeded = dict(initial_context or {})
        payload = _prepare_payload(message)
        for index in range(len(payload) - 1, -1, -1):
            item = payload[index]
            role = str(item.get("role") or "").strip().lower()
            if role not in {"user", "human"}:
                continue
            content = item.get("content")
            directive, _ = _extract_directive_from_text(_extract_text_from_content(content))
            if directive is None:
                break
            if directive.kind == "skill" and directive.skillId:
                skill = find_skill(settings.backend.skills_root, directive.skillId)
                if skill is not None and is_skill_allowed(skill, seeded):
                    seeded.pop("preferred_mcp_server", None)
                    activate_skill_context(seeded, skill)
            elif directive.kind == "mcp" and directive.serverId:
                seeded["preferred_mcp_server"] = directive.serverId
            break
        return seeded


    def _extract_tagged_files(content: str) -> List[str]:
        if not content:
            return []
        lines = content.splitlines()
        tagged: List[str] = []
        in_block = False
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if in_block:
                    break
                continue
            if stripped.startswith("Tagged files"):
                in_block = True
                continue
            if in_block:
                if stripped.startswith("-"):
                    candidate = stripped.lstrip("-").strip()
                    if candidate:
                        tagged.append(candidate)
                else:
                    break
        return tagged

    def _normalize_file_context_refs(raw_refs: Any) -> List[Dict[str, Any]]:
        normalized: List[Dict[str, Any]] = []
        if not isinstance(raw_refs, list):
            return normalized
        for item in raw_refs:
            if not isinstance(item, dict):
                continue
            source_name = str(item.get("sourceName") or "").strip()
            fingerprint = str(item.get("sourceVersionFingerprint") or "").strip()
            artifact_id = str(item.get("artifactId") or "").strip()
            if not source_name or not fingerprint or not artifact_id:
                continue
            normalized.append(
                {
                    "sourceFileId": item.get("sourceFileId"),
                    "sourceName": source_name,
                    "sourceMimeType": str(item.get("sourceMimeType") or "").strip() or None,
                    "sourceVersionFingerprint": fingerprint,
                    "artifactId": artifact_id,
                    "artifactVersion": item.get("artifactVersion"),
                    "derivedArtifactFileId": item.get("derivedArtifactFileId"),
                    "derivedArtifactPath": str(item.get("derivedArtifactPath") or "").strip() or None,
                    "effectiveMode": str(item.get("effectiveMode") or "part").strip() or "part",
                    "status": str(item.get("status") or "failed").strip() or "failed",
                    "summary": str(item.get("summary") or "").strip() or None,
                    "lastError": str(item.get("lastError") or "").strip() or None,
                }
            )
        return normalized

    def _load_mineru_text(workspace_id: str, tagged_paths: List[str]) -> str | None:
        output_root = rag_worker.store.config.raganything_output_dir
        for raw in tagged_paths:
            if not raw:
                continue
            name = Path(raw).name
            base = name
            if base.lower().endswith(".pdf"):
                base = base[:-4]
            candidates = [
                output_root / workspace_id / base / "docling" / f"{base}.md",
                output_root / workspace_id / base / "auto" / f"{base}.md",
            ]
            for md_path in candidates:
                if not md_path.exists():
                    continue
                try:
                    text = md_path.read_text(encoding="utf-8", errors="replace").strip()
                except Exception:
                    logger.exception("Failed reading parser markdown: %s", md_path)
                    continue
                if text:
                    return text
        return None

    def _load_tagged_html_outline(workspace_id: str, tagged_paths: List[str]) -> str | None:
        workspace_root = Path(settings.backend.workspace_root).resolve() / workspace_id
        for raw in tagged_paths:
            if not raw:
                continue
            normalized = str(raw).strip().lstrip("/").replace("\\", "/")
            if Path(normalized).suffix.lower() not in _TAGGED_HTML_EXTENSIONS:
                continue
            candidate = (workspace_root / normalized).resolve()
            if workspace_root not in candidate.parents and candidate != workspace_root:
                continue
            if not candidate.exists() or not candidate.is_file():
                continue
            outline = _extract_html_outline_from_path(candidate)
            if outline:
                return outline
        return None

    async def _prefetch_rag_context(
        workspace_id: str,
        prompt: str,
        tagged_paths_override: Sequence[str] | None = None,
    ) -> str | None:
        # Use extraction rather than relying on an exact marker string so backend text can evolve.
        if not prompt and not tagged_paths_override:
            return None
        tagged_paths = list(tagged_paths_override or _extract_tagged_files(prompt))
        rag_tagged_paths = _filter_rag_prefetchable_tagged_files(tagged_paths)
        if not rag_tagged_paths:
            return None
        keywords = _build_tagged_rag_keywords(prompt, rag_tagged_paths)
        rag_prompt = prompt
        if len(rag_tagged_paths) != len(tagged_paths):
            filtered_lines = ["Tagged files:"] + [f"- {path}" for path in rag_tagged_paths]
            rag_prompt = re.sub(
                r"(^|\n)Tagged files:\n(?:- .*(?:\n|$))+",
                ("\n" if "\nTagged files:" in prompt else "") + "\n".join(filtered_lines) + "\n",
                prompt,
                count=1,
            )
        try:
            response = await rag_worker.store.query_data(
                workspace_id,
                rag_prompt,
                mode="naive",
                include_references=False,
                hl_keywords=keywords,
                ll_keywords=keywords,
            )
            data = response.get("data") if isinstance(response, dict) else None
            chunks = data.get("chunks", []) if isinstance(data, dict) else []
            chunks = _filter_rag_chunks_to_tagged_paths(chunks, rag_tagged_paths)
            lines: List[str] = []
            for chunk in chunks:
                content = chunk.get("content") or ""
                if content and content.lstrip().startswith("SOURCE:"):
                    lines.append(content)
            if not lines:
                for chunk in chunks:
                    content = (chunk.get("content") or "").strip()
                    if not content:
                        continue
                    lines.append(content)
            if lines:
                non_textual = 0
                for content in lines:
                    lowered = content.lower()
                    if lowered.startswith("table analysis:") or lowered.startswith("discarded content analysis:"):
                        non_textual += 1
                if non_textual < len(lines):
                    compressed = _compress_tagged_context_lines(lines)
                    if compressed:
                        return compressed
            response = await rag_worker.store.query_data(
                workspace_id,
                rag_prompt,
                mode="hybrid",
                include_references=False,
                hl_keywords=keywords,
                ll_keywords=keywords,
            )
            data = response.get("data") if isinstance(response, dict) else None
            chunks = data.get("chunks", []) if isinstance(data, dict) else []
            chunks = _filter_rag_chunks_to_tagged_paths(chunks, rag_tagged_paths)
            lines = []
            for chunk in chunks:
                content = (chunk.get("content") or "").strip()
                if not content:
                    continue
                lines.append(content)
            if lines:
                non_textual = 0
                for content in lines:
                    lowered = content.lower()
                    if lowered.startswith("table analysis:") or lowered.startswith("discarded content analysis:"):
                        non_textual += 1
                if non_textual < len(lines):
                    compressed = _compress_tagged_context_lines(lines)
                    if compressed:
                        return compressed
            mineru_text = _load_mineru_text(workspace_id, rag_tagged_paths)
            if mineru_text:
                compressed = _compress_tagged_context_lines([mineru_text])
                if compressed:
                    return compressed
            html_outline = _load_tagged_html_outline(workspace_id, rag_tagged_paths)
            if html_outline:
                return html_outline
            return None
        except Exception:
            logger.exception("Failed to prefetch RAG context for tagged files.")
            return None

    def _get_thread_id(runtime: AgentRuntimeState, force_reset: bool) -> str:
        context = runtime.workspace_state.context
        if force_reset or not context.get("thread_id"):
            suffix = ""
            if isinstance(context, dict):
                user_id = context.get("user_id")
                if isinstance(user_id, str) and user_id.strip():
                    suffix = f":{user_id.strip()}"
            base = f"{runtime.agent_name}:{runtime.workspace_state.workspace_id}{suffix}"
            if force_reset:
                thread_id = f"{base}:{uuid4()}"
            else:
                thread_id = base
            context["thread_id"] = thread_id
        return context["thread_id"]

    def _build_agent_config(runtime: AgentRuntimeState, message: ChatRequest, callbacks=None) -> Dict[str, Any]:
        thread_id = _get_thread_id(runtime, message.forceReset)
        config: Dict[str, Any] = {"configurable": {"thread_id": thread_id}}
        if callbacks:
            config["callbacks"] = callbacks
        return config

    async def _invoke_agent(runtime: AgentRuntimeState, message: ChatRequest):
        agent = getattr(runtime, "agent", None)
        if agent is None:
            raise HTTPException(status_code=500, detail="Agent not initialized")
        context = getattr(runtime.workspace_state, "context", None)
        manager = context.get("data_agent_manager") if isinstance(context, dict) else None
        if manager and hasattr(manager, "reset_session"):
            manager.reset_session()
        payload = await _prepare_turn_payload(runtime, message, fresh_turn=True)
        lf = langfuse_langchain_callbacks()
        config = _build_agent_config(runtime, message, callbacks=lf or None)
        if hasattr(agent, "ainvoke"):
            return await agent.ainvoke({"messages": payload}, config=config)
        return agent.invoke({"messages": payload}, config=config)

    def _json_line(payload: Dict[str, Any]) -> bytes:
        return (json.dumps(payload, ensure_ascii=False) + "\n").encode("utf-8")

    def _message_to_text(message: Any) -> str:
        content = getattr(message, "content", None)
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            parts: List[str] = []
            for part in content:
                if isinstance(part, str):
                    parts.append(part)
                elif isinstance(part, dict):
                    if isinstance(part.get("text"), str):
                        parts.append(part["text"])
                    elif "content" in part:
                        parts.append(str(part["content"]))
                elif hasattr(part, "text"):
                    parts.append(str(part.text))
            if parts:
                return "".join(parts)
        if isinstance(message, dict):
            if isinstance(message.get("content"), str):
                return message["content"]
            if "text" in message and isinstance(message["text"], str):
                return message["text"]
        if hasattr(message, "text"):
            return str(message.text)
        return str(message)

    class _CallbackStreamingHandler(AsyncCallbackHandler):
        """Streams LangChain callback events into JSON payloads for the UI."""

        def __init__(self, text_fn, *, suppress_interrupt_tool_start: bool = False):
            super().__init__()
            self.queue: asyncio.Queue[Any] = asyncio.Queue()
            self._tool_names: Dict[str, str] = {}
            self._tool_meta: Dict[str, Any] = {}
            self._active_llm_runs: Set[str] = set()
            self._to_text = text_fn
            self._has_events = False
            self._has_assistant_text = False
            self._interrupt_emitted = False
            self._cancel_run: Optional[Callable[[], None]] = None
            self._suppress_interrupt_tool_start = suppress_interrupt_tool_start

        @property
        def has_events(self) -> bool:
            return self._has_events

        @property
        def has_assistant_text(self) -> bool:
            return self._has_assistant_text

        @property
        def interrupt_emitted(self) -> bool:
            return self._interrupt_emitted

        def attach_cancel(self, cancel_cb: Callable[[], None]) -> None:
            self._cancel_run = cancel_cb

        async def _emit(self, payload: Dict[str, Any]) -> None:
            self._has_events = True
            if payload.get("type") in {"token", "chunk"}:
                role = payload.get("role")
                if role is None or str(role).lower() == "assistant":
                    self._has_assistant_text = True
            await self.queue.put(payload)

        async def on_llm_new_token(
            self,
            token: str,
            *,
            run_id,
            **_: Any,
        ) -> None:
            if not token:
                return
            self._active_llm_runs.add(str(run_id))
            await self._emit({"type": "token", "content": token, "role": "assistant"})

        async def on_llm_end(self, response, *, run_id, **_: Any) -> None:
            run_key = str(run_id)
            if run_key in self._active_llm_runs:
                self._active_llm_runs.discard(run_key)
                return
            generations = getattr(response, "generations", None)
            if not generations:
                return
            text_parts: List[str] = []
            for generation in generations:
                if not generation:
                    continue
                candidate = generation[0]
                candidate_text = getattr(candidate, "text", None)
                if candidate_text:
                    text_parts.append(candidate_text)
            if text_parts:
                await self._emit(
                    {"type": "token", "content": "".join(text_parts), "role": "assistant"}
                )

        async def on_agent_action(self, action, **_: Any) -> None:
            log = getattr(action, "log", "")
            if log:
                await self._emit({"type": "thought", "content": log})

        async def on_agent_finish(self, finish, **_: Any) -> None:
            if self._has_assistant_text:
                return
            text = ""
            return_values = getattr(finish, "return_values", None)
            if isinstance(return_values, dict):
                candidate = return_values.get("output") or return_values.get("text")
                if isinstance(candidate, str):
                    text = candidate
            if not text:
                text = self._to_text(finish)
            if text:
                for piece in _chunk_text(text):
                    await self._emit({"type": "token", "content": piece, "role": "assistant"})

        async def on_tool_start(
            self,
            serialized,
            input_str,
            *,
            run_id,
            metadata: Dict[str, Any] | None = None,
            **_: Any,
        ) -> None:
            name = (serialized or {}).get("name") or (metadata or {}).get("name") or "tool"
            self._tool_names[str(run_id)] = name
            if name in _INTERRUPT_TOOL_NAMES:
                # On resumed clarification/action flows, let the tool consume the
                # resume payload first. Eagerly re-emitting the interrupt here can
                # cancel the resumed run before the answer is applied.
                if self._suppress_interrupt_tool_start:
                    return
                interrupt_payload = extract_interrupt_payload_from_tool_call(name, input_str)
                if interrupt_payload:
                    self._interrupt_emitted = True
                    await self._emit(interrupt_payload)
                    if self._cancel_run:
                        self._cancel_run()
                    return
            preview = input_str.strip()
            await self._emit(
                {
                    "type": "tool_start",
                    "name": name,
                    "content": preview[:200] if preview else "",
                }
            )

        async def on_tool_end(self, output, *, run_id, **_: Any) -> None:
            run_key = str(run_id)
            name = self._tool_names.pop(run_key, "tool")
            text = self._to_text(output)
            if name in _INTERRUPT_TOOL_NAMES:
                interrupt_payload = extract_interrupt_payload_from_tool_text(text)
                if interrupt_payload:
                    self._tool_meta.pop(run_key, None)
                    await self._emit(interrupt_payload)
                    return
            payload: Dict[str, Any] = {
                "type": "tool_end",
                "name": name,
                "content": text,
            }
            output_files = _extract_output_files_from_tool_result(name, text)
            meta = self._tool_meta.pop(run_key, None)
            if meta and meta.get("files"):
                output_files.extend(meta["files"])
            if output_files:
                dedup: Dict[str, Dict[str, Any]] = {}
                for item in output_files:
                    path = str(item.get("path") or "").strip()
                    if not path:
                        continue
                    dedup[path] = item
                payload["outputFiles"] = list(dedup.values())
            await self._emit(payload)

        async def on_tool_error(self, error, *, run_id, **_: Any) -> None:
            run_key = str(run_id)
            name = self._tool_names.pop(run_key, "tool")
            if _extract_interrupt_from_exception(error):
                self._tool_meta.pop(run_key, None)
                return
            await self._emit(
                {
                    "type": "tool_error",
                    "name": name,
                    "content": _format_exception(error),
                }
            )

        async def on_custom_event(
            self,
            name: str,
            data: Any,
            *,
            run_id,
            **_: Any,
        ) -> None:
            if name == "tool_artifacts" and isinstance(data, dict):
                self._tool_meta[str(run_id)] = data

    class _DeltaTracker:
        def __init__(self) -> None:
            self._latest: Dict[str, str] = {}

        def push(self, role: str, text: str) -> str:
            if not text:
                return ""
            previous = self._latest.get(role, "")
            if text.startswith(previous):
                delta = text[len(previous):]
            else:
                delta = text
            self._latest[role] = text
            return delta

    def _chunk_text(payload: str, max_chars: int = 60) -> Iterable[str]:
        if len(payload) <= max_chars:
            return [payload]

        chunks: List[str] = []
        start = 0
        text_length = len(payload)
        while start < text_length:
            end = min(text_length, start + max_chars)
            if end < text_length:
                newline = payload.rfind("\n", start, end)
                if newline > start + 40:
                    end = newline + 1
            chunk = payload[start:end]
            if chunk:
                chunks.append(chunk)
            start = max(end, start + 1)
        return chunks

    _INTERNAL_STREAM_TEXT_PATTERNS = (
        re.compile(r"^PLAN_(APPROVAL|EDIT|REJECTION|REJECT|CLARIFICATION|ACTION)_[A-Z_]+", re.IGNORECASE),
        re.compile(r"^Command\s*\(", re.IGNORECASE),
    )

    def _is_internal_stream_text(text: str) -> bool:
        normalized = (text or "").strip()
        if not normalized:
            return False
        return any(pattern.match(normalized) for pattern in _INTERNAL_STREAM_TEXT_PATTERNS)

    def _message_role(message: Any) -> str:
        for attr in ("type", "role"):
            value = getattr(message, attr, None)
            if isinstance(value, str):
                return value.lower()
            if isinstance(value, dict):
                role = value.get("role")
                if isinstance(role, str):
                    return role.lower()
        if isinstance(message, dict):
            role = message.get("role")
            if isinstance(role, str):
                return role.lower()
        return "assistant"

    def _parse_multi_mode_chunk(raw_chunk: Any) -> tuple[str | None, Any]:
        if isinstance(raw_chunk, tuple) and len(raw_chunk) == 2 and isinstance(raw_chunk[0], str):
            return raw_chunk[0], raw_chunk[1]
        return None, raw_chunk

    def _extract_messages(chunk: Any) -> List[Any] | None:
        if chunk is None:
            return None
        if isinstance(chunk, dict):
            if "messages" in chunk:
                return chunk.get("messages")  # type: ignore[return-value]
            output = chunk.get("output")
            if isinstance(output, dict) and "messages" in output:
                return output.get("messages")  # type: ignore[return-value]
        if isinstance(chunk, (list, tuple)):
            if len(chunk) == 2 and chunk[0] == "messages":
                candidate = chunk[1]
                if isinstance(candidate, (list, tuple)):
                    return list(candidate)
                if candidate is not None:
                    return [candidate]
                return None
            for item in chunk:
                if isinstance(item, dict) and "messages" in item:
                    return item.get("messages")  # type: ignore[return-value]
        return None

    def _build_interrupt_payload(raw: Any) -> Dict[str, Any] | None:
        if not raw or not isinstance(raw, (list, tuple)):
            return None
        first = raw[0] if raw else None
        if first is None:
            return None

        interrupt_value = None
        interrupt_id = None
        if isinstance(first, dict):
            interrupt_value = first.get("value")
            interrupt_id = first.get("id")
        else:
            interrupt_value = getattr(first, "value", None)
            interrupt_id = getattr(first, "id", None)

        if not isinstance(interrupt_value, dict):
            return None

        return normalize_interrupt_payload_value(interrupt_value, interrupt_id if isinstance(interrupt_id, str) else None)

    def _extract_interrupt_payload(chunk: Any) -> Dict[str, Any] | None:
        if not isinstance(chunk, dict):
            return None
        return _build_interrupt_payload(chunk.get("__interrupt__"))

    def _extract_interrupt_from_exception(error: BaseException) -> Dict[str, Any] | None:
        if isinstance(error, BaseExceptionGroup):
            for inner in error.exceptions:
                payload = _extract_interrupt_from_exception(inner)
                if payload:
                    return payload
            return None
        if isinstance(error, GraphInterrupt):
            return _build_interrupt_payload(error.args[0] if error.args else None)
        return None

    def _active_skill_policy(runtime: AgentRuntimeState) -> Dict[str, Any]:
        context = runtime.workspace_state.context or {}
        raw_policy = context.get("active_skill_policy") or {}
        if not isinstance(raw_policy, dict):
            raw_policy = {}
        raw_limit = raw_policy.get("pre_plan_search_limit", 0)
        raw_used = context.get("pre_plan_search_count", 0)
        try:
            pre_plan_search_limit = max(0, int(raw_limit or 0))
        except (TypeError, ValueError):
            pre_plan_search_limit = 0
        try:
            pre_plan_search_used = max(0, int(raw_used or 0))
        except (TypeError, ValueError):
            pre_plan_search_used = 0
        return {
            "skill": context.get("active_skill"),
            "requiresHitlPlan": bool(raw_policy.get("requires_hitl_plan", False)),
            "requiresArtifacts": bool(raw_policy.get("requires_workspace_artifacts", False)),
            "requiredArtifactsMode": raw_policy.get("required_artifacts_mode"),
            "prePlanSearchLimit": pre_plan_search_limit,
            "prePlanSearchUsed": pre_plan_search_used,
        }

    def _missing_required_artifacts(runtime: AgentRuntimeState) -> List[str]:
        context = runtime.workspace_state.context or {}
        policy = context.get("active_skill_policy") or {}
        if not isinstance(policy, dict):
            return []
        if not bool(policy.get("requires_workspace_artifacts", False)):
            return []
        root = runtime.workspace_state.root_path
        required = policy.get("required_artifacts") or []
        required_items = [str(item).strip() for item in required if str(item).strip()]
        if not required_items:
            return []
        missing: List[str] = []
        for item in required_items:
            if item.startswith("pattern:"):
                pattern = item[len("pattern:"):].lstrip("/")
                matched = False
                for child in root.rglob("*"):
                    if not child.is_file():
                        continue
                    rel = child.relative_to(root).as_posix()
                    if fnmatch.fnmatch(rel, pattern):
                        matched = True
                        break
                if not matched:
                    missing.append(item)
                continue
            rel = item.lstrip("/")
            if not (root / rel).exists():
                missing.append(item)
        return missing

    def _reset_turn_context(runtime: AgentRuntimeState) -> None:
        context = runtime.workspace_state.context or {}
        skip_plan_approvals = bool(context.get("skip_plan_approvals"))
        # Skill execution state is per top-level user task. Resumes should preserve it,
        # but a fresh user turn should not inherit approval or active-skill state.
        context.pop("active_skill", None)
        context.pop("active_skill_scope", None)
        context.pop("active_skill_policy", None)
        context.pop("last_plan_feedback", None)
        context.pop("last_plan_file_path", None)
        context.pop("preferred_mcp_server", None)
        context.pop("tagged_files", None)
        context.pop("tagged_rag_context", None)
        context.pop("loaded_skill_ids_this_turn", None)
        context.pop("skill_load_attempts_this_turn", None)
        context["tagged_files_only"] = False
        context["plan_approved"] = skip_plan_approvals
        context["pre_plan_search_count"] = 0

    async def _prepare_turn_payload(
        runtime: AgentRuntimeState,
        message: ChatRequest,
        *,
        fresh_turn: bool,
    ) -> List[Dict[str, Any]]:
        payload = _prepare_payload(message)
        if fresh_turn:
            _reset_turn_context(runtime)
        payload, latest_user_text = _apply_embedded_directives(runtime, payload)

        prompt_for_tagged_files = latest_user_text
        if not prompt_for_tagged_files:
            for index in range(len(payload) - 1, -1, -1):
                role = str(payload[index].get("role") or "").strip().lower()
                if role in {"user", "human"}:
                    prompt_for_tagged_files = _extract_text_from_content(payload[index].get("content"))
                    break
        if not prompt_for_tagged_files:
            prompt_for_tagged_files = message.message or ""
        message_file_context_refs = _normalize_file_context_refs(message.fileContextRefs)
        if message_file_context_refs:
            runtime.workspace_state.context["file_context_refs"] = message_file_context_refs
        active_file_context_refs = _normalize_file_context_refs(runtime.workspace_state.context.get("file_context_refs"))
        explicit_artifact_paths = [
            str(item.get("derivedArtifactPath") or "").strip()
            for item in active_file_context_refs
            if str(item.get("status") or "").strip().lower() in {"ready", "partial"}
            and str(item.get("derivedArtifactPath") or "").strip()
        ]
        pending_files = [
            str(item.get("sourceName") or "").strip()
            for item in active_file_context_refs
            if str(item.get("status") or "").strip().lower() == "pending"
        ]
        tagged_files = explicit_artifact_paths or _extract_tagged_files(prompt_for_tagged_files)
        guided_prompt = _append_tagged_file_guidance(prompt_for_tagged_files, tagged_files)
        guided_prompt = _append_artifact_first_guidance(
            guided_prompt,
            active_file_context_refs,
            tagged_files,
            multimodal_active=bool(message.messageContent),
        )
        if pending_files:
            pending_note = (
                "Attached files still being processed: "
                + ", ".join(pending_files)
                + ". Be explicit that understanding is still in progress."
            )
            guided_prompt = f"{guided_prompt.rstrip()}\n\n{pending_note}".strip()
        if guided_prompt != prompt_for_tagged_files:
            for index in range(len(payload) - 1, -1, -1):
                role = str(payload[index].get("role") or "").strip().lower()
                if role in {"user", "human"}:
                    payload[index]["content"] = _replace_content_text(payload[index].get("content"), guided_prompt)
                    break
            prompt_for_tagged_files = guided_prompt
        runtime.workspace_state.context["tagged_files"] = tagged_files
        env_tagged_files_rag_only = (os.getenv("TAGGED_FILES_RAG_ONLY", "false") or "false").strip().lower() in {
            "1",
            "true",
            "yes",
            "y",
            "on",
        }
        tagged_files_rag_only = _should_force_tagged_files_rag_only(
            env_tagged_files_rag_only=env_tagged_files_rag_only,
            explicit_artifact_paths=explicit_artifact_paths,
            message_content=message.messageContent,
            preferred_mcp_server=str(runtime.workspace_state.context.get("preferred_mcp_server") or "").strip() or None,
            prompt_for_tagged_files=prompt_for_tagged_files,
        )
        runtime.workspace_state.context["tagged_files_only"] = bool(tagged_files) and tagged_files_rag_only
        if tagged_files:
            rag_context = await _prefetch_rag_context(
                runtime.workspace_state.workspace_id,
                prompt_for_tagged_files,
                tagged_paths_override=tagged_files,
            )
            if rag_context:
                runtime.workspace_state.context["tagged_rag_context"] = rag_context
                if tagged_files_rag_only:
                    for index in range(len(payload) - 1, -1, -1):
                        role = str(payload[index].get("role") or "").strip().lower()
                        if role in {"user", "human"}:
                            payload[index]["content"] = (
                                _replace_content_text(
                                    payload[index].get("content"),
                                    (
                                        f"{prompt_for_tagged_files}\n\nRAG_CONTEXT:\n{rag_context}\n\n"
                                        "Use RAG_CONTEXT as the primary file-grounded source for attached artifacts. "
                                        "If the user explicitly asks to verify current facts, pricing, dates, or model versions, "
                                        "you may use the requested MCP server or other allowed fresh sources to validate or correct the artifact."
                                    ),
                                )
                            )
                            break
        _inject_host_datetime_context(payload)
        return payload

    _ASSISTANT_ROLES = {"assistant", "ai", "aimessagechunk"}
    _TOOL_ROLES = {"tool"}

    def _emit_text(role: str, text: str) -> Iterable[Dict[str, str]]:
        if not text:
            return []
        if role in _ASSISTANT_ROLES:
            if _is_internal_stream_text(text):
                return []
            return [
                {"type": "token", "content": piece, "role": "assistant"}
                for piece in _chunk_text(text)
            ]
        if role in _TOOL_ROLES:
            return [
                {"type": "thought", "content": text, "role": role}
            ]
        return []

    async def _stream_agent_response(
        runtime: AgentRuntimeState,
        message: ChatRequest,
        *,
        resume_decisions: Optional[List[Dict[str, Any]]] = None,
        resume_value: Any = None,
    ) -> AsyncGenerator[bytes, None]:
        agent = getattr(runtime, "agent", None)
        if agent is None:
            yield _json_line({"type": "error", "message": "Agent not initialized"})
            return
        context = getattr(runtime.workspace_state, "context", None)
        manager = context.get("data_agent_manager") if isinstance(context, dict) else None
        if manager and hasattr(manager, "reset_session"):
            manager.reset_session()

        if resume_decisions is None and resume_value is None:
            payload = await _prepare_turn_payload(runtime, message, fresh_turn=True)
        else:
            payload = _prepare_payload(message)
        handler = _CallbackStreamingHandler(
            _message_to_text,
            suppress_interrupt_tool_start=resume_decisions is not None or resume_value is not None,
        )
        sentinel = object()
        stream_started = asyncio.get_running_loop().time()
        saw_interrupt = False
        yield _json_line({"type": "policy", **_active_skill_policy(runtime)})
        logger.info(
            "Agent stream start: agent=%s workspace=%s",
            runtime.agent_name,
            runtime.workspace_state.workspace_id,
        )

        async def _agent_runner():
            try:
                nonlocal saw_interrupt
                stream_config = _build_agent_config(
                    runtime,
                    message,
                    callbacks=[handler, *langfuse_langchain_callbacks()],
                )
                stream_input: Any = {"messages": payload}
                if resume_decisions is not None:
                    stream_input = Command(resume={"decisions": resume_decisions})
                elif resume_value is not None:
                    stream_input = Command(resume=resume_value)
                final_result = None
                async for chunk in agent.astream(stream_input, config=stream_config, stream_mode="values"):
                    final_result = chunk
                    interrupt_payload = _extract_interrupt_payload(chunk)
                    if interrupt_payload:
                        saw_interrupt = True
                        await handler._emit(interrupt_payload)
                        return

                emitted = False
                interrupt_payload = _extract_interrupt_payload(final_result)
                if interrupt_payload:
                    saw_interrupt = True
                    emitted = True
                    await handler._emit(interrupt_payload)

                messages = _extract_messages(final_result)
                if messages and not handler.has_assistant_text:
                    tracker = _DeltaTracker()
                    for msg in messages:
                        text = _message_to_text(msg)
                        role = _message_role(msg)
                        delta = tracker.push(role, text)
                        if delta:
                            emitted = True
                            for event_payload in _emit_text(role, delta):
                                await handler._emit(event_payload)
                elif not handler.has_assistant_text:
                    text = _message_to_text(final_result)
                    if text and not _is_internal_stream_text(text):
                        emitted = True
                        for event_payload in _emit_text("assistant", text):
                            await handler._emit(event_payload)

                if not emitted and not handler.has_events:
                    await handler._emit(
                        {
                            "type": "thought",
                            "role": "assistant",
                            "content": "Model returned no output",
                        }
                    )
            except GraphInterrupt as exc:
                interrupt_payload = _extract_interrupt_from_exception(exc)
                if interrupt_payload:
                    saw_interrupt = True
                    await handler._emit(interrupt_payload)
                    return
                raise
            except asyncio.CancelledError:
                if handler.interrupt_emitted:
                    saw_interrupt = True
                    return
                raise
            except Exception as exc:  # pragma: no cover - streaming guard
                error_message = _format_exception(exc)
                logger.exception("Agent stream error: %s", error_message)
                await handler._emit({"type": "error", "message": error_message})
                raise
            finally:
                elapsed = asyncio.get_running_loop().time() - stream_started
                logger.info(
                    "Agent stream finished: agent=%s workspace=%s elapsed=%.2fs",
                    runtime.agent_name,
                    runtime.workspace_state.workspace_id,
                    elapsed,
                )
                await handler.queue.put(sentinel)

        task = asyncio.create_task(_agent_runner())
        handler.attach_cancel(task.cancel)
        try:
            while True:
                try:
                    event = await asyncio.wait_for(handler.queue.get(), timeout=15.0)
                except asyncio.TimeoutError:
                    yield _json_line({"type": "keepalive"})
                    continue
                if event is sentinel:
                    break
                yield _json_line(event)
            source_tracker.update_final_report(runtime.workspace_state)
            if not saw_interrupt:
                missing = _missing_required_artifacts(runtime)
                if missing:
                    runtime.workspace_state.context["artifact_contract_failed"] = True
                    yield _json_line(
                        {
                            "type": "contract_error",
                            "message": "Artifact contract not satisfied.",
                            "missing": missing,
                        }
                    )
            yield _json_line({"type": "done"})
        finally:
            await task

    @app.post("/agents/{agent_name}/workspace/{workspace_id}/chat", response_model=ChatResponse)
    async def chat(agent_name: str, workspace_id: str, chat_request: ChatRequest, request: Request):
        try:
            initial_context = _seed_initial_skill_context(_extract_request_context(request), chat_request)
            runtime = await registry.get_or_create(agent_name, workspace_id, initial_context=initial_context)
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc

        result = await _invoke_agent(runtime, chat_request)
        source_tracker.update_final_report(runtime.workspace_state)
        return ChatResponse(reply=result)

    @app.post("/agents/{agent_name}/workspace/{workspace_id}/chat/stream")
    async def chat_stream(agent_name: str, workspace_id: str, chat_request: ChatRequest, request: Request):
        try:
            initial_context = _seed_initial_skill_context(_extract_request_context(request), chat_request)
            runtime = await registry.get_or_create(agent_name, workspace_id, initial_context=initial_context)
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc

        stream = _stream_agent_response(runtime, chat_request)
        return StreamingResponse(stream, media_type="application/jsonl")

    @app.post("/agents/{agent_name}/workspace/{workspace_id}/chat/stream/resume")
    async def chat_stream_resume(
        agent_name: str,
        workspace_id: str,
        resume_request: ResumeChatRequest,
        request: Request,
    ):
        try:
            initial_context = _extract_request_context(request)
            runtime = await registry.get_or_create(agent_name, workspace_id, initial_context=initial_context)
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc

        decisions_payload: List[Dict[str, Any]] = []
        for item in resume_request.decisions:
            if hasattr(item, "model_dump"):
                decisions_payload.append(item.model_dump(exclude_none=True))  # type: ignore[attr-defined]
            else:
                decisions_payload.append(item.dict(exclude_none=True))  # type: ignore[attr-defined]
        placeholder = ChatRequest(message="", history=None, forceReset=False)
        stream = _stream_agent_response(runtime, placeholder, resume_decisions=decisions_payload)
        return StreamingResponse(stream, media_type="application/jsonl")

    @app.post("/agents/{agent_name}/workspace/{workspace_id}/chat/stream/respond")
    async def chat_stream_respond(
        agent_name: str,
        workspace_id: str,
        response_request: InterruptResponseRequest,
        request: Request,
    ):
        try:
            initial_context = _extract_request_context(request)
            runtime = await registry.get_or_create(agent_name, workspace_id, initial_context=initial_context)
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc

        if hasattr(response_request, "model_dump"):
            response_payload = response_request.model_dump(exclude_none=True)  # type: ignore[attr-defined]
        else:
            response_payload = response_request.dict(exclude_none=True)  # type: ignore[attr-defined]
        placeholder = ChatRequest(message="", history=None, forceReset=False)
        stream = _stream_agent_response(runtime, placeholder, resume_value=response_payload)
        return StreamingResponse(stream, media_type="application/jsonl")

    @app.post("/agents/{agent_name}/workspace/{workspace_id}/chat/stream/act")
    async def chat_stream_act(
        agent_name: str,
        workspace_id: str,
        action_request: InterruptActionRequest,
        request: Request,
    ):
        try:
            initial_context = _extract_request_context(request)
            runtime = await registry.get_or_create(agent_name, workspace_id, initial_context=initial_context)
        except ValueError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc

        if hasattr(action_request, "model_dump"):
            action_payload = action_request.model_dump(exclude_none=True)  # type: ignore[attr-defined]
        else:
            action_payload = action_request.dict(exclude_none=True)  # type: ignore[attr-defined]
        placeholder = ChatRequest(message="", history=None, forceReset=False)
        stream = _stream_agent_response(runtime, placeholder, resume_value=action_payload)
        return StreamingResponse(stream, media_type="application/jsonl")

    return app
