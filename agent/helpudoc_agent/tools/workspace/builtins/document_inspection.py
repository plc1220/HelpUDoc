"""Bounded, read-only inspection tools for workspace documents.

These tools deliberately operate on the original workspace file at agent run
time. They replace the requirement to pre-parse every upload before a chat can
start. Results are location-addressable so the model can cite pages,
paragraphs, sheets, and cells without loading an entire document into context.
"""
from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Any, Iterable, List, Optional

from langchain_core.tools import Tool, tool

from ....state import WorkspaceState
from .workspace_files import _display_path, _resolve_under_root, _workspace_root


_TEXT_SUFFIXES = {
    ".csv",
    ".html",
    ".htm",
    ".json",
    ".md",
    ".sql",
    ".toml",
    ".tsv",
    ".txt",
    ".yaml",
    ".yml",
}
_MAX_OUTPUT_CHARS = int(os.getenv("DOCUMENT_INSPECTION_MAX_OUTPUT_CHARS", "30000"))
_MAX_SEARCH_RESULTS = int(os.getenv("DOCUMENT_INSPECTION_MAX_SEARCH_RESULTS", "40"))
_MAX_SPREADSHEET_INSPECT_CELLS = int(os.getenv("DOCUMENT_INSPECTION_MAX_RANGE_CELLS", "10000"))
_MAX_SPREADSHEET_SCAN_CELLS = int(os.getenv("DOCUMENT_INSPECTION_MAX_SPREADSHEET_CELLS", "250000"))
_MAX_PPTX_INSPECT_SLIDES = int(os.getenv("DOCUMENT_INSPECTION_MAX_PPTX_SLIDES", "20"))
_MAX_PPTX_INVENTORY_SLIDES = int(os.getenv("DOCUMENT_INSPECTION_MAX_PPTX_INVENTORY_SLIDES", "50"))
_MAX_PPTX_SHAPES_PER_SLIDE = int(os.getenv("DOCUMENT_INSPECTION_MAX_PPTX_SHAPES_PER_SLIDE", "200"))
_MAX_PPTX_SCAN_UNITS = int(os.getenv("DOCUMENT_INSPECTION_MAX_PPTX_SCAN_UNITS", "250000"))
_DEFAULT_XLSX_RANGE = "A1:J25"

_SUPPORTED_TYPES_HINT = "Supported: PDF, DOCX, XLSX/XLSM, PPTX, and text/Markdown/CSV."

# Stable machine-readable error codes. Callers (skills, prompts, and the
# runtime loop guard) branch on these, so values must not be renamed.
ERROR_FILE_NOT_FOUND = "FILE_NOT_FOUND"
ERROR_AMBIGUOUS_FILE_PATH = "AMBIGUOUS_FILE_PATH"
ERROR_PATH_OUTSIDE_WORKSPACE = "PATH_OUTSIDE_WORKSPACE"
ERROR_UNSUPPORTED_DOCUMENT_TYPE = "UNSUPPORTED_DOCUMENT_TYPE"
ERROR_UNKNOWN_SHEET = "UNKNOWN_SHEET"
ERROR_INVALID_RANGE = "INVALID_RANGE"
ERROR_RANGE_TOO_LARGE = "RANGE_TOO_LARGE"
ERROR_MISSING_QUERY = "MISSING_QUERY"
ERROR_INVALID_ARGUMENT = "INVALID_ARGUMENT"
ERROR_DEPENDENCY_MISSING = "DEPENDENCY_MISSING"
ERROR_DOCUMENT_READ_FAILED = "DOCUMENT_READ_FAILED"

_INSPECT_FAILURE_PREFIX = "Document inspection failed"
_SEARCH_FAILURE_PREFIX = "Document search failed"


class DocumentToolError(Exception):
    """Error carrying a stable code and a recovery hint for the model."""

    def __init__(
        self,
        code: str,
        message: str,
        *,
        retryable: bool = False,
        suggested_next_call: Optional[str] = None,
    ) -> None:
        super().__init__(message)
        self.code = code
        self.retryable = retryable
        self.suggested_next_call = suggested_next_call


def _short_error(exc: BaseException) -> str:
    return str(exc).strip() or exc.__class__.__name__


def _default_next_call(code: str, tool_name: str) -> str:
    if code in {ERROR_FILE_NOT_FOUND, ERROR_AMBIGUOUS_FILE_PATH, ERROR_PATH_OUTSIDE_WORKSPACE}:
        return "ls(path='/') to confirm the exact workspace-relative file path"
    if code == ERROR_UNSUPPORTED_DOCUMENT_TYPE:
        return "none: ask the user for a supported export (PDF, DOCX, XLSX, PPTX, CSV, or text)"
    if code == ERROR_UNKNOWN_SHEET:
        return "inspect_document(file_path=<same file>) to list the available sheets"
    if code in {ERROR_INVALID_RANGE, ERROR_RANGE_TOO_LARGE}:
        return "inspect_document(file_path=<same file>, sheet_name=<sheet>, cell_range='A1:J25')"
    if code == ERROR_MISSING_QUERY:
        return "search_document(file_path=<same file>, query=<non-empty search text>)"
    if code == ERROR_DEPENDENCY_MISSING:
        return "none: report the missing document dependency instead of retrying"
    if code == ERROR_DOCUMENT_READ_FAILED:
        return f"{tool_name} once more with a smaller bounded range, then stop and report the failure"
    return "none: report the error instead of repeating the same call"


def _classify_error(exc: BaseException) -> tuple[str, bool]:
    if isinstance(exc, FileNotFoundError):
        return ERROR_FILE_NOT_FOUND, False
    if isinstance(exc, (ImportError, ModuleNotFoundError)):
        return ERROR_DEPENDENCY_MISSING, False
    if isinstance(exc, ValueError):
        text = str(exc).lower()
        if "remain inside the workspace" in text:
            return ERROR_PATH_OUTSIDE_WORKSPACE, False
        if "ambiguous" in text:
            return ERROR_AMBIGUOUS_FILE_PATH, False
        return ERROR_INVALID_ARGUMENT, False
    if isinstance(exc, OSError):
        return ERROR_DOCUMENT_READ_FAILED, True
    return ERROR_DOCUMENT_READ_FAILED, False


def _ok_envelope(payload: dict[str, Any]) -> str:
    """Success envelope: existing keys are preserved and ``status`` is added."""
    body = {"status": "ok"}
    body.update(payload)
    return _clip(json.dumps(body, ensure_ascii=False, default=str, indent=2))


def _error_envelope(tool_name: str, legacy_prefix: str, exc: BaseException) -> str:
    if isinstance(exc, DocumentToolError):
        code = exc.code
        retryable = exc.retryable
        suggested = exc.suggested_next_call or _default_next_call(code, tool_name)
    else:
        code, retryable = _classify_error(exc)
        suggested = _default_next_call(code, tool_name)
    body = {
        "status": "error",
        "tool": tool_name,
        # The legacy prefix is retained so older string-matching callers and
        # transcripts keep working after the envelope change.
        "message": f"{legacy_prefix}: {_short_error(exc)}",
        "errorCode": code,
        "retryable": retryable,
        "suggestedNextCall": suggested,
    }
    return _clip(json.dumps(body, ensure_ascii=False, default=str, indent=2))


def _clip(value: str, limit: int = _MAX_OUTPUT_CHARS) -> str:
    text = str(value or "")
    if limit > 0 and len(text) > limit:
        return text[:limit].rstrip() + "\n\n[Output truncated]"
    return text


def _clean_query(query: str) -> tuple[str, list[str]]:
    phrase = str(query or "").strip().lower()
    tokens = [token for token in re.findall(r"[\w.-]+", phrase) if len(token) > 1]
    return phrase, tokens


def _matches(text: str, phrase: str, tokens: Iterable[str]) -> bool:
    haystack = str(text or "").lower()
    if phrase and phrase in haystack:
        return True
    token_list = list(tokens)
    return bool(token_list) and all(token in haystack for token in token_list)


def _snippet(text: str, phrase: str, limit: int = 500) -> str:
    compact = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(compact) <= limit:
        return compact
    lowered = compact.lower()
    index = lowered.find(phrase) if phrase else -1
    if index < 0:
        index = 0
    start = max(0, index - limit // 3)
    end = min(len(compact), start + limit)
    prefix = "…" if start else ""
    suffix = "…" if end < len(compact) else ""
    return f"{prefix}{compact[start:end].strip()}{suffix}"


def _resolve_document(root: Path, file_path: str) -> Path:
    raw = str(file_path or "").strip().replace("\\", "/")
    if not raw:
        raise DocumentToolError(ERROR_INVALID_ARGUMENT, "file_path is required")
    candidate = _resolve_under_root(root, raw)
    if candidate.is_file():
        return candidate

    # A basename is convenient for @mentions, but it must resolve uniquely.
    basename = Path(raw).name
    matches = [
        path.resolve()
        for path in root.rglob(basename)
        if path.is_file() and ".system/document-cache/" not in path.as_posix()
    ]
    if len(matches) == 1:
        return matches[0]
    if len(matches) > 1:
        options = ", ".join(_display_path(root, item) for item in matches[:10])
        raise DocumentToolError(
            ERROR_AMBIGUOUS_FILE_PATH,
            f"File name is ambiguous. Use one of: {options}",
        )
    raise DocumentToolError(ERROR_FILE_NOT_FOUND, f"Workspace file not found: {raw}")


def _pdf_outline(reader: Any) -> list[str]:
    labels: list[str] = []

    def visit(items: Any, depth: int = 0) -> None:
        for item in items or []:
            if isinstance(item, list):
                visit(item, depth + 1)
                continue
            title = str(getattr(item, "title", "") or "").strip()
            if title:
                labels.append(f"{'  ' * depth}- {title}")
            if len(labels) >= 100:
                return

    try:
        visit(getattr(reader, "outline", []) or [])
    except Exception:
        return []
    return labels


def _inspect_pdf(path: Path, page_start: int, page_end: int) -> dict[str, Any]:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    page_count = len(reader.pages)
    if page_count == 0:
        return {"kind": "pdf", "pageCount": 0, "pages": []}
    start = max(1, page_start)
    end = min(page_count, max(start, page_end))
    pages = []
    for page_number in range(start, end + 1):
        try:
            content = reader.pages[page_number - 1].extract_text() or ""
        except Exception as exc:
            content = f"[Page extraction failed: {exc}]"
        pages.append({"page": page_number, "text": content.strip()})
    metadata = {
        str(key).lstrip("/"): str(value)
        for key, value in dict(getattr(reader, "metadata", {}) or {}).items()
        if value is not None
    }
    return {
        "kind": "pdf",
        "pageCount": page_count,
        "selectedPages": [start, end],
        "metadata": metadata,
        "outline": _pdf_outline(reader),
        "pages": pages,
    }


def _inspect_docx(path: Path, item_start: int, item_end: int) -> dict[str, Any]:
    from docx import Document  # type: ignore

    document = Document(str(path))
    headings = []
    items: list[dict[str, Any]] = []
    start = max(1, item_start)
    end = max(start, item_end)
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = str(paragraph.text or "").strip()
        if not text:
            continue
        style_name = str(getattr(paragraph.style, "name", "") or "")
        if style_name.lower().startswith("heading"):
            headings.append({"paragraph": index, "style": style_name, "text": text})
        if start <= index <= end:
            items.append({"kind": "paragraph", "index": index, "style": style_name, "text": text})
    for table_index, table in enumerate(document.tables, start=1):
        if table_index < start or table_index > end:
            continue
        rows = []
        for row_index, row in enumerate(table.rows, start=1):
            values = [str(cell.text or "").strip() for cell in row.cells]
            rows.append({"row": row_index, "values": values})
        items.append({"kind": "table", "index": table_index, "rows": rows})
    return {
        "kind": "docx",
        "paragraphCount": len(document.paragraphs),
        "tableCount": len(document.tables),
        "headings": headings[:200],
        "selectedRange": [start, end],
        "selectedItems": items,
    }


def _sheet_dimensions(sheet: Any, *, allow_force: bool) -> tuple[Optional[str], Optional[str]]:
    """Return ``(dimensions, error)`` for one sheet without ever raising.

    ``read_only`` workbooks produced by Google Drive/Sheets exports often omit
    the ``<dimension>`` element, and ``calculate_dimension()`` then raises. The
    ``force=True`` fallback scans the sheet, so it is only used when the caller
    is already paying for a full inventory pass.
    """
    try:
        dimensions = sheet.calculate_dimension()
        if dimensions:
            return str(dimensions), None
        first_error = "Worksheet returned an empty dimension"
    except Exception as exc:  # openpyxl raises ValueError for unsized sheets
        first_error = _short_error(exc)
    if not allow_force:
        return None, first_error
    try:
        return str(sheet.calculate_dimension(force=True)), None
    except TypeError:
        # Non read-only worksheets do not accept ``force``.
        return None, first_error
    except Exception as exc:
        return None, _short_error(exc)


def _sheet_inventory_entry(sheet: Any, *, allow_force: bool) -> dict[str, Any]:
    from openpyxl.utils.cell import range_boundaries  # type: ignore

    entry: dict[str, Any] = {"name": str(getattr(sheet, "title", "") or "")}
    dimensions, dimensions_error = _sheet_dimensions(sheet, allow_force=allow_force)
    max_row = getattr(sheet, "max_row", None)
    max_column = getattr(sheet, "max_column", None)
    if dimensions and (max_row is None or max_column is None):
        try:
            _, _, bounded_column, bounded_row = range_boundaries(str(dimensions))
            max_column = max_column if max_column is not None else bounded_column
            max_row = max_row if max_row is not None else bounded_row
        except Exception:
            pass
    entry["maxRow"] = int(max_row) if isinstance(max_row, int) else None
    entry["maxColumn"] = int(max_column) if isinstance(max_column, int) else None
    entry["dimensions"] = dimensions
    if dimensions_error:
        # Partial metadata is more useful than failing the whole inspection.
        entry["dimensionsError"] = dimensions_error
    return entry


def _xlsx_inventory(workbook: Any, *, allow_force: bool) -> list[dict[str, Any]]:
    inventory: list[dict[str, Any]] = []
    for sheet in getattr(workbook, "worksheets", []) or []:
        try:
            inventory.append(_sheet_inventory_entry(sheet, allow_force=allow_force))
        except Exception as exc:  # never let one sheet break the others
            inventory.append(
                {
                    "name": str(getattr(sheet, "title", "") or ""),
                    "maxRow": None,
                    "maxColumn": None,
                    "dimensions": None,
                    "dimensionsError": _short_error(exc),
                }
            )
    return inventory


def _read_xlsx_range(sheet: Any, requested_range: str) -> list[list[dict[str, Any]]]:
    from openpyxl.utils import get_column_letter  # type: ignore
    from openpyxl.utils.cell import range_boundaries  # type: ignore

    try:
        min_column, min_row, max_column, max_row = range_boundaries(requested_range)
    except ValueError as exc:
        raise DocumentToolError(
            ERROR_INVALID_RANGE,
            f"Invalid spreadsheet cell range: {requested_range}",
        ) from exc
    if any(value is None for value in (min_column, min_row, max_column, max_row)):
        raise DocumentToolError(
            ERROR_INVALID_RANGE,
            "Spreadsheet inspection ranges must include bounded rows and columns",
        )
    cell_count = (max_column - min_column + 1) * (max_row - min_row + 1)
    if cell_count > _MAX_SPREADSHEET_INSPECT_CELLS:
        raise DocumentToolError(
            ERROR_RANGE_TOO_LARGE,
            f"Spreadsheet inspection range is too large ({cell_count} cells); "
            f"request at most {_MAX_SPREADSHEET_INSPECT_CELLS} cells",
        )
    rows: list[list[dict[str, Any]]] = []
    # Explicit bounds keep this read independent of sheet dimension discovery,
    # and coordinates are derived from the bounds so padded/empty cells (which
    # carry no row/column of their own) stay addressable.
    iterator = sheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_column,
        max_col=max_column,
    )
    for row_offset, row in enumerate(iterator):
        row_number = min_row + row_offset
        cells: list[dict[str, Any]] = []
        for column_offset, cell in enumerate(row or ()):
            column_number = min_column + column_offset
            cells.append(
                {
                    "cell": f"{get_column_letter(column_number)}{row_number}",
                    "value": getattr(cell, "value", None),
                    "dataType": getattr(cell, "data_type", None),
                }
            )
        rows.append(cells)
    return rows


def _inspect_xlsx(path: Path, sheet_name: Optional[str], cell_range: Optional[str]) -> dict[str, Any]:
    from openpyxl import load_workbook  # type: ignore

    workbook = load_workbook(str(path), read_only=True, data_only=False)
    try:
        if not sheet_name:
            return {"kind": "xlsx", "sheets": _xlsx_inventory(workbook, allow_force=True)}
        if sheet_name not in workbook.sheetnames:
            raise DocumentToolError(
                ERROR_UNKNOWN_SHEET,
                f"Unknown sheet {sheet_name!r}. Available: {', '.join(workbook.sheetnames)}",
            )
        sheet = workbook[sheet_name]
        requested_range = str(cell_range or "").strip() or _DEFAULT_XLSX_RANGE
        # The targeted read runs first and never depends on dimension discovery.
        rows = _read_xlsx_range(sheet, requested_range)
        return {
            "kind": "xlsx",
            "sheets": _xlsx_inventory(workbook, allow_force=False),
            "selectedSheet": sheet_name,
            "selectedRange": requested_range,
            "cells": rows,
        }
    finally:
        workbook.close()


def _pptx_core_properties(presentation: Any) -> dict[str, Any]:
    properties = getattr(presentation, "core_properties", None)
    if properties is None:
        return {}
    metadata: dict[str, Any] = {}
    for field in (
        "title",
        "subject",
        "author",
        "keywords",
        "comments",
        "category",
        "language",
        "created",
        "modified",
        "last_modified_by",
        "revision",
        "version",
    ):
        try:
            value = getattr(properties, field, None)
        except Exception:
            continue
        if value not in (None, ""):
            serialized = value.isoformat() if hasattr(value, "isoformat") else value
            metadata[field] = _snippet(str(serialized), "", 500)
    return metadata


def _pptx_slide_title(slide: Any) -> Optional[str]:
    try:
        title_shape = slide.shapes.title
        text = str(getattr(title_shape, "text", "") or "").strip() if title_shape else ""
        return _snippet(text, "", 300) if text else None
    except Exception:
        return None


def _pptx_notes(slide: Any) -> list[dict[str, Any]]:
    try:
        if not bool(getattr(slide, "has_notes_slide", False)):
            return []
        text_frame = getattr(slide.notes_slide, "notes_text_frame", None)
        paragraphs = getattr(text_frame, "paragraphs", ()) or ()
    except Exception:
        return []
    notes: list[dict[str, Any]] = []
    for paragraph_index, paragraph in enumerate(paragraphs, start=1):
        text = str(getattr(paragraph, "text", "") or "").strip()
        if text:
            notes.append({"paragraph": paragraph_index, "text": text})
    return notes


def _pptx_alt_text(shape: Any) -> dict[str, str]:
    """Return accessible title/description from a shape's OOXML cNvPr node."""
    element = getattr(shape, "_element", None)
    if element is None:
        return {}
    try:
        for child in element.iter():
            if str(getattr(child, "tag", "")).rsplit("}", 1)[-1] != "cNvPr":
                continue
            values = {
                "title": str(child.get("title") or "").strip(),
                "description": str(child.get("descr") or "").strip(),
            }
            return {key: value for key, value in values.items() if value}
    except Exception:
        return {}
    return {}


def _walk_pptx_shapes(shapes: Any, index_prefix: tuple[int, ...] = ()) -> Iterable[tuple[Any, tuple[int, ...]]]:
    """Yield top-level and grouped shapes in deterministic XML order."""
    for index, shape in enumerate(shapes or (), start=1):
        index_path = index_prefix + (index,)
        yield shape, index_path
        children = getattr(shape, "shapes", None)
        if children is not None:
            yield from _walk_pptx_shapes(children, index_path)


def _pptx_shape_id(shape: Any, index_path: tuple[int, ...]) -> str:
    try:
        return str(int(shape.shape_id))
    except Exception:
        # Malformed decks may omit the non-visual shape ID. The XML-order path
        # remains deterministic and addressable within the inspected file.
        return "index-" + ".".join(str(item) for item in index_path)


def _pptx_shape_entry(
    shape: Any,
    *,
    slide_number: int,
    index_path: tuple[int, ...],
    remaining_chars: list[int],
) -> dict[str, Any]:
    shape_id = _pptx_shape_id(shape, index_path)
    location = f"slide:{slide_number}:shape:{shape_id}"
    entry: dict[str, Any] = {
        "location": location,
        "shapeId": shape_id,
        "shapeIndexPath": ".".join(str(item) for item in index_path),
        "name": _snippet(str(getattr(shape, "name", "") or ""), "", 200),
        "shapeType": str(getattr(shape, "shape_type", "") or ""),
    }
    bounds: dict[str, int] = {}
    for key in ("left", "top", "width", "height"):
        try:
            value = getattr(shape, key, None)
            if value is not None:
                bounds[key] = int(value)
        except Exception:
            continue
    if bounds:
        entry["boundsEmu"] = bounds

    raw_alt_text = _pptx_alt_text(shape)
    if raw_alt_text:
        alt_text = {key: _snippet(value, "", 1200) for key, value in raw_alt_text.items()}
        entry["altText"] = alt_text
        remaining_chars[0] -= sum(len(value) for value in alt_text.values())

    if bool(getattr(shape, "has_text_frame", False)):
        paragraphs = []
        try:
            source_paragraphs = shape.text_frame.paragraphs
        except Exception:
            source_paragraphs = ()
        for paragraph_index, paragraph in enumerate(source_paragraphs, start=1):
            text = str(getattr(paragraph, "text", "") or "").strip()
            if not text:
                continue
            clipped = _snippet(text, "", min(1200, max(120, remaining_chars[0])))
            paragraphs.append(
                {
                    "location": f"{location}:paragraph:{paragraph_index}",
                    "paragraph": paragraph_index,
                    "text": clipped,
                }
            )
            remaining_chars[0] -= len(clipped) + 100
            if remaining_chars[0] <= 0:
                entry["contentTruncated"] = True
                break
        if paragraphs:
            entry["paragraphs"] = paragraphs

    if bool(getattr(shape, "has_table", False)) and remaining_chars[0] > 0:
        try:
            table = shape.table
            rows = []
            for row_index, row in enumerate(table.rows, start=1):
                values = []
                for cell in row.cells:
                    value = _snippet(str(getattr(cell, "text", "") or "").strip(), "", 500)
                    values.append(value)
                    remaining_chars[0] -= len(value) + 20
                rows.append(
                    {
                        "location": f"{location}:table:row:{row_index}",
                        "row": row_index,
                        "values": values,
                    }
                )
                if remaining_chars[0] <= 0:
                    entry["contentTruncated"] = True
                    break
            entry["table"] = {
                "rowCount": len(table.rows),
                "columnCount": len(table.columns),
                "rows": rows,
            }
        except Exception as exc:
            entry["tableReadError"] = _short_error(exc)
    return entry


def _pptx_slide_inventory(slide: Any, slide_number: int) -> dict[str, Any]:
    notes = _pptx_notes(slide)
    try:
        layout_name = _snippet(str(getattr(slide.slide_layout, "name", "") or ""), "", 120)
    except Exception:
        layout_name = ""
    try:
        slide_id = int(slide.slide_id)
    except Exception:
        slide_id = None
    try:
        hidden = str(slide._element.get("show") or "1") == "0"
    except Exception:
        hidden = False
    return {
        "location": f"slide:{slide_number}",
        "slide": slide_number,
        "slideId": slide_id,
        "title": _pptx_slide_title(slide),
        "layout": layout_name,
        "shapeCount": len(slide.shapes),
        "hasSpeakerNotes": bool(notes),
        "hidden": hidden,
    }


def _trim_pptx_inspection_payload(payload: dict[str, Any]) -> None:
    """Trim optional PPTX detail while preserving a valid JSON envelope."""
    if _MAX_OUTPUT_CHARS <= 0:
        return
    target = max(2000, _MAX_OUTPUT_CHARS - 1500)

    def serialized_size() -> int:
        return len(json.dumps(payload, ensure_ascii=False, default=str))

    trimmed = False
    while serialized_size() > target:
        inventory = payload.get("slideInventory")
        if isinstance(inventory, list) and inventory:
            inventory.pop()
            payload["inventoryTruncated"] = True
            trimmed = True
            continue
        removed = False
        for slide in reversed(payload.get("slides") or []):
            shapes = slide.get("shapes") if isinstance(slide, dict) else None
            if isinstance(shapes, list) and shapes:
                shapes.pop()
                slide["contentTruncated"] = True
                removed = trimmed = True
                break
        if removed:
            continue
        for slide in reversed(payload.get("slides") or []):
            notes = slide.get("speakerNotes") if isinstance(slide, dict) else None
            if isinstance(notes, list) and notes:
                notes.pop()
                slide["contentTruncated"] = True
                removed = trimmed = True
                break
        if removed:
            continue
        metadata = payload.get("metadata")
        if isinstance(metadata, dict) and metadata:
            metadata.pop(next(reversed(metadata)))
            trimmed = True
            continue
        break
    if trimmed:
        payload["outputTruncated"] = True


def _inspect_pptx(path: Path, slide_start: int, slide_end: int) -> dict[str, Any]:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    slide_count = len(presentation.slides)
    width = int(presentation.slide_width)
    height = int(presentation.slide_height)
    inventory_limit = max(0, _MAX_PPTX_INVENTORY_SLIDES)
    inventory = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number > inventory_limit:
            break
        inventory.append(_pptx_slide_inventory(slide, slide_number))
    payload: dict[str, Any] = {
        "kind": "pptx",
        "slideCount": slide_count,
        "slideSize": {
            "widthEmu": width,
            "heightEmu": height,
            "widthInches": round(width / 914400, 3),
            "heightInches": round(height / 914400, 3),
        },
        "metadata": _pptx_core_properties(presentation),
        "slideInventory": inventory,
        "inventoryTruncated": slide_count > len(inventory),
    }
    if slide_count == 0:
        payload.update({"selectedSlides": [], "slides": []})
        return payload

    start = max(1, int(slide_start))
    if start > slide_count:
        raise DocumentToolError(
            ERROR_INVALID_RANGE,
            f"PPTX slide_start {start} exceeds the deck's {slide_count} slides",
            suggested_next_call=(
                f"inspect_document(file_path=<same file>, slide_start=1, "
                f"slide_end={min(slide_count, 5)})"
            ),
        )
    end = min(slide_count, max(start, int(slide_end)))
    requested_count = end - start + 1
    if requested_count > _MAX_PPTX_INSPECT_SLIDES:
        raise DocumentToolError(
            ERROR_RANGE_TOO_LARGE,
            f"PPTX inspection range is too large ({requested_count} slides); "
            f"request at most {_MAX_PPTX_INSPECT_SLIDES} slides",
            suggested_next_call=(
                f"inspect_document(file_path=<same file>, slide_start={start}, "
                f"slide_end={start + _MAX_PPTX_INSPECT_SLIDES - 1})"
            ),
        )

    # Keep the serialized JSON below the shared tool-output ceiling while
    # returning valid, addressable structure rather than clipping mid-JSON.
    remaining_chars = [max(2000, min(10000, _MAX_OUTPUT_CHARS // 3))]
    selected = []
    selected_numbers = list(range(start, end + 1))
    for slide_number in selected_numbers:
        slide = presentation.slides[slide_number - 1]
        summary = _pptx_slide_inventory(slide, slide_number)
        shapes = []
        shapes_truncated = False
        for shape_number, (shape, index_path) in enumerate(_walk_pptx_shapes(slide.shapes), start=1):
            if shape_number > _MAX_PPTX_SHAPES_PER_SLIDE or remaining_chars[0] <= 0:
                shapes_truncated = True
                break
            shapes.append(
                _pptx_shape_entry(
                    shape,
                    slide_number=slide_number,
                    index_path=index_path,
                    remaining_chars=remaining_chars,
                )
            )
            remaining_chars[0] -= 180
        notes = []
        for note in _pptx_notes(slide):
            if remaining_chars[0] <= 0:
                shapes_truncated = True
                break
            text = _snippet(note["text"], "", min(1200, max(120, remaining_chars[0])))
            notes.append(
                {
                    "location": f"slide:{slide_number}:note:paragraph:{note['paragraph']}",
                    "paragraph": note["paragraph"],
                    "text": text,
                }
            )
            remaining_chars[0] -= len(text) + 100
        summary.update(
            {
                "shapes": shapes,
                "speakerNotes": notes,
                "contentTruncated": shapes_truncated,
            }
        )
        selected.append(summary)
    payload.update(
        {
            "selectedSlides": [start, end],
            "slides": selected,
            "visualInspection": {
                "rendered": False,
                "renderRequiredForVisualQuestions": True,
                "slides": selected_numbers,
                "reason": (
                    "Text and OOXML extraction do not preserve final layout, clipping, "
                    "overlap, chart appearance, or image content. Render these selected "
                    "slides with the PPTX skill QA flow for appearance-sensitive questions."
                ),
            },
        }
    )
    _trim_pptx_inspection_payload(payload)
    return payload


def _inspect_text(path: Path, line_start: int, line_end: int) -> dict[str, Any]:
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    if not lines:
        return {"kind": "text", "lineCount": 0, "lines": []}
    start = max(1, line_start)
    end = min(len(lines), max(start, line_end))
    return {
        "kind": "text",
        "lineCount": len(lines),
        "selectedLines": [start, end],
        "lines": [{"line": index, "text": lines[index - 1]} for index in range(start, end + 1)],
    }


def _search_pdf(path: Path, phrase: str, tokens: list[str], limit: int) -> list[dict[str, Any]]:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    results = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            continue
        if _matches(text, phrase, tokens):
            results.append(
                {
                    "location": f"page:{page_number}",
                    "page": page_number,
                    "snippet": _snippet(text, phrase),
                }
            )
        if len(results) >= limit:
            break
    return results


def _search_docx(path: Path, phrase: str, tokens: list[str], limit: int) -> list[dict[str, Any]]:
    from docx import Document  # type: ignore

    document = Document(str(path))
    results = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = str(paragraph.text or "")
        if _matches(text, phrase, tokens):
            results.append(
                {
                    "location": f"paragraph:{index}",
                    "paragraph": index,
                    "style": str(getattr(paragraph.style, "name", "") or ""),
                    "snippet": _snippet(text, phrase),
                }
            )
        if len(results) >= limit:
            return results
    for table_index, table in enumerate(document.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            text = " | ".join(str(cell.text or "").strip() for cell in row.cells)
            if _matches(text, phrase, tokens):
                results.append(
                    {
                        "location": f"table:{table_index}:row:{row_index}",
                        "table": table_index,
                        "row": row_index,
                        "snippet": _snippet(text, phrase),
                    }
                )
            if len(results) >= limit:
                return results
    return results


def _search_xlsx(path: Path, phrase: str, tokens: list[str], limit: int) -> tuple[list[dict[str, Any]], bool]:
    from openpyxl import load_workbook  # type: ignore

    workbook = load_workbook(str(path), read_only=True, data_only=False)
    results: list[dict[str, Any]] = []
    scanned = 0
    truncated = False
    try:
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    scanned += 1
                    if scanned > _MAX_SPREADSHEET_SCAN_CELLS:
                        truncated = True
                        return results, truncated
                    if cell.value is None:
                        continue
                    text = str(cell.value)
                    if _matches(text, phrase, tokens):
                        results.append(
                            {
                                "location": f"sheet:{sheet.title}:cell:{cell.coordinate}",
                                "sheet": sheet.title,
                                "cell": cell.coordinate,
                                "value": cell.value,
                                "dataType": cell.data_type,
                            }
                        )
                    if len(results) >= limit:
                        return results, truncated
    finally:
        workbook.close()
    return results, truncated


def _search_pptx(path: Path, phrase: str, tokens: list[str], limit: int) -> tuple[list[dict[str, Any]], bool]:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    results: list[dict[str, Any]] = []
    scanned = 0

    def scan_unit(cost: int = 1) -> bool:
        nonlocal scanned
        scanned += max(1, cost)
        return scanned <= _MAX_PPTX_SCAN_UNITS

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape, index_path in _walk_pptx_shapes(slide.shapes):
            if not scan_unit():
                return results, True
            shape_id = _pptx_shape_id(shape, index_path)
            base_location = f"slide:{slide_number}:shape:{shape_id}"
            shape_name = _snippet(str(getattr(shape, "name", "") or ""), "", 200)

            if bool(getattr(shape, "has_text_frame", False)):
                try:
                    paragraphs = shape.text_frame.paragraphs
                except Exception:
                    paragraphs = ()
                for paragraph_index, paragraph in enumerate(paragraphs, start=1):
                    if not scan_unit():
                        return results, True
                    text = str(getattr(paragraph, "text", "") or "")
                    if _matches(text, phrase, tokens):
                        results.append(
                            {
                                "kind": "slideText",
                                "location": f"{base_location}:paragraph:{paragraph_index}",
                                "slide": slide_number,
                                "shapeId": shape_id,
                                "shapeName": shape_name,
                                "paragraph": paragraph_index,
                                "snippet": _snippet(text, phrase, 350),
                            }
                        )
                    if len(results) >= limit:
                        return results, False

            if bool(getattr(shape, "has_table", False)):
                try:
                    rows = shape.table.rows
                except Exception:
                    rows = ()
                for row_index, row in enumerate(rows, start=1):
                    cells = list(getattr(row, "cells", ()) or ())
                    if not scan_unit(len(cells)):
                        return results, True
                    text = " | ".join(str(getattr(cell, "text", "") or "").strip() for cell in cells)
                    if _matches(text, phrase, tokens):
                        results.append(
                            {
                                "kind": "tableRow",
                                "location": f"{base_location}:table:row:{row_index}",
                                "slide": slide_number,
                                "shapeId": shape_id,
                                "shapeName": shape_name,
                                "row": row_index,
                                "snippet": _snippet(text, phrase, 350),
                            }
                        )
                    if len(results) >= limit:
                        return results, False

            alt_text = _pptx_alt_text(shape)
            if alt_text:
                if not scan_unit(len(alt_text)):
                    return results, True
                text = " | ".join(value for value in alt_text.values() if value)
                if _matches(text, phrase, tokens):
                    results.append(
                        {
                            "kind": "altText",
                            "location": f"{base_location}:alt-text",
                            "slide": slide_number,
                            "shapeId": shape_id,
                            "shapeName": shape_name,
                            "altTextFields": list(alt_text),
                            "snippet": _snippet(text, phrase, 350),
                        }
                    )
                if len(results) >= limit:
                    return results, False

        for note in _pptx_notes(slide):
            if not scan_unit():
                return results, True
            text = note["text"]
            if _matches(text, phrase, tokens):
                results.append(
                    {
                        "kind": "speakerNote",
                        "location": f"slide:{slide_number}:note:paragraph:{note['paragraph']}",
                        "slide": slide_number,
                        "paragraph": note["paragraph"],
                        "snippet": _snippet(text, phrase, 350),
                    }
                )
            if len(results) >= limit:
                return results, False
    return results, False


def _search_text(path: Path, phrase: str, tokens: list[str], limit: int) -> list[dict[str, Any]]:
    results = []
    for line_number, line in enumerate(
        path.read_text(encoding="utf-8", errors="replace").splitlines(),
        start=1,
    ):
        if _matches(line, phrase, tokens):
            results.append(
                {
                    "location": f"line:{line_number}",
                    "line": line_number,
                    "snippet": _snippet(line, phrase),
                }
            )
        if len(results) >= limit:
            break
    return results


def build_inspect_document_tool(workspace_state: WorkspaceState) -> Tool:
    root = _workspace_root(workspace_state)

    @tool
    def inspect_document(
        file_path: str,
        page_start: int = 1,
        page_end: int = 5,
        slide_start: int = 1,
        slide_end: int = 5,
        item_start: int = 1,
        item_end: int = 40,
        sheet_name: Optional[str] = None,
        cell_range: Optional[str] = None,
    ) -> str:
        """Inspect a bounded range of a PDF, DOCX, XLSX, PPTX, or text workspace file."""
        try:
            path = _resolve_document(root, file_path)
            suffix = path.suffix.lower()
            if suffix == ".pdf":
                payload = _inspect_pdf(path, page_start, page_end)
            elif suffix == ".docx":
                payload = _inspect_docx(path, item_start, item_end)
            elif suffix in {".xlsx", ".xlsm"}:
                payload = _inspect_xlsx(path, sheet_name, cell_range)
            elif suffix == ".pptx":
                payload = _inspect_pptx(path, slide_start, slide_end)
            elif suffix in _TEXT_SUFFIXES:
                payload = _inspect_text(path, item_start, item_end)
            else:
                raise DocumentToolError(
                    ERROR_UNSUPPORTED_DOCUMENT_TYPE,
                    f"Unsupported document type {suffix or '[none]'}. {_SUPPORTED_TYPES_HINT}",
                )
            payload["file"] = _display_path(root, path)
            return _ok_envelope(payload)
        except Exception as exc:
            return _error_envelope("inspect_document", _INSPECT_FAILURE_PREFIX, exc)

    inspect_document.name = "inspect_document"
    inspect_document.description = (
        "Read a bounded range from an original workspace PDF, DOCX, XLSX/XLSM, "
        "PPTX, or text file. Start with metadata/headings/sheet or slide inventory, "
        "then request only the pages, paragraphs, slides, speaker notes, sheet, or "
        "cells needed. PPTX inspection returns render-needed metadata for visual "
        "questions. Returns JSON with "
        "status='ok', or status='error' plus errorCode, retryable, and "
        "suggestedNextCall. Do not repeat a call whose error is not retryable."
    )
    return inspect_document


def build_search_document_tool(workspace_state: WorkspaceState) -> Tool:
    root = _workspace_root(workspace_state)

    @tool
    def search_document(file_path: str, query: str, max_results: int = 20) -> str:
        """Search an original workspace document and return location-addressable matches."""
        phrase, tokens = _clean_query(query)
        if not phrase:
            return _error_envelope(
                "search_document",
                _SEARCH_FAILURE_PREFIX,
                DocumentToolError(ERROR_MISSING_QUERY, "query is required"),
            )
        limit = max(1, min(int(max_results or 20), _MAX_SEARCH_RESULTS))
        try:
            path = _resolve_document(root, file_path)
            suffix = path.suffix.lower()
            truncated = False
            if suffix == ".pdf":
                results = _search_pdf(path, phrase, tokens, limit)
            elif suffix == ".docx":
                results = _search_docx(path, phrase, tokens, limit)
            elif suffix in {".xlsx", ".xlsm"}:
                results, truncated = _search_xlsx(path, phrase, tokens, limit)
            elif suffix == ".pptx":
                results, truncated = _search_pptx(path, phrase, tokens, limit)
            elif suffix in _TEXT_SUFFIXES:
                results = _search_text(path, phrase, tokens, limit)
            else:
                raise DocumentToolError(
                    ERROR_UNSUPPORTED_DOCUMENT_TYPE,
                    f"Unsupported document type {suffix or '[none]'}. {_SUPPORTED_TYPES_HINT}",
                )
            payload = {
                "file": _display_path(root, path),
                "query": query,
                "resultCount": len(results),
                "results": results,
                "scanTruncated": truncated,
            }
            return _ok_envelope(payload)
        except Exception as exc:
            return _error_envelope("search_document", _SEARCH_FAILURE_PREFIX, exc)

    search_document.name = "search_document"
    search_document.description = (
        "Search inside an original PDF, DOCX, XLSX/XLSM, PPTX, or text workspace file. "
        "Returns page, paragraph, table-row, slide/shape, speaker-note, alt-text, "
        "sheet/cell, or line locations for grounded follow-up inspection. "
        "Returns JSON with status='ok', or status='error' plus errorCode, retryable, "
        "and suggestedNextCall. Do not repeat a call whose error is not retryable."
    )
    return search_document


def build_document_inspection_tools(workspace_state: WorkspaceState) -> List[Tool]:
    return [
        build_inspect_document_tool(workspace_state),
        build_search_document_tool(workspace_state),
    ]
